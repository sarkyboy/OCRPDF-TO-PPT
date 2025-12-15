"""
导出功能模块 - PPT/PDF/图片导出

说明：
- 这里的函数以 editor 实例作为第一个参数（等价于原先的 self）。
"""

from __future__ import annotations

import os
import tempfile
import threading

import tkinter as tk
from tkinter import filedialog, messagebox

from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..constants import (
    COLOR_GRAY,
    COLOR_GREEN,
    COLOR_THEME,
    COLOR_WHITE,
    FONT_FAMILY,
    Px,
)


def _get_page_background_image(editor, page) -> Image.Image:
    """
    获取页面导出用的底图（包含背景/原图 + 图层合成，不包含文本渲染）。
    """
    if hasattr(editor, "get_page_composited_background"):
        try:
            img = editor.get_page_composited_background(page)
            if img is not None:
                return img
        except Exception:
            pass

    if page.get("bg_path") and os.path.exists(page["bg_path"]):
        return Image.open(page["bg_path"])

    return page["image"].copy()


def generate_multi_page_ppt(editor) -> None:
    if not editor.pages:
        editor.update_status("请先导入图片")
        return

    editor.save_current_page()

    save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint文件", "*.pptx")])
    if not save_path:
        return

    temp_paths: list[str] = []

    try:
        prs = Presentation()

        for page_idx, page in enumerate(editor.pages):
            editor.update_status(f"生成第 {page_idx+1}/{len(editor.pages)} 页...")

            bg_img = _get_page_background_image(editor, page)
            if bg_img.mode == "RGBA":
                bg_img = bg_img.convert("RGB")
            img_w, img_h = bg_img.size

            if page_idx == 0:
                prs.slide_width = Px(img_w)
                prs.slide_height = Px(img_h)

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            temp_bg = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            temp_bg_path = temp_bg.name
            temp_bg.close()
            bg_img.save(temp_bg_path)
            temp_paths.append(temp_bg_path)

            slide.shapes.add_picture(temp_bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            for box_data in page.get("text_boxes", []):
                if not box_data.get("text"):
                    continue

                x = box_data["x"]
                y = box_data["y"]
                w = box_data["width"]
                h = box_data["height"]
                font_size = box_data.get("font_size", 16)

                textbox = slide.shapes.add_textbox(Px(x), Px(y), Px(w), Px(h))
                tf = textbox.text_frame
                tf.word_wrap = False
                tf.margin_left = Px(2)
                tf.margin_right = Px(2)
                tf.margin_top = Px(1)
                tf.margin_bottom = Px(1)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = tf.paragraphs[0]
                p.text = box_data["text"]

                align = box_data.get("align", "left")
                p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
                    align, PP_ALIGN.LEFT
                )

                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(font_size)
                    run.font.name = box_data.get("font_name", "微软雅黑")
                    run.font.bold = box_data.get("bold", False)
                    run.font.italic = box_data.get("italic", False)

                    color_hex = box_data.get("font_color", "#000000").lstrip("#")
                    run.font.color.rgb = RGBColor(
                        int(color_hex[0:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16),
                    )

                p.line_spacing = 1.0
                p.space_before = Pt(0)
                p.space_after = Pt(0)

        prs.save(save_path)
        messagebox.showinfo("成功", f"PPT已保存！\n共 {len(editor.pages)} 页\n{save_path}")
        editor.update_status(f"PPT生成成功！共 {len(editor.pages)} 页 ✓")

    except Exception as e:
        messagebox.showerror("失败", f"生成失败: {e}")

    finally:
        for p in temp_paths:
            try:
                os.remove(p)
            except Exception:
                pass


def export_as_pdf(editor) -> None:
    """导出为PDF文件"""
    if not editor.pages:
        editor.update_status("没有可导出的内容")
        messagebox.showwarning("提示", "请先导入图片或PDF")
        return

    save_path = filedialog.asksaveasfilename(
        defaultextension=".pdf",
        filetypes=[("PDF文件", "*.pdf")],
        initialfile="output.pdf",
    )
    if not save_path:
        return

    editor.update_status("正在生成PDF...")

    def export_pdf():
        try:
            editor.root.after(0, editor.save_current_page)

            pdf_images: list[Image.Image] = []

            for page_idx, page in enumerate(editor.pages):
                editor.root.after(
                    0,
                    lambda idx=page_idx + 1, total=len(editor.pages): editor.update_status(
                        f"正在渲染第 {idx}/{total} 页..."
                    ),
                )

                bg_image = _get_page_background_image(editor, page)

                preview_img = bg_image.copy()
                if preview_img.mode != "RGB":
                    preview_img = preview_img.convert("RGB")

                draw = ImageDraw.Draw(preview_img)

                for box_data in page.get("text_boxes", []):
                    if not box_data.get("text"):
                        continue

                    try:
                        pixel_font_size = int(box_data.get("font_size", 16) * 96 / 72)
                        font_path = editor._get_font_path(box_data.get("font_name", "微软雅黑"))

                        if font_path and os.path.exists(font_path):
                            font = ImageFont.truetype(font_path, pixel_font_size)
                        else:
                            font = ImageFont.load_default()

                        color_hex = box_data.get("font_color", "#000000").lstrip("#")
                        r = int(color_hex[0:2], 16)
                        g = int(color_hex[2:4], 16)
                        b = int(color_hex[4:6], 16)

                        x, y = box_data["x"], box_data["y"]
                        w, h = box_data["width"], box_data["height"]

                        try:
                            bbox = draw.textbbox((0, 0), box_data["text"], font=font)
                            text_width = bbox[2] - bbox[0]
                            text_height = bbox[3] - bbox[1]
                            bbox_x0, bbox_y0 = bbox[0], bbox[1]
                        except Exception:
                            text_width = len(box_data["text"]) * pixel_font_size * 0.6
                            text_height = pixel_font_size
                            bbox_x0, bbox_y0 = 0, 0

                        align = box_data.get("align", "left")
                        if align == "center":
                            text_x = x + (w - text_width) // 2 - bbox_x0
                        elif align == "right":
                            text_x = x + w - text_width - 3 - bbox_x0
                        else:
                            text_x = x + 3 - bbox_x0

                        text_y = y + (h - text_height) // 2 - bbox_y0
                        draw.text((text_x, text_y), box_data["text"], font=font, fill=(r, g, b))

                    except Exception as e:
                        print(f"绘制文字失败 (页{page_idx+1}): {e}")
                        continue

                pdf_images.append(preview_img)

            if pdf_images:
                editor.root.after(0, lambda: editor.update_status("正在保存PDF文件..."))
                pdf_images[0].save(
                    save_path,
                    "PDF",
                    save_all=True,
                    append_images=pdf_images[1:],
                    resolution=100.0,
                )

                editor.root.after(
                    0,
                    lambda: messagebox.showinfo(
                        "成功",
                        f"PDF导出成功！\n\n" f"共 {len(editor.pages)} 页\n" f"保存位置：\n{save_path}",
                    ),
                )
                editor.root.after(0, lambda: editor.update_status("PDF导出成功！"))

        except Exception as e:
            import traceback

            error_msg = traceback.format_exc()
            print(f"PDF导出失败:\n{error_msg}")
            editor.root.after(0, lambda: messagebox.showerror("错误", f"PDF导出失败:\n\n{str(e)}"))
            editor.root.after(0, lambda: editor.update_status("PDF导出失败"))

    threading.Thread(target=export_pdf, daemon=True).start()


def export_as_images(editor) -> None:
    """导出为图片序列"""
    if not editor.pages:
        editor.update_status("没有可导出的内容")
        messagebox.showwarning("提示", "请先导入图片或PDF")
        return

    folder_path = filedialog.askdirectory(title="选择导出目录")
    if not folder_path:
        return

    show_image_format_dialog(editor, folder_path)


def _show_image_format_dialog(editor, folder_path: str) -> None:
    return show_image_format_dialog(editor, folder_path)


def show_image_format_dialog(editor, folder_path: str) -> None:
    """显示图片格式选择对话框"""
    format_dialog = tk.Toplevel(editor.root)
    format_dialog.title("选择图片格式")
    format_dialog.geometry("350x220")
    format_dialog.configure(bg=COLOR_WHITE)
    format_dialog.transient(editor.root)
    format_dialog.grab_set()

    format_dialog.update_idletasks()
    x = (format_dialog.winfo_screenwidth() - 350) // 2
    y = (format_dialog.winfo_screenheight() - 220) // 2
    format_dialog.geometry(f"+{x}+{y}")

    title_frame = tk.Frame(format_dialog, bg=COLOR_THEME, height=40)
    title_frame.pack(fill=tk.X)
    title_frame.pack_propagate(False)
    tk.Label(
        title_frame,
        text="  选择图片格式",
        bg=COLOR_THEME,
        fg="white",
        font=(FONT_FAMILY, 11, "bold"),
    ).pack(side=tk.LEFT, pady=8)

    content = tk.Frame(format_dialog, bg=COLOR_WHITE, padx=20, pady=15)
    content.pack(fill=tk.BOTH, expand=True)

    format_var = tk.StringVar(value="PNG")
    quality_var = tk.IntVar(value=95)

    format_frame = tk.Frame(content, bg=COLOR_WHITE)
    format_frame.pack(pady=10, fill=tk.X)

    tk.Radiobutton(
        format_frame,
        text="PNG - 无损压缩，高质量（推荐）",
        variable=format_var,
        value="PNG",
        bg=COLOR_WHITE,
        font=(FONT_FAMILY, 10),
    ).pack(anchor="w", pady=3)
    tk.Radiobutton(
        format_frame,
        text="JPEG - 有损压缩，文件较小",
        variable=format_var,
        value="JPEG",
        bg=COLOR_WHITE,
        font=(FONT_FAMILY, 10),
    ).pack(anchor="w", pady=3)

    quality_frame = tk.Frame(content, bg=COLOR_WHITE)
    quality_frame.pack(pady=10, fill=tk.X)

    tk.Label(quality_frame, text="JPEG质量 (1-100):", bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(
        side=tk.LEFT
    )

    quality_scale = tk.Scale(
        quality_frame,
        from_=50,
        to=100,
        orient=tk.HORIZONTAL,
        variable=quality_var,
        bg=COLOR_WHITE,
        length=150,
    )
    quality_scale.pack(side=tk.LEFT, padx=10)

    btn_frame = tk.Frame(content, bg=COLOR_WHITE)
    btn_frame.pack(pady=15)

    def start_export():
        img_format = format_var.get()
        quality = quality_var.get()
        format_dialog.destroy()
        do_export_images(editor, folder_path, img_format, quality)

    tk.Button(
        btn_frame,
        text="开始导出",
        command=start_export,
        bg=COLOR_GREEN,
        fg="white",
        font=(FONT_FAMILY, 10),
        padx=20,
        pady=5,
        cursor="hand2",
        relief=tk.FLAT,
    ).pack(side=tk.LEFT, padx=5)

    tk.Button(
        btn_frame,
        text="取消",
        command=format_dialog.destroy,
        bg=COLOR_GRAY,
        fg="white",
        font=(FONT_FAMILY, 10),
        padx=20,
        pady=5,
        cursor="hand2",
        relief=tk.FLAT,
    ).pack(side=tk.LEFT, padx=5)


def do_export_images(editor, folder_path: str, img_format: str, quality: int) -> None:
    """执行图片导出"""
    editor.update_status("正在导出图片...")

    def export_images():
        try:
            editor.root.after(0, editor.save_current_page)

            for page_idx, page in enumerate(editor.pages):
                editor.root.after(
                    0,
                    lambda idx=page_idx + 1, total=len(editor.pages): editor.update_status(
                        f"正在导出第 {idx}/{total} 页..."
                    ),
                )

                bg_image = _get_page_background_image(editor, page)

                preview_img = bg_image.copy()
                if preview_img.mode not in ["RGB", "RGBA"]:
                    preview_img = preview_img.convert("RGB")

                draw = ImageDraw.Draw(preview_img)

                for box_data in page.get("text_boxes", []):
                    if not box_data.get("text"):
                        continue

                    try:
                        pixel_font_size = int(box_data.get("font_size", 16) * 96 / 72)
                        font_path = editor._get_font_path(box_data.get("font_name", "微软雅黑"))

                        if font_path and os.path.exists(font_path):
                            font = ImageFont.truetype(font_path, pixel_font_size)
                        else:
                            font = ImageFont.load_default()

                        color_hex = box_data.get("font_color", "#000000").lstrip("#")
                        r = int(color_hex[0:2], 16)
                        g = int(color_hex[2:4], 16)
                        b = int(color_hex[4:6], 16)

                        x, y = box_data["x"], box_data["y"]
                        w, h = box_data["width"], box_data["height"]

                        try:
                            bbox = draw.textbbox((0, 0), box_data["text"], font=font)
                            text_width = bbox[2] - bbox[0]
                            text_height = bbox[3] - bbox[1]
                            bbox_x0, bbox_y0 = bbox[0], bbox[1]
                        except Exception:
                            text_width = len(box_data["text"]) * pixel_font_size * 0.6
                            text_height = pixel_font_size
                            bbox_x0, bbox_y0 = 0, 0

                        align = box_data.get("align", "left")
                        if align == "center":
                            text_x = x + (w - text_width) // 2 - bbox_x0
                        elif align == "right":
                            text_x = x + w - text_width - 3 - bbox_x0
                        else:
                            text_x = x + 3 - bbox_x0

                        text_y = y + (h - text_height) // 2 - bbox_y0
                        draw.text((text_x, text_y), box_data["text"], font=font, fill=(r, g, b))

                    except Exception as e:
                        print(f"绘制文字失败 (页{page_idx+1}): {e}")
                        continue

                ext = ".png" if img_format == "PNG" else ".jpg"
                save_path = os.path.join(folder_path, f"page_{page_idx+1:03d}{ext}")

                if img_format == "PNG":
                    if preview_img.mode == "RGBA":
                        preview_img.save(save_path, "PNG")
                    else:
                        preview_img.convert("RGB").save(save_path, "PNG")
                else:
                    if preview_img.mode == "RGBA":
                        preview_img = preview_img.convert("RGB")
                    preview_img.save(save_path, "JPEG", quality=quality)

            editor.root.after(
                0,
                lambda: messagebox.showinfo(
                    "成功",
                    f"图片导出成功！\n\n"
                    f"共导出 {len(editor.pages)} 张图片\n"
                    f"格式：{img_format}\n"
                    f"保存位置：\n{folder_path}",
                ),
            )
            editor.root.after(0, lambda: editor.update_status(f"图片导出成功！共 {len(editor.pages)} 张"))

        except Exception as e:
            import traceback

            error_msg = traceback.format_exc()
            print(f"图片导出失败:\n{error_msg}")
            editor.root.after(0, lambda: messagebox.showerror("错误", f"图片导出失败:\n\n{str(e)}"))
            editor.root.after(0, lambda: editor.update_status("图片导出失败"))

    threading.Thread(target=export_images, daemon=True).start()


def _do_export_images(editor, folder_path: str, img_format: str, quality: int) -> None:
    return do_export_images(editor, folder_path, img_format, quality)
