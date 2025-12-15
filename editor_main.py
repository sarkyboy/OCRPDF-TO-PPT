"""
ç°ä»£åŒ–PPTç¼–è¾‘å™¨ - ä»¿PowerPointç•Œé¢
UIé£æ ¼ï¼šå‚è€ƒPowerPointçš„ç°ä»£å¸ƒå±€
- é¡¶éƒ¨ï¼šçº¢è‰²æ ‡é¢˜æ  + åŒè¡Œå·¥å…·æ 
- å·¦ä¾§ï¼šé¡µé¢ç¼©ç•¥å›¾å¯¼èˆª
- ä¸­é—´ï¼šä¸»ç¼–è¾‘ç”»å¸ƒ
- å³ä¾§ï¼šå±æ€§é¢æ¿
- åº•éƒ¨ï¼šçº¢è‰²çŠ¶æ€æ 
"""

import tkinter as tk
from tkinter import filedialog, messagebox, ttk, colorchooser, simpledialog
from PIL import Image, ImageTk, ImageDraw, ImageFont, ImageFilter, ImageChops, ImageOps
import json
import os
import threading
import logging
import cv2
import numpy as np
import tempfile
import copy
import math
from datetime import datetime
import requests
import base64
from io import BytesIO
import uuid

# AIå›¾ç‰‡ç”ŸæˆAPIæ”¯æŒ
from .ai_image_api_module import AIImageAPIManager, blend_images

# PDFæ”¯æŒ - ä½¿ç”¨PyMuPDFï¼Œæ›´ç®€å•ï¼Œä¸éœ€è¦Poppler
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("æç¤º: å®‰è£… PyMuPDF å¯æ”¯æŒPDFå¯¼å…¥")
    print("      pip install PyMuPDF")


logging.getLogger("ppocr").setLevel(logging.WARNING)

from .config import get_base_dir, load_config, save_config
from .constants import (
    COLOR_THEME,
    COLOR_THEME_HOVER,
    COLOR_RIBBON_BG,
    COLOR_RIBBON_ROW2,
    COLOR_CANVAS_BG,
    COLOR_SIDEBAR_BG,
    COLOR_WHITE,
    COLOR_TEXT,
    COLOR_BLUE,
    COLOR_GREEN,
    COLOR_ORANGE,
    COLOR_PURPLE,
    COLOR_RED,
    COLOR_GRAY,
    FONT_FAMILY,
    Px,
)
from .textbox import TextBox

from .core import history as history_core
from .core import ocr as ocr_core
from .core.font_fit import fit_font_size_pt
from .core import page_manager as page_manager_core
from .features import export as export_feature
from .features import project as project_feature
from .features import inpaint as inpaint_feature
from .features import ai_replace as ai_replace_feature

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class ModernPPTEditor:
    def __init__(self, root):
        self.root = root
        self.root.title("PPTç¼–è¾‘å™¨ä¸“ä¸šç‰ˆ - å¢å¼ºç‰ˆ")
        self.root.geometry("1500x900")
        self.root.configure(bg=COLOR_RIBBON_BG)

        # åŠ è½½é…ç½®
        self.config = load_config()

        # å¤šé¡µæ”¯æŒ
        self.pages = []
        self.current_page_index = 0

        # å½“å‰é¡µæ•°æ®
        self.original_img_path = None
        self.clean_bg_path = None
        self.original_image = None
        self.display_image = None
        self.tk_image = None
        self.scale = 1.0

        # æ–‡æœ¬æ¡†
        self.text_boxes = []
        # å›¾å±‚ï¼ˆæ¯é¡µæŒä¹…åŒ–åœ¨ page["layers"]ï¼Œè¿™é‡Œæ˜¯å½“å‰é¡µå¼•ç”¨ï¼‰
        self.layers = []
        self.selected_layer_index = -1
        self.selected_box_index = -1
        self.selected_boxes = []

        # é¢„è§ˆæ¨¡å¼
        # raw: åªçœ‹åŸå§‹ç¼–è¾‘åº•å›¾ï¼ˆä¸å åŠ å›¾å±‚ï¼‰
        # edit: ç¼–è¾‘è§†å›¾ï¼ˆå åŠ èƒŒæ™¯/å›¾å±‚ + æ¡†ï¼‰
        # ppt: PPTæ•ˆæœï¼ˆå åŠ èƒŒæ™¯/å›¾å±‚ + æ¸²æŸ“æ–‡å­—ï¼‰
        self.current_preview_mode = "raw"
        self.ppt_preview_image = None

        # æ’¤é”€/é‡åš - å…¨å±€å†å²ç³»ç»Ÿ
        self.history = []  # æ ¼å¼: [{"type": "xxx", "data": {...}, "page_index": N}, ...]
        self.history_index = -1
        self.max_history = 50

        # ç»˜åˆ¶çŠ¶æ€
        self.is_drawing = False
        self.draw_start_x = 0
        self.draw_start_y = 0
        self.temp_rect_id = None
        self.is_dragging = False
        self.drag_start_x = 0
        self.drag_start_y = 0
        self.is_resizing = False
        self.resize_handle = None
        self.is_selecting = False  # æ¡†é€‰æ¨¡å¼
        self.select_start_x = 0
        self.select_start_y = 0
        # å›¾å±‚æ‹–åŠ¨ï¼ˆä»å›¾å±‚é¢æ¿é€‰ä¸­åï¼Œåœ¨ç”»å¸ƒä¸Šæ‹–åŠ¨ï¼‰
        self.is_layer_dragging = False
        self._layer_drag_start_canvas = None
        self._layer_drag_origin_xy = None

        # ç»˜åˆ¶æ¨¡å¼
        self.draw_mode = True

        # OCRæ¨¡å‹
        self.ocr = None

        # ç¼©ç•¥å›¾
        self.thumbnail_images = []
        # å¤åˆ¶ç²˜è´´æ”¯æŒ
        self.clipboard_boxes = []

        # æ¶‚æŠ¹æ¨¡å¼ç›¸å…³
        self.inpaint_mode = False  # æ˜¯å¦å¤„äºæ¶‚æŠ¹æ¨¡å¼
        self.inpaint_tool = "brush"  # brush æˆ– rect
        self.inpaint_brush_size = 30  # ç¬”åˆ·å¤§å°
        self.inpaint_mask_layer = None  # PIL Image (Læ¨¡å¼)ï¼Œç™½è‰²=éœ€è¦ä¿®å¤çš„åŒºåŸŸ
        self.inpaint_draw_layer = None  # ImageDrawå¯¹è±¡
        self.inpaint_last_pos = None  # ç¬”åˆ·ä¸Šä¸€ä¸ªä½ç½®
        self.inpaint_rect_start = None  # çŸ©å½¢æ¡†é€‰èµ·å§‹ç‚¹
        self.inpaint_temp_rect_id = None  # ä¸´æ—¶çŸ©å½¢è§†è§‰ID
        self.inpaint_strokes = []  # æ¶‚æŠ¹å†å²è®°å½•ï¼ˆç”¨äºæ’¤é”€ï¼‰

        # AIå›¾ç‰‡æ›¿æ¢æ¨¡å¼ç›¸å…³
        self.ai_replace_mode = False  # æ˜¯å¦å¤„äºAIæ›¿æ¢æ¨¡å¼
        self.ai_replace_rect_start = None  # æ¡†é€‰èµ·å§‹ç‚¹
        self.ai_replace_rect_end = None  # æ¡†é€‰ç»“æŸç‚¹
        self.ai_replace_selection = None  # å½“å‰é€‰ä¸­çš„åŒºåŸŸ (x1, y1, x2, y2)
        self.ai_replace_rect_id = None  # é€‰æ¡†çš„canvas ID

        # AIå›¾ç‰‡APIç®¡ç†å™¨
        self.ai_api_manager = AIImageAPIManager()
        # åŠ è½½AI APIé…ç½®
        if "ai_image_api" in self.config:
            self.ai_api_manager.load_config(self.config)

        # è‡ªåŠ¨ä¿å­˜
        self.autosave_timer = None
        self.project_file_path = None
        self.has_unsaved_changes = False

        # åˆ›å»ºè‡ªåŠ¨ä¿å­˜ç›®å½•
        AUTOSAVE_DIR = os.path.join(get_base_dir(), "autosave")
        os.makedirs(AUTOSAVE_DIR, exist_ok=True)
        self.autosave_dir = AUTOSAVE_DIR


        # åˆ›å»ºç•Œé¢
        self.create_ui()

        # ç»‘å®šå¿«æ·é”®
        self.bind_shortcuts()

        # åå°åŠ è½½ OCRï¼ˆæŒ‰éœ€ä½¿ç”¨ï¼›ä¸é˜»å¡ UIï¼‰
        if self.config.get("ocr_autoload", True):
            threading.Thread(target=self.init_ocr, daemon=True).start()

        # å¯åŠ¨è‡ªåŠ¨ä¿å­˜
        if self.config.get("autosave_enabled", True):
            self.start_autosave()

        # çª—å£å…³é—­äº‹ä»¶
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)


    def _imread_unicode(self, filepath):
        """
        å®‰å…¨è¯»å–åŒ…å«ä¸­æ–‡è·¯å¾„çš„å›¾ç‰‡
        è§£å†³OpenCVæ— æ³•è¯»å–ä¸­æ–‡è·¯å¾„çš„é—®é¢˜
        """
        try:
            # ä½¿ç”¨numpyè¯»å–æ–‡ä»¶ï¼Œç„¶åè§£ç ä¸ºå›¾ç‰‡
            img_array = np.fromfile(filepath, dtype=np.uint8)
            img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
            return img
        except Exception as e:
            print(f"è¯»å–å›¾ç‰‡å¤±è´¥: {filepath}, é”™è¯¯: {e}")
            return None

    def init_ocr(self):
        return ocr_core.init_ocr(self)

    def create_ui(self):
        """åˆ›å»ºç•Œé¢"""
        # === é¡¶éƒ¨æ ‡é¢˜æ  ===
        self.create_title_bar()

        # === å·¥å…·æ  ===
        self.create_toolbar()

        # === ä¸»å†…å®¹åŒº ===
        self.main_container = tk.Frame(self.root, bg=COLOR_CANVAS_BG)
        self.main_container.pack(fill=tk.BOTH, expand=True)

        # å·¦ä¾§ï¼šé¡µé¢ç¼©ç•¥å›¾
        self.create_thumbnail_panel()

        # ä¸­é—´ï¼šä¸»ç¼–è¾‘åŒº
        self.create_canvas_area()

        # å³ä¾§ï¼šå±æ€§é¢æ¿
        self.create_property_panel()

        # === åº•éƒ¨çŠ¶æ€æ  ===
        self.create_status_bar()

    def create_title_bar(self):
        """åˆ›å»ºé¡¶éƒ¨æ ‡é¢˜æ  - PowerPointçº¢è‰²é£æ ¼"""
        title_bar = tk.Frame(self.root, bg=COLOR_THEME, height=32)
        title_bar.pack(fill=tk.X, side=tk.TOP)
        title_bar.pack_propagate(False)

        # å·¦ä¾§æ ‡é¢˜
        title_label = tk.Label(title_bar, text="PPTç¼–è¾‘å™¨ä¸“ä¸šç‰ˆ",
                              bg=COLOR_THEME, fg="white",
                              font=(FONT_FAMILY, 11, "bold"))
        title_label.pack(side=tk.LEFT, padx=15)

        # å³ä¾§é¡µç ä¿¡æ¯
        self.title_page_label = tk.Label(title_bar, text="ç¬¬ 0/0 é¡µ",
                                         bg=COLOR_THEME, fg="white",
                                         font=(FONT_FAMILY, 10))
        self.title_page_label.pack(side=tk.RIGHT, padx=15)

        # è‡ªåŠ¨ä¿å­˜çŠ¶æ€æŒ‡ç¤ºå™¨
        self.autosave_indicator = tk.Label(title_bar, text="â—",
                                          bg=COLOR_THEME, fg="#4CAF50",
                                          font=(FONT_FAMILY, 16))
        self.autosave_indicator.pack(side=tk.RIGHT, padx=5)

    def create_toolbar(self):
        """åˆ›å»ºé¡¶éƒ¨å·¥å…·æ  - ä¸‰è¡Œå¸ƒå±€ï¼ˆé€‚é…å°å±å¹•ï¼‰"""
        toolbar = tk.Frame(self.root, bg=COLOR_RIBBON_BG, relief=tk.FLAT)
        toolbar.pack(fill=tk.X, side=tk.TOP)

        # åº•éƒ¨è¾¹æ¡†çº¿
        border_line = tk.Frame(toolbar, bg="#ddd", height=1)
        border_line.pack(fill=tk.X, side=tk.BOTTOM)

        # === ç¬¬ä¸€è¡Œï¼šæ–‡ä»¶ã€æ£€æµ‹ã€è¯†åˆ« ===
        row1 = tk.Frame(toolbar, bg=COLOR_RIBBON_BG)
        row1.pack(fill=tk.X, padx=10, pady=(6, 2))

        # æ–‡ä»¶ç»„
        tk.Label(row1, text="æ–‡ä»¶:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "å¯¼å…¥å›¾ç‰‡", self.load_multiple_images, COLOR_GREEN)
        self.create_tool_btn(row1, "å¯¼å…¥èƒŒæ™¯", self.load_multiple_backgrounds, COLOR_BLUE)
        self.create_tool_btn(row1, "æ–°å»ºç©ºç™½", self.create_blank_page, "#2196F3")
        if PDF_SUPPORT:
            self.create_tool_btn(row1, "å¯¼å…¥PDF", self.import_pdf, "#D32F2F")

        self.create_tool_btn(row1, "ä¿å­˜é¡¹ç›®", self.save_project, COLOR_GRAY)
        self.create_tool_btn(row1, "æ‰“å¼€é¡¹ç›®", self.load_project, COLOR_GRAY)

        self.create_separator(row1)

        # æ£€æµ‹ç»„
        tk.Label(row1, text="æ£€æµ‹:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "å½“å‰é¡µ", self.auto_detect_text_regions, COLOR_ORANGE)
        self.create_tool_btn(row1, "å…¨éƒ¨é¡µ", self.auto_detect_all_pages, "#EF6C00")

        self.create_separator(row1)

        # è¯†åˆ«ç»„
        tk.Label(row1, text="è¯†åˆ«:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "å½“å‰é¡µ", self.ocr_all_boxes, COLOR_PURPLE)
        self.create_tool_btn(row1, "å…¨éƒ¨é¡µ", self.ocr_all_pages, "#6A1B9A")

        self.create_separator(row1)

        # è‡ªåŠ¨å­—å·ç»„
        tk.Label(row1, text="å­—å·:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row1, "å½“å‰é¡µ", self.auto_font_size_all, "#00ACC1")
        self.create_tool_btn(row1, "å…¨éƒ¨é¡µ", self.auto_font_size_all_pages, "#00838F")

        # å³ä¾§ï¼šå¯¼å‡ºå’Œè®¾ç½®
        settings_btn = tk.Button(row1, text="âš™ è®¾ç½®", command=self.show_settings_dialog,
                                bg="#546E7A", fg="white", font=(FONT_FAMILY, 9),
                                padx=8, pady=2, cursor="hand2", relief=tk.FLAT, bd=0)
        settings_btn.pack(side=tk.RIGHT, padx=5)

        self.create_tool_btn_right(row1, "å¯¼å‡ºå›¾ç‰‡", self.export_as_images, "#F57C00")
        self.create_tool_btn_right(row1, "å¯¼å‡ºPDF", self.export_as_pdf, "#C62828")
        self.create_tool_btn_right(row1, "ç”ŸæˆPPT", self.generate_multi_page_ppt, COLOR_RED)

        # === ç¬¬äºŒè¡Œï¼šæ¶‚æŠ¹ã€AIæ›¿æ¢ã€èƒŒæ™¯ç”Ÿæˆ ===
        row2 = tk.Frame(toolbar, bg=COLOR_RIBBON_ROW2)
        row2.pack(fill=tk.X, padx=10, pady=(2, 2))

        # æ¶‚æŠ¹å·¥å…·ç»„
        tk.Label(row2, text="æ¶‚æŠ¹:", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        # æ¶‚æŠ¹æ¨¡å¼å¼€å…³
        self.inpaint_mode_btn = self.create_tool_btn(row2, "è¿›å…¥æ¶‚æŠ¹", self.toggle_inpaint_mode, "#FF6F00", bg=COLOR_RIBBON_ROW2)

        # å·¥å…·é€‰æ‹©ï¼ˆåˆå§‹éšè—ï¼‰
        self.inpaint_tools_frame = tk.Frame(row2, bg=COLOR_RIBBON_ROW2)
        self.inpaint_tools_frame.pack(side=tk.LEFT)

        self.brush_btn = tk.Button(self.inpaint_tools_frame, text="ç¬”åˆ·",
                                   command=lambda: self.switch_inpaint_tool("brush"),
                                   bg="#FFE0B2", relief=tk.SUNKEN, font=(FONT_FAMILY, 9),
                                   padx=8, pady=3, cursor="hand2")
        self.brush_btn.pack(side=tk.LEFT, padx=2)

        self.rect_btn = tk.Button(self.inpaint_tools_frame, text="æ¡†é€‰",
                                  command=lambda: self.switch_inpaint_tool("rect"),
                                  bg=COLOR_RIBBON_ROW2, relief=tk.RAISED, font=(FONT_FAMILY, 9),
                                  padx=8, pady=3, cursor="hand2")
        self.rect_btn.pack(side=tk.LEFT, padx=2)

        # ç¬”åˆ·å¤§å°ï¼ˆåˆå§‹éšè—ï¼‰
        self.brush_size_frame = tk.Frame(row2, bg=COLOR_RIBBON_ROW2)
        self.brush_size_frame.pack(side=tk.LEFT)

        tk.Label(self.brush_size_frame, text="å¤§å°:", bg=COLOR_RIBBON_ROW2,
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT, padx=3)
        self.brush_size_scale = tk.Scale(self.brush_size_frame, from_=5, to=100,
                                         orient=tk.HORIZONTAL, length=80,
                                         command=self.update_brush_size,
                                         bg=COLOR_RIBBON_ROW2, highlightthickness=0)
        self.brush_size_scale.set(30)
        self.brush_size_scale.pack(side=tk.LEFT)

        # æ¶‚æŠ¹æ“ä½œæŒ‰é’®ï¼ˆåˆå§‹éšè—ï¼‰
        self.inpaint_actions_frame = tk.Frame(row2, bg=COLOR_RIBBON_ROW2)
        self.inpaint_actions_frame.pack(side=tk.LEFT)

        tk.Button(self.inpaint_actions_frame, text="æ¸…ç©º",
                 command=self.clear_inpaint_mask,
                 bg="#FFCDD2", font=(FONT_FAMILY, 9), padx=8, pady=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        tk.Button(self.inpaint_actions_frame, text="ç”Ÿæˆå›¾å±‚",
                 command=self.generate_bg_from_custom_mask,
                 bg="#A5D6A7", font=(FONT_FAMILY, 9, "bold"), padx=12, pady=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # åˆå§‹éšè—å·¥å…·æ 
        self.inpaint_tools_frame.pack_forget()
        self.brush_size_frame.pack_forget()
        self.inpaint_actions_frame.pack_forget()

        self.create_separator(row2, bg=COLOR_RIBBON_ROW2)

        # AIæ›¿æ¢å·¥å…·ç»„
        tk.Label(row2, text="AIæ›¿æ¢:", bg=COLOR_RIBBON_ROW2, fg="#666",
                 font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        # AIæ›¿æ¢æ¨¡å¼å¼€å…³
        self.ai_replace_mode_btn = self.create_tool_btn(
            row2,
            "AIæ›¿æ¢",
            self.toggle_ai_replace_mode,
            "#E91E63",
            bg=COLOR_RIBBON_ROW2,
        )

        # AI æ–‡å­—ç”Ÿå›¾ï¼ˆä¸éœ€è¦æ¡†é€‰ï¼Œç”Ÿæˆåä½œä¸ºå›¾å±‚ï¼‰
        self.create_tool_btn(row2, "æ–‡å­—ç”Ÿå›¾", self.ai_text_to_image_layer, "#7B1FA2", bg=COLOR_RIBBON_ROW2)

        # AIæ•´é¡µç”ŸæˆèƒŒæ™¯ï¼ˆæŠŠå½“å‰é¡µæ•´å›¾å‘ç»™AIç”Ÿæˆï¼Œè¿”å›è®¾ä¸ºèƒŒæ™¯ï¼‰
        self.create_tool_btn(row2, "æ•´é¡µç”Ÿæˆ", self.ai_generate_fullpage_background, "#6A1B9A", bg=COLOR_RIBBON_ROW2)

        # AI APIé…ç½®æŒ‰é’®
        self.create_tool_btn(row2, "APIè®¾ç½®", self.open_ai_api_settings, "#9C27B0", bg=COLOR_RIBBON_ROW2)

        self.create_separator(row2, bg=COLOR_RIBBON_ROW2)

        # IOPaint å»å­—ï¼ˆç»“æœä½œä¸ºå›¾å±‚å åŠ ï¼Œä¸æ›¿æ¢åŸå›¾/èƒŒæ™¯ï¼‰
        tk.Label(row2, text="å»å­—(å±‚):", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row2, "å½“å‰é¡µ", self.auto_generate_background_current, "#E91E63", bg=COLOR_RIBBON_ROW2)
        self.create_tool_btn(row2, "å…¨éƒ¨é¡µ", self.auto_generate_background_all, "#C2185B", bg=COLOR_RIBBON_ROW2)

        self.create_separator(row2, bg=COLOR_RIBBON_ROW2)

        # é¢„è§ˆæ¨¡å¼
        tk.Label(row2, text="é¢„è§ˆ:", bg=COLOR_RIBBON_ROW2, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.preview_mode_var = tk.StringVar(value="raw")
        self.preview_orig_btn = tk.Button(row2, text="åŸå›¾", command=lambda: self.set_preview_mode("raw"),
                                          bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                                          padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.preview_orig_btn.pack(side=tk.LEFT, padx=2)

        self.preview_edit_btn = tk.Button(row2, text="å åŠ ", command=lambda: self.set_preview_mode("edit"),
                                          bg="#757575", fg="white", font=(FONT_FAMILY, 9),
                                          padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.preview_edit_btn.pack(side=tk.LEFT, padx=2)

        self.preview_ppt_btn = tk.Button(row2, text="PPTæ•ˆæœ", command=lambda: self.set_preview_mode("ppt"),
                                         bg="#757575", fg="white", font=(FONT_FAMILY, 9),
                                         padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.preview_ppt_btn.pack(side=tk.LEFT, padx=2)

        # === ç¬¬ä¸‰è¡Œï¼šç¼–è¾‘å·¥å…·å’Œè§†å›¾ ===
        row3 = tk.Frame(toolbar, bg=COLOR_RIBBON_BG)
        row3.pack(fill=tk.X, padx=10, pady=(2, 6))

        # ç¼–è¾‘å·¥å…·
        tk.Label(row3, text="ç¼–è¾‘:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.draw_mode_var = tk.BooleanVar(value=True)
        self.draw_btn = tk.Button(row3, text="ç”»æ¡†æ¨¡å¼", command=self.toggle_draw_mode_btn,
                                  bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9),
                                  padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        self.draw_btn.pack(side=tk.LEFT, padx=2)


        self.create_tool_btn(row3, "å¤åˆ¶", self.copy_boxes, "#009688")
        self.create_tool_btn(row3, "ç²˜è´´", self.paste_boxes, "#00ACC1")
        self.create_tool_btn(row3, "åˆ é™¤æ¡†", self.delete_selected_box, COLOR_RED)
        self.create_tool_btn(row3, "æ¸…ç©ºå…¨éƒ¨", self.clear_all_boxes, "#795548")
        self.create_tool_btn(row3, "æ’¤é”€", self.undo, "#78909C")
        self.create_tool_btn(row3, "é‡åš", self.redo, "#78909C")

        self.create_separator(row3)

        # ç¼©æ”¾æ§åˆ¶
        tk.Label(row3, text="è§†å›¾:", bg=COLOR_RIBBON_BG, fg="#666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.create_tool_btn(row3, "é€‚åº”çª—å£", self.fit_image_to_canvas, "#455A64")
        self.create_tool_btn(row3, "100%", self.zoom_to_100, "#455A64")

        self.zoom_label = tk.Label(row3, text="100%", bg=COLOR_RIBBON_BG, fg="#333",
                                   font=(FONT_FAMILY, 9), padx=10)
        self.zoom_label.pack(side=tk.LEFT)

        # å¿«æ·é”®æç¤º
        tk.Label(row3, text="Ctrl+æ»šè½®ç¼©æ”¾ | åŒå‡»ç¼–è¾‘ | Ctrl+ç‚¹å‡»å¤šé€‰",
                bg=COLOR_RIBBON_BG, fg="#999", font=(FONT_FAMILY, 8)).pack(side=tk.LEFT, padx=10)

    def create_tool_btn(self, parent, text, command, color, bg=None):
        """åˆ›å»ºå·¥å…·æ æŒ‰é’®"""
        if bg is None:
            bg = COLOR_RIBBON_BG
        btn = tk.Button(parent, text=text, command=command,
                       bg=color, fg="white", font=(FONT_FAMILY, 9),
                       padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        btn.pack(side=tk.LEFT, padx=2)
        return btn

    def create_tool_btn_right(self, parent, text, command, color, bg=None):
        """åˆ›å»ºé å³å¯¹é½çš„å·¥å…·æ æŒ‰é’®"""
        if bg is None:
            bg = COLOR_RIBBON_BG
        btn = tk.Button(parent, text=text, command=command,
                       bg=color, fg="white", font=(FONT_FAMILY, 9),
                       padx=8, cursor="hand2", relief=tk.FLAT, bd=0)
        btn.pack(side=tk.RIGHT, padx=2)
        return btn

    def create_separator(self, parent, bg=None):
        """åˆ›å»ºåˆ†éš”çº¿"""
        if bg is None:
            bg = COLOR_RIBBON_BG
        sep_frame = tk.Frame(parent, bg=bg)
        sep_frame.pack(side=tk.LEFT, padx=6)
        sep_line = tk.Frame(sep_frame, bg="#ccc", width=1, height=20)
        sep_line.pack()

    def create_icon_button(self, parent, text, command, color, icon=""):
        """åˆ›å»ºå›¾æ ‡æŒ‰é’®"""
        btn_text = f"{icon}\n{text}" if icon else text
        btn = tk.Button(parent, text=btn_text, command=command,
                       bg=color, fg="white", font=("å¾®è½¯é›…é»‘", 8),
                       width=5, height=2, cursor="hand2", relief=tk.GROOVE, bd=2)
        btn.pack(side=tk.LEFT, padx=2, pady=2)

        # æ‚¬åœæ•ˆæœ
        def on_enter(e):
            btn.config(relief=tk.RAISED)
        def on_leave(e):
            btn.config(relief=tk.GROOVE)
        btn.bind("<Enter>", on_enter)
        btn.bind("<Leave>", on_leave)
        return btn

    def toggle_draw_mode_btn(self):
        """åˆ‡æ¢ç»˜åˆ¶æ¨¡å¼"""
        self.draw_mode = not self.draw_mode
        self.draw_mode_var.set(self.draw_mode)
        if self.draw_mode:
            self.draw_btn.config(bg=COLOR_GREEN, text="ç”»æ¡†æ¨¡å¼")
            self.canvas.config(cursor="crosshair")
        else:
            self.draw_btn.config(bg="#9E9E9E", text="é€‰æ‹©æ¨¡å¼")
            self.canvas.config(cursor="")

    def set_preview_mode(self, mode):
        """è®¾ç½®é¢„è§ˆæ¨¡å¼"""
        self.preview_mode_var.set(mode)
        self.current_preview_mode = mode
        # é¢œè‰²çŠ¶æ€
        self.preview_orig_btn.config(bg=COLOR_BLUE if mode == "raw" else "#757575", fg="white")
        if hasattr(self, "preview_edit_btn"):
            self.preview_edit_btn.config(bg=COLOR_BLUE if mode == "edit" else "#757575", fg="white")
        self.preview_ppt_btn.config(bg=COLOR_BLUE if mode == "ppt" else "#757575", fg="white")
        self.refresh_canvas()

    def create_thumbnail_panel(self):
        """åˆ›å»ºå·¦ä¾§ç¼©ç•¥å›¾é¢æ¿"""
        self.thumbnail_panel = tk.Frame(self.main_container, bg=COLOR_SIDEBAR_BG, width=180)
        self.thumbnail_panel.pack(side=tk.LEFT, fill=tk.Y)
        self.thumbnail_panel.pack_propagate(False)

        # æ ‡é¢˜æ 
        title_frame = tk.Frame(self.thumbnail_panel, bg=COLOR_BLUE, height=30)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  é¡µé¢åˆ—è¡¨", bg=COLOR_BLUE, fg="white",
                font=(FONT_FAMILY, 10, "bold"), anchor="w").pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # ç¼©ç•¥å›¾å®¹å™¨ï¼ˆå¯æ»šåŠ¨ï¼‰
        container = tk.Frame(self.thumbnail_panel, bg=COLOR_SIDEBAR_BG)
        container.pack(fill=tk.BOTH, expand=True)

        self.thumbnail_canvas = tk.Canvas(container, bg=COLOR_SIDEBAR_BG, highlightthickness=0, width=160)
        scrollbar = tk.Scrollbar(container, orient=tk.VERTICAL, command=self.thumbnail_canvas.yview)

        self.thumbnail_frame = tk.Frame(self.thumbnail_canvas, bg=COLOR_SIDEBAR_BG)

        self.thumbnail_canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.thumbnail_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        self.thumbnail_window = self.thumbnail_canvas.create_window((0, 0), window=self.thumbnail_frame, anchor=tk.NW)

        self.thumbnail_frame.bind("<Configure>",
            lambda e: self.thumbnail_canvas.configure(scrollregion=self.thumbnail_canvas.bbox("all")))

        # é¼ æ ‡æ»šè½®
        self.thumbnail_canvas.bind("<MouseWheel>",
            lambda e: self.thumbnail_canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        # é¡µé¢å¯¼èˆªæŒ‰é’®
        nav_frame = tk.Frame(self.thumbnail_panel, bg="#f5f5f5", height=40)
        nav_frame.pack(fill=tk.X, side=tk.BOTTOM)
        nav_frame.pack_propagate(False)

        tk.Button(nav_frame, text="ä¸Šä¸€é¡µ", command=self.prev_page,
                 bg="#e0e0e0", font=(FONT_FAMILY, 9), width=6, cursor="hand2",
                 relief=tk.FLAT).pack(side=tk.LEFT, padx=5, pady=5)

        self.page_label = tk.Label(nav_frame, text="0/0", bg="#f5f5f5",
                                   font=(FONT_FAMILY, 10, "bold"))
        self.page_label.pack(side=tk.LEFT, expand=True)

        tk.Button(nav_frame, text="ä¸‹ä¸€é¡µ", command=self.next_page,
                 bg="#e0e0e0", font=(FONT_FAMILY, 9), width=6, cursor="hand2",
                 relief=tk.FLAT).pack(side=tk.RIGHT, padx=5, pady=5)

    def create_canvas_area(self):
        """åˆ›å»ºä¸­é—´ç”»å¸ƒåŒºåŸŸ"""
        canvas_container = tk.Frame(self.main_container, bg=COLOR_CANVAS_BG)
        canvas_container.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # ç”»å¸ƒ
        self.canvas = tk.Canvas(canvas_container, bg="#c0c0c0", highlightthickness=0)

        # æ»šåŠ¨æ¡
        v_scroll = tk.Scrollbar(canvas_container, orient=tk.VERTICAL, command=self.canvas.yview)
        h_scroll = tk.Scrollbar(canvas_container, orient=tk.HORIZONTAL, command=self.canvas.xview)

        self.canvas.config(xscrollcommand=h_scroll.set, yscrollcommand=v_scroll.set)

        v_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        h_scroll.pack(side=tk.BOTTOM, fill=tk.X)
        self.canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # ç»‘å®šç”»å¸ƒäº‹ä»¶
        self.canvas.bind("<ButtonPress-1>", self.on_canvas_press)
        self.canvas.bind("<Control-ButtonPress-1>", self.on_canvas_ctrl_click)
        self.canvas.bind("<B1-Motion>", self.on_canvas_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_canvas_release)
        self.canvas.bind("<Configure>", self.on_canvas_resize)
        self.canvas.bind("<Double-Button-1>", self.on_canvas_double_click)
        self.canvas.bind("<Button-3>", self.on_canvas_right_click)  # å³é”®èœå•

        # Ctrl+æ»šè½®ç¼©æ”¾
        self.canvas.bind("<Control-MouseWheel>", self.on_canvas_zoom)
        # æ™®é€šæ»šè½®æ»šåŠ¨
        self.canvas.bind("<MouseWheel>", self.on_canvas_scroll)

        # å ä½æç¤º
        self.placeholder_label = tk.Label(self.canvas,
            text="ç‚¹å‡»ä¸Šæ–¹ã€Œå¯¼å…¥å›¾ç‰‡ã€æŒ‰é’®å¼€å§‹\n\næ”¯æŒæ‰¹é‡å¯¼å…¥å¤šå¼ å›¾ç‰‡",
            bg="#c0c0c0", fg="#666666", font=(FONT_FAMILY, 14), justify=tk.CENTER)
        self.canvas.create_window(400, 300, window=self.placeholder_label)

    def create_property_panel(self):
        """åˆ›å»ºå³ä¾§å±æ€§é¢æ¿"""
        self.right_panel = tk.Frame(self.main_container, bg=COLOR_WHITE, width=280)
        self.right_panel.pack(side=tk.RIGHT, fill=tk.Y)
        self.right_panel.pack_propagate(False)

        # æ ‡é¢˜
        title_frame = tk.Frame(self.right_panel, bg=COLOR_BLUE, height=30)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  å±æ€§è®¾ç½®", bg=COLOR_BLUE, fg="white",
                font=(FONT_FAMILY, 10, "bold"), anchor="w").pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # åˆ†é¡µï¼šå±æ€§ / å›¾å±‚ï¼ˆæ›´æ¥è¿‘ PS çš„é¢æ¿ä½“éªŒï¼‰
        self.right_notebook = ttk.Notebook(self.right_panel)
        self.right_notebook.pack(side=tk.TOP, fill=tk.BOTH, expand=True)

        props_tab = tk.Frame(self.right_notebook, bg=COLOR_WHITE)
        layers_tab = tk.Frame(self.right_notebook, bg=COLOR_WHITE)
        self.right_notebook.add(props_tab, text="å±æ€§")
        self.right_notebook.add(layers_tab, text="å›¾å±‚")
        self.layers_tab = layers_tab

        # å¯æ»šåŠ¨å®¹å™¨ï¼ˆå±æ€§é¡µï¼‰
        canvas = tk.Canvas(props_tab, bg=COLOR_WHITE, highlightthickness=0)
        self.prop_canvas = canvas
        scrollbar = tk.Scrollbar(props_tab, orient=tk.VERTICAL, command=canvas.yview)

        self.prop_frame = tk.Frame(canvas, bg=COLOR_WHITE)

        canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        canvas_window = canvas.create_window((0, 0), window=self.prop_frame, anchor=tk.NW)

        self.prop_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.bind("<Configure>", lambda e: canvas.itemconfig(canvas_window, width=e.width))
        canvas.bind("<MouseWheel>", lambda e: canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        # === æ–‡æœ¬æ¡†åˆ—è¡¨ ===
        self.create_section_header(self.prop_frame, "æ–‡æœ¬æ¡†åˆ—è¡¨")

        list_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        list_frame.pack(fill=tk.X, padx=10, pady=5)

        self.box_listbox = tk.Listbox(list_frame, height=5, bg="#f5f5f5",
                                       font=(FONT_FAMILY, 9), selectbackground=COLOR_BLUE,
                                       selectforeground="white", relief=tk.FLAT, bd=1)
        self.box_listbox.pack(fill=tk.X)
        self.box_listbox.bind("<<ListboxSelect>>", self.on_listbox_select)

        # === æ–‡æœ¬å†…å®¹ ===
        self.create_section_header(self.prop_frame, "æ–‡æœ¬å†…å®¹")

        text_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        text_frame.pack(fill=tk.X, padx=10, pady=5)

        self.text_entry = tk.Text(text_frame, height=3, bg="#f5f5f5",
                                  font=(FONT_FAMILY, 10), relief=tk.FLAT, bd=1, wrap=tk.WORD)
        self.text_entry.pack(fill=tk.X)
        self.text_entry.bind("<KeyRelease>", self.on_text_change)

        # OCRè¯†åˆ«æŒ‰é’®
        ocr_btn_frame = tk.Frame(text_frame, bg=COLOR_WHITE)
        ocr_btn_frame.pack(fill=tk.X, pady=5)

        tk.Button(ocr_btn_frame, text="ğŸ” OCRè¯†åˆ«æ­¤æ¡†", command=self.ocr_single_box,
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9, "bold"),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X)

        # === ä½ç½®å’Œå¤§å° ===
        self.create_section_header(self.prop_frame, "ä½ç½®å’Œå¤§å°")

        pos_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        pos_frame.pack(fill=tk.X, padx=10, pady=5)

        # X, Y
        row1 = tk.Frame(pos_frame, bg=COLOR_WHITE)
        row1.pack(fill=tk.X, pady=2)

        tk.Label(row1, text="X:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT)
        self.x_entry = tk.Entry(row1, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.x_entry.pack(side=tk.LEFT, padx=2)
        self.x_entry.bind("<KeyRelease>", self.on_position_change)

        tk.Label(row1, text="Y:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT, padx=(10, 0))
        self.y_entry = tk.Entry(row1, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.y_entry.pack(side=tk.LEFT, padx=2)
        self.y_entry.bind("<KeyRelease>", self.on_position_change)

        # å®½, é«˜
        row2 = tk.Frame(pos_frame, bg=COLOR_WHITE)
        row2.pack(fill=tk.X, pady=2)

        tk.Label(row2, text="å®½:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT)
        self.w_entry = tk.Entry(row2, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.w_entry.pack(side=tk.LEFT, padx=2)
        self.w_entry.bind("<KeyRelease>", self.on_position_change)

        tk.Label(row2, text="é«˜:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9), width=3).pack(side=tk.LEFT, padx=(10, 0))
        self.h_entry = tk.Entry(row2, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.h_entry.pack(side=tk.LEFT, padx=2)
        self.h_entry.bind("<KeyRelease>", self.on_position_change)

        # === å­—ä½“æ ·å¼ ===
        self.create_section_header(self.prop_frame, "å­—ä½“æ ·å¼")

        font_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        font_frame.pack(fill=tk.X, padx=10, pady=5)

        # å­—ä½“å’Œå­—å·
        row3 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row3.pack(fill=tk.X, pady=2)

        self.fontname_var = tk.StringVar(value="å¾®è½¯é›…é»‘")
        font_combo = ttk.Combobox(row3, textvariable=self.fontname_var, width=10,
                                  values=["å¾®è½¯é›…é»‘", "å®‹ä½“", "é»‘ä½“", "æ¥·ä½“", "ä»¿å®‹", "Arial"])
        font_combo.pack(side=tk.LEFT, padx=2)
        font_combo.bind("<<ComboboxSelected>>", self.on_font_change)

        self.fontsize_var = tk.StringVar(value="16")
        size_combo = ttk.Combobox(row3, textvariable=self.fontsize_var, width=5,
                                  values=["8", "10", "12", "14", "16", "18", "20", "24", "28", "32", "36", "48", "60", "72", "80", "100", "120", "150", "200"])
        size_combo.pack(side=tk.LEFT, padx=2)
        size_combo.bind("<<ComboboxSelected>>", self.on_font_change)

        # æ ·å¼æŒ‰é’®
        row4 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row4.pack(fill=tk.X, pady=5)

        self.bold_var = tk.BooleanVar(value=False)
        self.bold_btn = tk.Button(row4, text="B åŠ ç²—", command=self.toggle_bold,
                                  bg="#e0e0e0", font=(FONT_FAMILY, 9),
                                  width=6, cursor="hand2", relief=tk.FLAT)
        self.bold_btn.pack(side=tk.LEFT, padx=2)

        self.italic_var = tk.BooleanVar(value=False)
        self.italic_btn = tk.Button(row4, text="I æ–œä½“", command=self.toggle_italic,
                                    bg="#e0e0e0", font=(FONT_FAMILY, 9),
                                    width=6, cursor="hand2", relief=tk.FLAT)
        self.italic_btn.pack(side=tk.LEFT, padx=2)

        self.color_btn = tk.Button(row4, text="é¢œè‰²", command=self.choose_color,
                                   bg="#000000", fg="white", width=5, cursor="hand2", relief=tk.FLAT)
        self.color_btn.pack(side=tk.LEFT, padx=2)

        # è‡ªåŠ¨å­—å·æŒ‰é’®
        tk.Button(row4, text="è‡ªåŠ¨å­—å·", command=self.auto_font_size,
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=5)

        # å¯¹é½æŒ‰é’®
        row5 = tk.Frame(font_frame, bg=COLOR_WHITE)
        row5.pack(fill=tk.X, pady=5)

        tk.Label(row5, text="å¯¹é½:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.align_var = tk.StringVar(value="left")

        align_btn_frame = tk.Frame(row5, bg=COLOR_WHITE)
        align_btn_frame.pack(side=tk.LEFT, padx=5)

        self.align_left_btn = tk.Button(align_btn_frame, text="å·¦", command=lambda: self.set_align("left"),
                                        bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=3,
                                        cursor="hand2", relief=tk.FLAT)
        self.align_left_btn.pack(side=tk.LEFT, padx=1)

        self.align_center_btn = tk.Button(align_btn_frame, text="ä¸­", command=lambda: self.set_align("center"),
                                          bg="#e0e0e0", fg="#333", font=(FONT_FAMILY, 9), width=3,
                                          cursor="hand2", relief=tk.FLAT)
        self.align_center_btn.pack(side=tk.LEFT, padx=1)

        self.align_right_btn = tk.Button(align_btn_frame, text="å³", command=lambda: self.set_align("right"),
                                         bg="#e0e0e0", fg="#333", font=(FONT_FAMILY, 9), width=3,
                                         cursor="hand2", relief=tk.FLAT)
        self.align_right_btn.pack(side=tk.LEFT, padx=1)

        # === æ‰¹é‡åº”ç”¨ ===
        self.create_section_header(self.prop_frame, "æ‰¹é‡åº”ç”¨")

        batch_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        batch_frame.pack(fill=tk.X, padx=10, pady=5)

        tk.Label(batch_frame, text="Ctrl+ç‚¹å‡»å¤šé€‰ï¼Œå‹¾é€‰è¦åº”ç”¨çš„å±æ€§ï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8), wraplength=220).pack(anchor="w")

        # å‹¾é€‰é¡¹
        check_row1 = tk.Frame(batch_frame, bg=COLOR_WHITE)
        check_row1.pack(fill=tk.X, pady=2)

        self.apply_fontsize_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="å­—å·", variable=self.apply_fontsize_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_fontname_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="å­—ä½“", variable=self.apply_fontname_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_color_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row1, text="é¢œè‰²", variable=self.apply_color_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        check_row2 = tk.Frame(batch_frame, bg=COLOR_WHITE)
        check_row2.pack(fill=tk.X, pady=2)

        self.apply_bold_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="åŠ ç²—", variable=self.apply_bold_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_italic_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="æ–œä½“", variable=self.apply_italic_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        self.apply_align_var = tk.BooleanVar(value=False)
        tk.Checkbutton(check_row2, text="å¯¹é½", variable=self.apply_align_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        tk.Button(batch_frame, text="åº”ç”¨åˆ°é€‰ä¸­æ¡†", command=self.apply_style_to_selected,
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X, pady=5)

        # === å¯¹é½å·¥å…· ===
        self.create_section_header(self.prop_frame, "å¤šæ¡†å¯¹é½")

        align_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        align_frame.pack(fill=tk.X, padx=10, pady=5)

        # å…¨é€‰æŒ‰é’®
        select_all_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        select_all_frame.pack(fill=tk.X, pady=(0, 5))

        tk.Button(select_all_frame, text="å…¨é€‰å½“å‰é¡µæ‰€æœ‰æ¡† (Ctrl+A)", command=self.select_all_boxes,
                 bg="#FF9800", fg="white", font=(FONT_FAMILY, 9, "bold"),
                 cursor="hand2", relief=tk.FLAT).pack(fill=tk.X)

        tk.Label(align_frame, text="Ctrl+ç‚¹å‡»é€‰ä¸­å¤šä¸ªæ¡†ï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w", pady=(5, 0))

        # æ°´å¹³å¯¹é½
        h_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        h_align_frame.pack(fill=tk.X, pady=3)

        tk.Label(h_align_frame, text="æ°´å¹³:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        tk.Button(h_align_frame, text="å·¦", command=lambda: self.align_boxes("left"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_align_frame, text="ä¸­", command=lambda: self.align_boxes("center_h"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_align_frame, text="å³", command=lambda: self.align_boxes("right"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # å‚ç›´å¯¹é½
        v_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        v_align_frame.pack(fill=tk.X, pady=3)

        tk.Label(v_align_frame, text="å‚ç›´:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        tk.Button(v_align_frame, text="ä¸Š", command=lambda: self.align_boxes("top"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_align_frame, text="ä¸­", command=lambda: self.align_boxes("center_v"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_align_frame, text="ä¸‹", command=lambda: self.align_boxes("bottom"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 8), width=4,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # åˆ†éš”çº¿
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # å‡åŒ€åˆ†å¸ƒ
        tk.Label(align_frame, text="å‡åŒ€åˆ†å¸ƒï¼ˆéœ€è¦3ä¸ªæˆ–ä»¥ä¸Šï¼‰ï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        dist_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        dist_frame.pack(fill=tk.X, pady=3)

        tk.Button(dist_frame, text="æ°´å¹³ç­‰é—´è·", command=lambda: self.distribute_boxes("horizontal"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(dist_frame, text="å‚ç›´ç­‰é—´è·", command=lambda: self.distribute_boxes("vertical"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # åˆ†éš”çº¿
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # å°ºå¯¸ç»Ÿä¸€
        tk.Label(align_frame, text="å°ºå¯¸ç»Ÿä¸€ï¼ˆä»¥ç¬¬ä¸€ä¸ªé€‰ä¸­æ¡†ä¸ºåŸºå‡†ï¼‰ï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        size_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        size_frame.pack(fill=tk.X, pady=3)

        tk.Button(size_frame, text="ç»Ÿä¸€å®½", command=lambda: self.unify_size("width"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=7,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="ç»Ÿä¸€é«˜", command=lambda: self.unify_size("height"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=7,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="ç»Ÿä¸€å¤§å°", command=lambda: self.unify_size("both"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 8), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # åˆ†éš”çº¿
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        # å¯¹é½åˆ°ç”»å¸ƒ
        tk.Label(align_frame, text="å¯¹é½åˆ°ç”»å¸ƒä¸­å¿ƒï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        canvas_align_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        canvas_align_frame.pack(fill=tk.X, pady=3)

        tk.Button(canvas_align_frame, text="æ°´å¹³å±…ä¸­", command=lambda: self.align_to_canvas("h"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_align_frame, text="å‚ç›´å±…ä¸­", command=lambda: self.align_to_canvas("v"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_align_frame, text="å®Œå…¨å±…ä¸­", command=lambda: self.align_to_canvas("center"),
                 bg="#D32F2F", fg="white", font=(FONT_FAMILY, 8), width=9,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # æ‰¹é‡ä½ç§»
        tk.Frame(align_frame, bg="#e0e0e0", height=1).pack(fill=tk.X, pady=8)

        tk.Label(align_frame, text="æ‰¹é‡ä½ç§»ï¼ˆåƒç´ ï¼‰ï¼š",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")

        # ä½ç§»è¾“å…¥æ¡†
        offset_input_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        offset_input_frame.pack(fill=tk.X, pady=3)

        tk.Label(offset_input_frame, text="ç§»åŠ¨:", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        self.offset_px_var = tk.StringVar(value="10")
        offset_entry = tk.Entry(offset_input_frame, textvariable=self.offset_px_var,
                               width=5, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        offset_entry.pack(side=tk.LEFT, padx=3)

        tk.Label(offset_input_frame, text="px", bg=COLOR_WHITE, font=(FONT_FAMILY, 8), fg="#666").pack(side=tk.LEFT)

        # æ–¹å‘æŒ‰é’®
        offset_btn_frame = tk.Frame(align_frame, bg=COLOR_WHITE)
        offset_btn_frame.pack(fill=tk.X, pady=3)

        # ä¸ŠæŒ‰é’®
        tk.Button(offset_btn_frame, text="â†‘", command=lambda: self.batch_offset(0, -1),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # ä¸‹æŒ‰é’®
        tk.Button(offset_btn_frame, text="â†“", command=lambda: self.batch_offset(0, 1),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # å·¦æŒ‰é’®
        tk.Button(offset_btn_frame, text="â†", command=lambda: self.batch_offset(-1, 0),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # å³æŒ‰é’®
        tk.Button(offset_btn_frame, text="â†’", command=lambda: self.batch_offset(1, 0),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"), width=3,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=1)

        # === å½“å‰é¡µèƒŒæ™¯ ===
        self.create_section_header(self.prop_frame, "å½“å‰é¡µèƒŒæ™¯")

        bg_frame = tk.Frame(self.prop_frame, bg=COLOR_WHITE)
        bg_frame.pack(fill=tk.X, padx=10, pady=5)

        tk.Label(bg_frame, text="èƒŒæ™¯å›¾ä¼šè‡ªåŠ¨è°ƒæ•´ä¸ºä¸åŸå›¾ç›¸åŒå¤§å°",
                bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8), wraplength=220).pack(anchor="w")

        bg_btn_frame = tk.Frame(bg_frame, bg=COLOR_WHITE)
        bg_btn_frame.pack(fill=tk.X, pady=5)

        tk.Button(bg_btn_frame, text="è®¾ç½®èƒŒæ™¯", command=self.load_current_page_background,
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        tk.Button(bg_btn_frame, text="æ¸…é™¤èƒŒæ™¯", command=self.clear_current_page_background,
                 bg=COLOR_RED, fg="white", font=(FONT_FAMILY, 9),
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # èƒŒæ™¯çŠ¶æ€æ˜¾ç¤º
        self.bg_status_label = tk.Label(bg_frame, text="æœªè®¾ç½®èƒŒæ™¯",
                                        bg=COLOR_WHITE, fg="#999", font=(FONT_FAMILY, 8))
        self.bg_status_label.pack(anchor="w", pady=2)

        # å›¾å±‚é¢æ¿å·²ç§»åŠ¨åˆ°å³ä¾§â€œå›¾å±‚â€Tabï¼ˆæ›´æ˜“æ‰¾åˆ°ï¼Œä¹Ÿæ›´åƒ PSï¼‰ã€‚

        # å›¾å±‚é¡µ
        self.create_layers_panel(layers_tab)
        self.update_layer_listbox()

    def create_layers_panel(self, parent):
        """åˆ›å»ºå›¾å±‚é¢æ¿ï¼ˆç‹¬ç«‹ Tabï¼‰"""
        # === å›¾å±‚ ===
        header = tk.Frame(parent, bg=COLOR_WHITE)
        header.pack(fill=tk.X, padx=10, pady=(10, 6))
        tk.Label(header, text="å›¾å±‚", bg=COLOR_WHITE, fg=COLOR_TEXT, font=(FONT_FAMILY, 10, "bold")).pack(
            side=tk.LEFT
        )
        tk.Button(
            header,
            text="åˆ·æ–°",
            command=lambda: (self.update_layer_listbox(), self.refresh_canvas()),
            bg="#455A64",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.RIGHT)

        layer_frame = tk.Frame(parent, bg=COLOR_WHITE)
        layer_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=(0, 10))

        tree_frame = tk.Frame(layer_frame, bg=COLOR_WHITE)
        tree_frame.pack(fill=tk.BOTH, expand=True)

        columns = ("vis", "opacity", "lock")
        self.layer_tree = ttk.Treeview(
            tree_frame,
            columns=columns,
            show="tree headings",
            height=12,
            selectmode="browse",
        )
        self.layer_tree.heading("#0", text="å›¾å±‚")
        self.layer_tree.heading("vis", text="æ˜¾")
        self.layer_tree.heading("opacity", text="é€æ˜")
        self.layer_tree.heading("lock", text="é”")

        self.layer_tree.column("#0", width=160, anchor=tk.W)
        self.layer_tree.column("vis", width=38, anchor=tk.CENTER)
        self.layer_tree.column("opacity", width=60, anchor=tk.CENTER)
        self.layer_tree.column("lock", width=45, anchor=tk.CENTER)

        self.layer_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self.layer_tree.bind("<<TreeviewSelect>>", self.on_layer_select)
        self.layer_tree.bind("<Double-1>", lambda e: self.rename_selected_layer())
        # å›¾å±‚é¢æ¿äº¤äº’ï¼šç‚¹å‡»â€œæ˜¾â€åˆ—å¿«é€Ÿæ˜¾ç¤º/éšè—ï¼Œæ‹–æ‹½è¡Œè°ƒæ•´å›¾å±‚é¡ºåº
        self.layer_tree.bind("<Button-1>", self.on_layer_tree_click, add=True)
        self.layer_tree.bind("<ButtonPress-1>", self.on_layer_drag_start, add=True)
        self.layer_tree.bind("<B1-Motion>", self.on_layer_drag_motion, add=True)
        self.layer_tree.bind("<ButtonRelease-1>", self.on_layer_drag_release, add=True)

        tree_scroll = ttk.Scrollbar(tree_frame, orient="vertical", command=self.layer_tree.yview)
        self.layer_tree.configure(yscrollcommand=tree_scroll.set)
        tree_scroll.pack(side=tk.RIGHT, fill=tk.Y)

        layer_btn_frame = tk.Frame(layer_frame, bg=COLOR_WHITE)
        layer_btn_frame.pack(fill=tk.X, pady=(8, 6))

        # å…ˆæ”¾å³ä¾§æŒ‰é’®ï¼Œé¿å…è¢«å·¦ä¾§æŒ‰é’®æŒ¤å‡ºå¯è§†åŒºåŸŸ
        tk.Button(
            layer_btn_frame,
            text="åˆ é™¤",
            command=self.delete_selected_layer,
            bg=COLOR_RED,
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.RIGHT, padx=2)

        tk.Button(
            layer_btn_frame,
            text="å¯¼å…¥å›¾å±‚",
            command=self.import_layer_from_file,
            bg=COLOR_BLUE,
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.RIGHT, padx=2)

        tk.Button(
            layer_btn_frame,
            text="æ˜¾éš",
            command=self.toggle_selected_layer,
            bg="#607D8B",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.LEFT, padx=2)
        tk.Button(
            layer_btn_frame,
            text="é¢„è§ˆ",
            command=self.preview_selected_layer,
            bg="#455A64",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.LEFT, padx=2)
        tk.Button(
            layer_btn_frame,
            text="æ”¹å",
            command=self.rename_selected_layer,
            bg="#455A64",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.LEFT, padx=2)

        tk.Button(
            layer_btn_frame,
            text="é”å®š",
            command=self.toggle_selected_layer_lock,
            bg="#6D4C41",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        ).pack(side=tk.LEFT, padx=2)

        # å›¾å±‚é¡ºåºè°ƒæ•´ï¼ˆç‹¬ç«‹è¡Œï¼Œé¿å…è¢«æŒ¤å‡ºï¼‰
        layer_order_frame = tk.Frame(layer_frame, bg=COLOR_WHITE)
        layer_order_frame.pack(fill=tk.X, pady=(4, 6))

        tk.Label(layer_order_frame, text="å›¾å±‚é¡ºåº:", bg=COLOR_WHITE, fg="#666666",
                font=(FONT_FAMILY, 9, "bold")).pack(side=tk.LEFT, padx=(0, 5))

        tk.Button(
            layer_order_frame,
            text="â†‘ ä¸Šç§»",
            command=self.move_layer_up,
            bg="#607D8B",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
            width=8
        ).pack(side=tk.LEFT, padx=2)

        tk.Button(
            layer_order_frame,
            text="â†“ ä¸‹ç§»",
            command=self.move_layer_down,
            bg="#607D8B",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
            width=8
        ).pack(side=tk.LEFT, padx=2)

        tk.Label(layer_order_frame, text="(è°ƒæ•´å›¾å±‚å æ”¾é¡ºåº)", bg=COLOR_WHITE, fg="#999",
                font=(FONT_FAMILY, 8)).pack(side=tk.LEFT, padx=5)

        # å›¾å±‚å˜æ¢ï¼ˆä½ç½®/ç¼©æ”¾/è£å‰ª/é”å®šï¼‰
        transform_frame = tk.Frame(layer_frame, bg=COLOR_WHITE)
        transform_frame.pack(fill=tk.X, pady=(2, 6))

        self.layer_x_var = tk.IntVar(value=0)
        self.layer_y_var = tk.IntVar(value=0)
        self.layer_scale_var = tk.IntVar(value=100)  # %
        self.layer_lock_var = tk.BooleanVar(value=False)

        xy_row = tk.Frame(transform_frame, bg=COLOR_WHITE)
        xy_row.pack(fill=tk.X, pady=(0, 4))
        tk.Label(xy_row, text="X:", bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.layer_x_entry = tk.Entry(xy_row, textvariable=self.layer_x_var, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.layer_x_entry.pack(side=tk.LEFT, padx=(2, 8))
        tk.Label(xy_row, text="Y:", bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)
        self.layer_y_entry = tk.Entry(xy_row, textvariable=self.layer_y_var, width=6, font=(FONT_FAMILY, 9), relief=tk.FLAT, bg="#f5f5f5")
        self.layer_y_entry.pack(side=tk.LEFT, padx=(2, 8))
        self.layer_lock_check = tk.Checkbutton(xy_row, text="é”å®š", variable=self.layer_lock_var, bg=COLOR_WHITE, font=(FONT_FAMILY, 9), command=self._on_layer_lock_toggle)
        self.layer_lock_check.pack(side=tk.LEFT)

        self.layer_x_entry.bind("<Return>", lambda e: self._apply_layer_transform_from_ui())
        self.layer_y_entry.bind("<Return>", lambda e: self._apply_layer_transform_from_ui())

        tk.Label(transform_frame, text="ç¼©æ”¾(%):", bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")
        self.layer_scale_slider = tk.Scale(
            transform_frame,
            from_=10,
            to=300,
            orient=tk.HORIZONTAL,
            length=240,
            variable=self.layer_scale_var,
            bg=COLOR_WHITE,
            highlightthickness=0,
            command=lambda v: self._on_layer_scale_change(v),
        )
        self.layer_scale_slider.pack(fill=tk.X)
        self.layer_scale_slider.bind("<ButtonPress-1>", self._begin_layer_scale_drag, add=True)
        self.layer_scale_slider.bind("<ButtonRelease-1>", self._end_layer_scale_drag, add=True)

        crop_row = tk.Frame(transform_frame, bg=COLOR_WHITE)
        crop_row.pack(fill=tk.X, pady=(4, 0))
        self.layer_crop_btn = tk.Button(crop_row, text="è£å‰ª", command=self.crop_selected_layer, bg=COLOR_ORANGE, fg="white",
                                        font=(FONT_FAMILY, 9), cursor="hand2", relief=tk.FLAT)
        self.layer_crop_btn.pack(side=tk.LEFT, padx=2)
        self.layer_reset_crop_btn = tk.Button(crop_row, text="é‡ç½®è£å‰ª", command=self.reset_selected_layer_crop, bg="#607D8B", fg="white",
                                              font=(FONT_FAMILY, 9), cursor="hand2", relief=tk.FLAT)
        self.layer_reset_crop_btn.pack(side=tk.LEFT, padx=2)
        self.layer_cutout_btn = tk.Button(
            crop_row,
            text="çº¯è‰²æŠ å›¾",
            command=self.solid_color_cutout_selected_layer,
            bg="#00897B",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        )
        self.layer_cutout_btn.pack(side=tk.LEFT, padx=2)

        # OCRæ£€æµ‹è¯†åˆ«æŒ‰é’®
        ocr_row = tk.Frame(transform_frame, bg=COLOR_WHITE)
        ocr_row.pack(fill=tk.X, pady=(4, 0))
        tk.Label(ocr_row, text="OCR:", bg=COLOR_WHITE, fg="#666666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT, padx=(0, 2))

        self.layer_ocr_detect_btn = tk.Button(
            ocr_row,
            text="æ£€æµ‹",
            command=self.detect_text_in_selected_layers,
            bg="#FF9800",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        )
        self.layer_ocr_detect_btn.pack(side=tk.LEFT, padx=2)

        self.layer_ocr_recognize_btn = tk.Button(
            ocr_row,
            text="è¯†åˆ«",
            command=self.recognize_text_in_selected_layers,
            bg="#FF6F00",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        )
        self.layer_ocr_recognize_btn.pack(side=tk.LEFT, padx=2)

        # å»é™¤æ–‡æœ¬èƒŒæ™¯æŒ‰é’®
        remove_bg_row = tk.Frame(transform_frame, bg=COLOR_WHITE)
        remove_bg_row.pack(fill=tk.X, pady=(4, 0))
        tk.Label(remove_bg_row, text="å»å­—:", bg=COLOR_WHITE, fg="#666666",
                font=(FONT_FAMILY, 9)).pack(side=tk.LEFT, padx=(0, 2))

        self.layer_remove_text_bg_btn = tk.Button(
            remove_bg_row,
            text="å»é™¤æ–‡æœ¬èƒŒæ™¯",
            command=self.remove_text_background_from_layer,
            bg="#E91E63",
            fg="white",
            font=(FONT_FAMILY, 9),
            cursor="hand2",
            relief=tk.FLAT,
        )
        self.layer_remove_text_bg_btn.pack(side=tk.LEFT, padx=2)

        tk.Label(layer_frame, text="é€æ˜åº¦:", bg=COLOR_WHITE, fg="#666666", font=(FONT_FAMILY, 8)).pack(anchor="w")
        self.layer_opacity_scale = tk.Scale(
            layer_frame,
            from_=0,
            to=100,
            orient=tk.HORIZONTAL,
            length=240,
            command=self.on_layer_opacity_change,
            bg=COLOR_WHITE,
            highlightthickness=0,
        )
        self.layer_opacity_scale.set(100)
        self.layer_opacity_scale.pack(fill=tk.X, pady=(0, 5))
        self.layer_opacity_scale.bind("<ButtonPress-1>", self._begin_layer_opacity_drag, add=True)
        self.layer_opacity_scale.bind("<ButtonRelease-1>", self._end_layer_opacity_drag, add=True)

    def create_section_header(self, parent, text):
        """åˆ›å»ºå±æ€§é¢æ¿åˆ†ç»„æ ‡é¢˜"""
        header = tk.Frame(parent, bg="#e3f2fd")
        header.pack(fill=tk.X, pady=(10, 5))

        label = tk.Label(header, text=text, bg="#e3f2fd", fg="#1565C0",
                        font=(FONT_FAMILY, 9, "bold"), padx=10, pady=3)
        label.pack(fill=tk.X)
        return header

    def scroll_to_layers(self):
        """åˆ‡æ¢åˆ°å³ä¾§â€œå›¾å±‚â€Tabï¼ˆæˆ–å›é€€åˆ°æ—§çš„æ»šåŠ¨å®šä½ï¼‰ã€‚"""
        try:
            nb = getattr(self, "right_notebook", None)
            tab = getattr(self, "layers_tab", None)
            if nb is not None and tab is not None:
                nb.select(tab)
                return
        except Exception:
            pass

        # å…¼å®¹æ—§å¸ƒå±€ï¼šåœ¨æ»šåŠ¨å±æ€§é¢æ¿ä¸­å®šä½åˆ°â€œå›¾å±‚â€
        try:
            canvas = getattr(self, "prop_canvas", None)
            header = getattr(self, "layers_section_header", None)
            frame = getattr(self, "prop_frame", None)
            if canvas is None or header is None or frame is None:
                return
            canvas.update_idletasks()
            frame.update_idletasks()
            y = header.winfo_y()
            total = max(1, frame.winfo_height())
            canvas.yview_moveto(y / total)
        except Exception:
            pass

    def select_layer_by_id(self, layer_id: str):
        """åœ¨å›¾å±‚é¢æ¿ä¸­é€‰ä¸­æŒ‡å®šå›¾å±‚ï¼Œå¹¶åŒæ­¥ UIã€‚"""
        try:
            if not self.pages or not hasattr(self, "layer_tree"):
                return
            page = self.pages[self.current_page_index]
            layers = page.get("layers", [])
            idx = -1
            for i, layer in enumerate(layers):
                if layer and layer.get("id") == layer_id:
                    idx = i
                    break
            self.selected_layer_index = idx
            self.update_layer_listbox()
            try:
                self.layer_tree.selection_set(layer_id)
                self.layer_tree.focus(layer_id)
            except Exception:
                pass
        except Exception:
            pass

    def create_status_bar(self):
        """åˆ›å»ºåº•éƒ¨çŠ¶æ€æ  - PowerPointçº¢è‰²ä¸»é¢˜"""
        self.status_bar = tk.Frame(self.root, bg=COLOR_THEME, height=28)
        self.status_bar.pack(fill=tk.X, side=tk.BOTTOM)
        self.status_bar.pack_propagate(False)

        self.status_label = tk.Label(self.status_bar, text="å°±ç»ª - è¯·å¯¼å…¥å›¾ç‰‡å¼€å§‹ç¼–è¾‘",
                                     bg=COLOR_THEME, fg="white",
                                     font=(FONT_FAMILY, 9), padx=10)
        self.status_label.pack(side=tk.LEFT)

        self.status_info = tk.Label(self.status_bar, text="",
                                    bg=COLOR_THEME, fg="white",
                                    font=(FONT_FAMILY, 9), padx=10)
        self.status_info.pack(side=tk.RIGHT)

    def update_status(self, text):
        """æ›´æ–°çŠ¶æ€æ """
        self.status_label.config(text=text)

    def bind_shortcuts(self):
        """ç»‘å®šå¿«æ·é”®"""
        self.root.bind("<Control-z>", lambda e: self.undo())
        self.root.bind("<Control-y>", lambda e: self.redo())
        self.root.bind("<Delete>", lambda e: self.delete_selected_box())
        self.root.bind("<Left>", lambda e: self.prev_page())
        self.root.bind("<Right>", lambda e: self.next_page())
        self.root.bind("<Control-s>", lambda e: self.save_project())
        self.root.bind("<Control-o>", lambda e: self.load_project())
        # æ–°å¢å¿«æ·é”®
        self.root.bind("<Control-a>", lambda e: self.select_all_boxes())
        self.root.bind("<Control-c>", lambda e: self.copy_boxes())
        self.root.bind("<Control-v>", lambda e: self.paste_boxes())
        self.root.bind("<Left>", lambda e: self.move_box_by_key(-10, 0))
        self.root.bind("<Right>", lambda e: self.move_box_by_key(10, 0))
        self.root.bind("<Up>", lambda e: self.move_box_by_key(0, -10))
        self.root.bind("<Down>", lambda e: self.move_box_by_key(0, 10))
        self.root.bind("<Control-Left>", lambda e: self.move_box_by_key(-1, 0))
        self.root.bind("<Control-Right>", lambda e: self.move_box_by_key(1, 0))
        self.root.bind("<Control-Up>", lambda e: self.move_box_by_key(0, -1))
        self.root.bind("<Control-Down>", lambda e: self.move_box_by_key(0, 1))
        self.root.bind("<Prior>", lambda e: self.prev_page())
        self.root.bind("<Next>", lambda e: self.next_page())


    # ==================== é¡µé¢ç®¡ç† ====================

    # ç¼–è¾‘ç”¨çš„æœ€å¤§å›¾ç‰‡å°ºå¯¸ï¼ˆè¶…è¿‡æ­¤å°ºå¯¸ä¼šç¼©æ”¾ä»¥æé«˜æ€§èƒ½ï¼‰
    MAX_EDIT_SIZE = 1920

    def _resize_image_for_edit(self, img):
        """ç¼©æ”¾å›¾ç‰‡ç”¨äºç¼–è¾‘ï¼Œè¿”å›ç¼©æ”¾åçš„å›¾ç‰‡å’Œç¼©æ”¾æ¯”ä¾‹"""
        w, h = img.size
        if max(w, h) <= self.MAX_EDIT_SIZE:
            return img, 1.0

        scale = self.MAX_EDIT_SIZE / max(w, h)
        new_w = int(w * scale)
        new_h = int(h * scale)
        resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
        return resized, scale

    def load_multiple_images(self):
        """æ‰¹é‡åŠ è½½å¤šå¼ åŸå›¾"""
        file_paths = filedialog.askopenfilenames(
            title="é€‰æ‹©å¤šå¼ åŸå›¾ï¼ˆæŒ‰é¡ºåºé€‰æ‹©ï¼‰",
            filetypes=[("å›¾ç‰‡æ–‡ä»¶", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_paths:
            return

        if self.pages:
            self.save_current_page()

        clear_existing = False
        if self.pages:
            result = messagebox.askyesnocancel(
                "æç¤º", f"å·²æœ‰ {len(self.pages)} é¡µï¼Œæ˜¯å¦æ¸…ç©ºï¼Ÿ\n\næ˜¯ - æ¸…ç©ºåå¯¼å…¥\nå¦ - è¿½åŠ \nå–æ¶ˆ - å–æ¶ˆ"
            )
            if result is None:
                return
            elif result:
                self.pages = []
                clear_existing = True

        start_index = len(self.pages)

        for path in file_paths:
            original_img = Image.open(path)
            original_size = original_img.size  # ä¿å­˜åŸå§‹å°ºå¯¸

            # ç¼©æ”¾å›¾ç‰‡ç”¨äºç¼–è¾‘
            edit_img, edit_scale = self._resize_image_for_edit(original_img)

            page_data = {
                "original_path": path,
                "original_size": original_size,  # åŸå§‹å°ºå¯¸
                "edit_scale": edit_scale,  # ç¼–è¾‘ç¼©æ”¾æ¯”ä¾‹
                "bg_path": None,
                "bg_original_path": None,  # èƒŒæ™¯åŸå›¾è·¯å¾„
                "image": edit_img,  # ç¼–è¾‘ç”¨çš„ç¼©æ”¾å›¾ç‰‡
                "text_boxes": [],
                "layers": []
            }
            self.pages.append(page_data)

        self.current_page_index = start_index
        self.load_current_page()
        self.update_page_label()
        self.update_thumbnails()

        # éšè—å ä½ç¬¦
        self.placeholder_label.place_forget()

        # æ˜¾ç¤ºæ˜¯å¦æœ‰ç¼©æ”¾
        any_scaled = any(p["edit_scale"] < 1.0 for p in self.pages[start_index:])
        if any_scaled:
            self.update_status(f"å·²å¯¼å…¥ {len(file_paths)} å¼ å›¾ç‰‡ï¼ˆå¤§å›¾å·²è‡ªåŠ¨ç¼©æ”¾ä»¥æé«˜æ€§èƒ½ï¼‰ï¼Œå…± {len(self.pages)} é¡µ")
        else:
            self.update_status(f"å·²å¯¼å…¥ {len(file_paths)} å¼ å›¾ç‰‡ï¼Œå…± {len(self.pages)} é¡µ")

    def load_multiple_backgrounds(self):
        """æ‰¹é‡åŠ è½½èƒŒæ™¯å›¾ - è‡ªåŠ¨è°ƒæ•´å¤§å°ä¸ç¼–è¾‘å›¾ä¸€è‡´"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥åŸå›¾")
            return

        file_paths = filedialog.askopenfilenames(
            title="é€‰æ‹©èƒŒæ™¯å›¾",
            filetypes=[("å›¾ç‰‡æ–‡ä»¶", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_paths:
            return

        matched = 0
        for bg_path in file_paths:
            bg_name = os.path.splitext(os.path.basename(bg_path))[0].lower()
            for page in self.pages:
                orig_name = os.path.splitext(os.path.basename(page["original_path"]))[0].lower()
                if bg_name == orig_name or orig_name in bg_name or bg_name in orig_name:
                    # è°ƒæ•´èƒŒæ™¯å›¾å¤§å°ä¸ç¼–è¾‘å›¾ä¸€è‡´
                    resized_bg_path = self._resize_bg_to_match(bg_path, page["image"].size)
                    page["bg_path"] = resized_bg_path
                    matched += 1
                    break

        # å¦‚æœæ²¡æœ‰åŒ¹é…åˆ°ï¼ŒæŒ‰é¡ºåºåˆ†é…
        if matched == 0 and len(file_paths) == len(self.pages):
            for i, bg_path in enumerate(file_paths):
                resized_bg_path = self._resize_bg_to_match(bg_path, self.pages[i]["image"].size)
                self.pages[i]["bg_path"] = resized_bg_path
            matched = len(file_paths)

        # æ›´æ–°å½“å‰é¡µèƒŒæ™¯è·¯å¾„
        if self.pages and self.current_page_index < len(self.pages):
            self.clean_bg_path = self.pages[self.current_page_index].get("bg_path")

        # åˆ·æ–°æ˜¾ç¤º
        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"å·²åŒ¹é… {matched}/{len(self.pages)} å¼ èƒŒæ™¯å›¾")

    def create_blank_page(self):
        """åˆ›å»ºç©ºç™½é¡µé¢"""
        # åˆ›å»ºå¯¹è¯æ¡†
        dialog = tk.Toplevel(self.root)
        dialog.title("æ–°å»ºç©ºç™½å›¾")
        dialog.geometry("450x400")
        dialog.transient(self.root)
        dialog.grab_set()

        tk.Label(dialog, text="æ–°å»ºç©ºç™½å›¾", font=(FONT_FAMILY, 14, "bold")).pack(pady=(20, 10))

        # æ¯”ä¾‹é€‰æ‹©
        ratio_frame = tk.LabelFrame(dialog, text="é€‰æ‹©æ¯”ä¾‹", font=(FONT_FAMILY, 10, "bold"),
                                    padx=15, pady=15)
        ratio_frame.pack(fill=tk.X, padx=20, pady=(0, 15))

        ratio_var = tk.StringVar(value="16:9")

        ratios = [
            ("16:9 (1920Ã—1080)", "16:9"),
            ("9:16 (1080Ã—1920)", "9:16"),
            ("4:3 (1600Ã—1200)", "4:3"),
            ("3:4 (1200Ã—1600)", "3:4"),
            ("1:1 (1200Ã—1200)", "1:1"),
        ]

        for text, value in ratios:
            tk.Radiobutton(ratio_frame, text=text, variable=ratio_var, value=value,
                          font=(FONT_FAMILY, 10)).pack(anchor="w", pady=2)

        # é¢œè‰²é€‰æ‹©
        color_frame = tk.LabelFrame(dialog, text="èƒŒæ™¯é¢œè‰²", font=(FONT_FAMILY, 10, "bold"),
                                    padx=15, pady=15)
        color_frame.pack(fill=tk.X, padx=20, pady=(0, 15))

        selected_color = tk.StringVar(value="#FFFFFF")

        color_display_row = tk.Frame(color_frame)
        color_display_row.pack(fill=tk.X, pady=5)

        tk.Label(color_display_row, text="å½“å‰é¢œè‰²:", font=(FONT_FAMILY, 9)).pack(side=tk.LEFT)

        color_display = tk.Label(color_display_row, text="      ", bg="#FFFFFF",
                                relief=tk.RIDGE, borderwidth=2, width=10)
        color_display.pack(side=tk.LEFT, padx=10)

        color_label = tk.Label(color_display_row, text="#FFFFFF", font=(FONT_FAMILY, 9))
        color_label.pack(side=tk.LEFT)

        def choose_color():
            color = colorchooser.askcolor(title="é€‰æ‹©èƒŒæ™¯é¢œè‰²",
                                         initialcolor=selected_color.get())
            if color[1]:
                selected_color.set(color[1])
                color_display.config(bg=color[1])
                color_label.config(text=color[1])

        tk.Button(color_frame, text="é€‰æ‹©é¢œè‰²", command=choose_color,
                 bg=COLOR_THEME, fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10), padx=15, pady=5).pack(pady=5)

        # æŒ‰é’®
        btn_frame = tk.Frame(dialog)
        btn_frame.pack(pady=20)

        def on_create():
            ratio = ratio_var.get()
            color_hex = selected_color.get()

            # è®¡ç®—å°ºå¯¸
            ratio_sizes = {
                "16:9": (1920, 1080),
                "9:16": (1080, 1920),
                "4:3": (1600, 1200),
                "3:4": (1200, 1600),
                "1:1": (1200, 1200),
            }

            width, height = ratio_sizes.get(ratio, (1920, 1080))

            try:
                # åˆ›å»ºç©ºç™½å›¾ç‰‡
                blank_img = Image.new("RGB", (width, height), color_hex)

                # ä¿å­˜åˆ°ä¸´æ—¶æ–‡ä»¶
                temp_dir = os.path.join(get_base_dir(), "temp_blank_pages")
                os.makedirs(temp_dir, exist_ok=True)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                temp_path = os.path.join(temp_dir, f"blank_{ratio.replace(':', 'x')}_{timestamp}.png")
                blank_img.save(temp_path)

                # åˆ›å»ºé¡µé¢æ•°æ®
                if self.pages:
                    self.save_current_page()

                # è¯¢é—®æ˜¯è¿½åŠ è¿˜æ˜¯æ¸…ç©º
                clear_existing = False
                if self.pages:
                    result = messagebox.askyesnocancel(
                        "æç¤º", f"å·²æœ‰ {len(self.pages)} é¡µï¼Œæ˜¯å¦æ¸…ç©ºï¼Ÿ\n\næ˜¯ - æ¸…ç©ºåæ·»åŠ \nå¦ - è¿½åŠ \nå–æ¶ˆ - å–æ¶ˆ"
                    )
                    if result is None:
                        dialog.destroy()
                        return
                    elif result:
                        self.pages = []
                        clear_existing = True

                start_index = len(self.pages)

                # ç¼©æ”¾å›¾ç‰‡ç”¨äºç¼–è¾‘
                edit_img, edit_scale = self._resize_image_for_edit(blank_img)

                page_data = {
                    "original_path": temp_path,
                    "original_size": (width, height),
                    "edit_scale": edit_scale,
                    "bg_path": None,
                    "bg_original_path": None,
                    "image": edit_img,
                    "text_boxes": [],
                    "layers": []
                }
                self.pages.append(page_data)

                self.current_page_index = start_index
                self.load_current_page()
                self.update_page_label()
                self.update_thumbnails()

                # éšè—å ä½ç¬¦
                self.placeholder_label.place_forget()

                dialog.destroy()
                self.update_status(f"å·²åˆ›å»ºç©ºç™½å›¾ ({ratio})ï¼Œå…± {len(self.pages)} é¡µ")
                messagebox.showinfo("æˆåŠŸ", f"ç©ºç™½å›¾åˆ›å»ºæˆåŠŸï¼\næ¯”ä¾‹: {ratio}\né¢œè‰²: {color_hex}")

            except Exception as e:
                messagebox.showerror("é”™è¯¯", f"åˆ›å»ºç©ºç™½å›¾å¤±è´¥:\n{e}")

        def on_cancel():
            dialog.destroy()

        tk.Button(btn_frame, text="åˆ›å»º", command=on_create,
                 bg="#2196F3", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10, "bold"), padx=25, pady=10).pack(side=tk.LEFT, padx=5)

        tk.Button(btn_frame, text="å–æ¶ˆ", command=on_cancel,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10), padx=25, pady=10).pack(side=tk.LEFT, padx=5)

    def load_current_page_background(self):
        """ä¸ºå½“å‰é¡µå•ç‹¬è®¾ç½®èƒŒæ™¯å›¾"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥åŸå›¾")
            return

        file_path = filedialog.askopenfilename(
            title=f"é€‰æ‹©ç¬¬ {self.current_page_index + 1} é¡µçš„èƒŒæ™¯å›¾",
            filetypes=[("å›¾ç‰‡æ–‡ä»¶", "*.jpg *.jpeg *.png *.bmp")]
        )
        if not file_path:
            return

        page = self.pages[self.current_page_index]
        edit_size = page["image"].size

        # è°ƒæ•´èƒŒæ™¯å›¾å¤§å°ä¸ç¼–è¾‘å›¾ä¸€è‡´
        resized_bg_path = self._resize_bg_to_match(file_path, edit_size)
        page["bg_path"] = resized_bg_path
        self.clean_bg_path = resized_bg_path

        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"ç¬¬ {self.current_page_index + 1} é¡µèƒŒæ™¯å·²è®¾ç½®")

    def _resize_bg_to_match(self, bg_path, target_size):
        """è°ƒæ•´èƒŒæ™¯å›¾å¤§å°ä¸ç›®æ ‡å°ºå¯¸ä¸€è‡´ï¼Œè¿”å›è°ƒæ•´åçš„å›¾ç‰‡è·¯å¾„"""
        bg_img = Image.open(bg_path)

        # å¦‚æœå¤§å°å·²ç»ä¸€è‡´ï¼Œç›´æ¥è¿”å›åŸè·¯å¾„
        if bg_img.size == target_size:
            return bg_path

        # è°ƒæ•´å¤§å°
        resized_img = bg_img.resize(target_size, Image.Resampling.LANCZOS)

        # ä¿å­˜åˆ°ä¸´æ—¶æ–‡ä»¶
        bg_dir = os.path.dirname(bg_path)
        bg_name = os.path.splitext(os.path.basename(bg_path))[0]
        bg_ext = os.path.splitext(bg_path)[1]

        # åˆ›å»ºè°ƒæ•´åçš„æ–‡ä»¶å
        resized_path = os.path.join(bg_dir, f"{bg_name}_resized_{target_size[0]}x{target_size[1]}{bg_ext}")

        # å¦‚æœå·²å­˜åœ¨åŒåè°ƒæ•´åçš„æ–‡ä»¶ï¼Œæ£€æŸ¥æ˜¯å¦éœ€è¦é‡æ–°ç”Ÿæˆ
        if not os.path.exists(resized_path):
            if resized_img.mode == 'RGBA' and bg_ext.lower() in ['.jpg', '.jpeg']:
                resized_img = resized_img.convert('RGB')
            resized_img.save(resized_path, quality=95)

        return resized_path

    def clear_current_page_background(self):
        """æ¸…é™¤å½“å‰é¡µèƒŒæ™¯"""
        if not self.pages:
            return

        self.pages[self.current_page_index]["bg_path"] = None
        self.clean_bg_path = None
        self.update_bg_status()
        self.update_thumbnails()
        self.refresh_canvas()
        self.update_status(f"ç¬¬ {self.current_page_index + 1} é¡µèƒŒæ™¯å·²æ¸…é™¤")

    def update_bg_status(self):
        return page_manager_core.update_bg_status(self)

    def save_current_page(self):
        return page_manager_core.save_current_page(self)

    def load_current_page(self):
        return page_manager_core.load_current_page(self)

    def prev_page(self):
        return page_manager_core.prev_page(self)

    def next_page(self):
        return page_manager_core.next_page(self)

    def update_page_label(self):
        return page_manager_core.update_page_label(self)

    def update_status_info(self):
        return page_manager_core.update_status_info(self)

    def update_thumbnails(self):
        return page_manager_core.update_thumbnails(self)

    def show_thumbnail_menu(self, event, page_index):
        return page_manager_core.show_thumbnail_menu(self, event, page_index)

    def set_page_background(self, page_index):
        return page_manager_core.set_page_background(self, page_index)

    def clear_page_background(self, page_index):
        return page_manager_core.clear_page_background(self, page_index)

    def delete_page(self, page_index):
        return page_manager_core.delete_page(self, page_index)

    def highlight_current_thumbnail(self):
        return page_manager_core.highlight_current_thumbnail(self)

    def go_to_page(self, index):
        return page_manager_core.go_to_page(self, index)

    # ==================== ç”»å¸ƒæ“ä½œ ====================

    def fit_image_to_canvas(self):
        return page_manager_core.fit_image_to_canvas(self)

    def on_canvas_resize(self, event):
        return page_manager_core.on_canvas_resize(self, event)

    def on_canvas_zoom(self, event):
        return page_manager_core.on_canvas_zoom(self, event)

    def on_canvas_scroll(self, event):
        return page_manager_core.on_canvas_scroll(self, event)

    def zoom_to_100(self):
        return page_manager_core.zoom_to_100(self)

    def refresh_canvas(self):
        """åˆ·æ–°ç”»å¸ƒ"""
        if not self.original_image:
            return

        if self.current_preview_mode == "ppt":
            self._draw_ppt_preview()
        elif self.current_preview_mode == "edit":
            self._draw_original_with_boxes()
        else:
            self._draw_raw_with_boxes()

        self.update_status_info()

    def _draw_raw_with_boxes(self):
        """ç»˜åˆ¶åŸå§‹åº•å›¾ï¼ˆä¸å åŠ èƒŒæ™¯/å›¾å±‚ï¼‰+æ¡†"""
        self.canvas.delete("all")

        base_img = None
        try:
            if self.pages and 0 <= self.current_page_index < len(self.pages):
                base_img = self.pages[self.current_page_index].get("image")
        except Exception:
            base_img = None
        if base_img is None:
            base_img = self.original_image
        if base_img is None:
            return

        img_w, img_h = base_img.size
        display_w = int(img_w * self.scale)
        display_h = int(img_h * self.scale)

        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()
        offset_x = max(0, (canvas_w - display_w) // 2)
        offset_y = max(0, (canvas_h - display_h) // 2)

        self.display_image = base_img.resize((display_w, display_h), Image.Resampling.LANCZOS)
        self.tk_image = ImageTk.PhotoImage(self.display_image)
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self.tk_image, tags="image")

        self.canvas_offset_x = offset_x
        self.canvas_offset_y = offset_y

        for idx, box in enumerate(self.text_boxes):
            self.draw_box(idx, box, offset_x, offset_y)

        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        self.canvas.config(
            scrollregion=(
                0,
                0,
                max(canvas_w, display_w + offset_x * 2),
                max(canvas_h, display_h + offset_y * 2),
            )
        )
    def _draw_original_with_boxes(self):
        """ç»˜åˆ¶åŸå›¾+æ¡†"""
        self.canvas.delete("all")

        base_img = self.get_current_page_composited_background()
        if base_img is None:
            return

        img_w, img_h = base_img.size
        display_w = int(img_w * self.scale)
        display_h = int(img_h * self.scale)

        # å±…ä¸­æ˜¾ç¤º
        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()
        offset_x = max(0, (canvas_w - display_w) // 2)
        offset_y = max(0, (canvas_h - display_h) // 2)

        self.display_image = base_img.resize((display_w, display_h), Image.Resampling.LANCZOS)
        self.tk_image = ImageTk.PhotoImage(self.display_image)
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self.tk_image, tags="image")

        # ä¿å­˜åç§»é‡ç”¨äºåæ ‡è½¬æ¢
        self.canvas_offset_x = offset_x
        self.canvas_offset_y = offset_y

        for idx, box in enumerate(self.text_boxes):
            self.draw_box(idx, box, offset_x, offset_y)

        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        self.canvas.config(scrollregion=(0, 0, max(canvas_w, display_w + offset_x * 2),
                                          max(canvas_h, display_h + offset_y * 2)))

    def _draw_ppt_preview(self):
        """ç»˜åˆ¶PPTé¢„è§ˆ"""
        base_img = self.get_current_page_composited_background()
        if base_img is None:
            return

        preview_img = base_img.copy().convert("RGBA")
        img_w, img_h = preview_img.size

        try:
            draw = ImageDraw.Draw(preview_img)

            for box in self.text_boxes:
                if not box.text:
                    continue

                pixel_font_size = int(box.font_size * 96 / 72)

                try:
                    font_path = self._get_font_path(box.font_name)
                    if font_path:
                        font = ImageFont.truetype(font_path, pixel_font_size)
                    else:
                        font = ImageFont.load_default()
                except:
                    font = ImageFont.load_default()

                color_hex = box.font_color.lstrip('#')
                r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)

                # Pillow é»˜è®¤ä»¥å·¦ä¸Šè§’ä¸ºåŸç‚¹ç»˜åˆ¶ï¼›ä¸åŒå­—ä½“ä¼šæœ‰ ascent/descent å¯¼è‡´è§†è§‰ä¸Šâ€œåä¸‹â€ã€‚
                # ä¼˜å…ˆä½¿ç”¨ anchor ä»¥æ–‡æœ¬è‡ªèº«ä¸­çº¿å¯¹é½ï¼Œå®ç°å‚ç›´å±…ä¸­ã€‚
                center_y = box.y + box.height // 2
                if box.align == "center":
                    text_x = box.x + box.width // 2
                    anchor = "mm"
                elif box.align == "right":
                    text_x = box.x + box.width - 3
                    anchor = "rm"
                else:
                    text_x = box.x + 3
                    anchor = "lm"

                try:
                    draw.text((text_x, center_y), box.text, font=font, fill=(r, g, b, 255), anchor=anchor)
                except TypeError:
                    # å…¼å®¹æ—§ç‰ˆ Pillowï¼ˆä¸æ”¯æŒ anchor å‚æ•°ï¼‰ï¼šä½¿ç”¨ bbox åç§»çŸ«æ­£åˆ°å‚ç›´å±…ä¸­
                    try:
                        bbox = draw.textbbox((0, 0), box.text, font=font)
                        text_w = bbox[2] - bbox[0]
                        text_h = bbox[3] - bbox[1]
                        y = box.y + (box.height - text_h) // 2 - bbox[1]
                        if box.align == "center":
                            x = box.x + (box.width - text_w) // 2 - bbox[0]
                        elif box.align == "right":
                            x = box.x + box.width - text_w - 3 - bbox[0]
                        else:
                            x = box.x + 3 - bbox[0]
                        draw.text((x, y), box.text, font=font, fill=(r, g, b, 255))
                    except Exception:
                        draw.text((box.x + 3, box.y + 2), box.text, font=font, fill=(r, g, b, 255))

        except Exception as e:
            print(f"ç»˜åˆ¶æ–‡å­—å¤±è´¥: {e}")

        preview_img = preview_img.convert("RGB")

        canvas_w = self.canvas.winfo_width()
        canvas_h = self.canvas.winfo_height()

        # ä¿æŒå½“å‰ç¼©æ”¾æ¯”ä¾‹ï¼Œä¸å¼ºåˆ¶é‡ç½®
        display_w = int(img_w * self.scale)
        display_h = int(img_h * self.scale)

        offset_x = max(0, (canvas_w - display_w) // 2)
        offset_y = max(0, (canvas_h - display_h) // 2)

        self.canvas_offset_x = offset_x
        self.canvas_offset_y = offset_y

        preview_img = preview_img.resize((display_w, display_h), Image.Resampling.LANCZOS)
        self.ppt_preview_image = ImageTk.PhotoImage(preview_img)

        self.canvas.delete("all")
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self.ppt_preview_image)

        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        for idx, box in enumerate(self.text_boxes):
            self._draw_ppt_edit_box(idx, box, offset_x, offset_y)

        self.canvas.config(scrollregion=(0, 0, max(canvas_w, display_w + offset_x * 2),
                                          max(canvas_h, display_h + offset_y * 2)))

    def get_current_page_composited_background(self):
        """
        è·å–å½“å‰é¡µâ€œåº•å›¾â€ï¼ˆèƒŒæ™¯/åŸå›¾ï¼‰+å›¾å±‚åˆæˆåçš„å›¾ç‰‡ï¼ˆä¸å«æ–‡æœ¬æ¸²æŸ“ï¼‰ã€‚
        """
        if not self.pages:
            return self.original_image.copy() if self.original_image else None
        page = self.pages[self.current_page_index]
        return self.get_page_composited_background(page)

    def get_page_composited_background(self, page):
        """
        è·å–æŒ‡å®šé¡µâ€œåº•å›¾â€ï¼ˆèƒŒæ™¯/åŸå›¾ï¼‰+å›¾å±‚åˆæˆåçš„å›¾ç‰‡ï¼ˆä¸å«æ–‡æœ¬æ¸²æŸ“ï¼‰ã€‚
        åæ ‡ç³»ä»¥ page["image"] ä¸ºå‡†ï¼ˆç¼–è¾‘å°ºå¯¸ï¼‰ã€‚
        """
        base_img = None
        bg_path = page.get("bg_path")
        if bg_path and os.path.exists(bg_path):
            try:
                base_img = Image.open(bg_path)
            except Exception:
                base_img = None

        if base_img is None:
            if page.get("image") is not None:
                base_img = page["image"].copy()
            else:
                return self.original_image.copy() if self.original_image else None

        edit_img = page.get("image")
        if edit_img is not None and base_img.size != edit_img.size:
            base_img = base_img.resize(edit_img.size, Image.Resampling.LANCZOS)

        base_rgba = base_img.convert("RGBA")

        layers = page.get("layers", [])
        # PSä¹ æƒ¯ï¼šåˆ—è¡¨é¡¶éƒ¨ä¸ºæœ€ä¸Šå±‚ï¼›åˆæˆæ—¶åº”ä»åº•åˆ°é¡¶ç»˜åˆ¶ï¼ˆåå‘éå†ï¼‰
        for layer in reversed(layers):
            if not layer or not layer.get("visible", True):
                continue
            path = layer.get("path")
            if not path or not os.path.exists(path):
                continue

            try:
                overlay = Image.open(path).convert("RGBA")
            except Exception:
                continue

            # è£å‰ªï¼ˆä»¥å›¾å±‚åŸå›¾åæ ‡ç³»ä¸ºå‡†ï¼‰
            crop = layer.get("crop")
            if crop:
                try:
                    if isinstance(crop, dict):
                        x0 = int(crop.get("x0", 0))
                        y0 = int(crop.get("y0", 0))
                        x1 = int(crop.get("x1", overlay.size[0]))
                        y1 = int(crop.get("y1", overlay.size[1]))
                    else:
                        x0, y0, x1, y1 = [int(v) for v in crop]
                    x0 = max(0, min(overlay.size[0], x0))
                    y0 = max(0, min(overlay.size[1], y0))
                    x1 = max(0, min(overlay.size[0], x1))
                    y1 = max(0, min(overlay.size[1], y1))
                    if x1 > x0 and y1 > y0:
                        overlay = overlay.crop((x0, y0, x1, y1))
                except Exception:
                    pass

            # ç¼©æ”¾
            try:
                scale = float(layer.get("scale", 1.0))
            except Exception:
                scale = 1.0
            if scale <= 0:
                scale = 1.0
            if abs(scale - 1.0) > 1e-6:
                try:
                    new_w = max(1, int(round(overlay.size[0] * scale)))
                    new_h = max(1, int(round(overlay.size[1] * scale)))
                    overlay = overlay.resize((new_w, new_h), Image.Resampling.LANCZOS)
                except Exception:
                    pass

            opacity = float(layer.get("opacity", 1.0))
            opacity = max(0.0, min(opacity, 1.0))
            if opacity < 1.0:
                r, g, b, a = overlay.split()
                a = a.point(lambda v: int(v * opacity))
                overlay = Image.merge("RGBA", (r, g, b, a))

            x = int(layer.get("x", 0))
            y = int(layer.get("y", 0))
            base_rgba.paste(overlay, (x, y), overlay)

        return base_rgba.convert("RGB")

    def _ensure_page_layers(self, page):
        return page.setdefault("layers", [])

    def add_image_layer(self, page, image, name="AIå›¾å±‚", x=0, y=0, opacity=1.0, visible=True):
        """
        å°†ä¸€å¼  PIL Image ä¿å­˜ä¸ºå›¾å±‚å¹¶åŠ å…¥åˆ°é¡µé¢ layersã€‚
        """
        layers = self._ensure_page_layers(page)

        temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
        os.makedirs(temp_dir, exist_ok=True)

        layer_id = uuid.uuid4().hex[:10]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = os.path.join(temp_dir, f"layer_{layer_id}_{timestamp}.png")

        img = image.convert("RGBA") if image.mode != "RGBA" else image
        img.save(path)

        layer = {
            "id": layer_id,
            "name": name,
            "path": path,
            "x": int(x),
            "y": int(y),
            "scale": float(1.0),
            "crop": None,
            "locked": False,
            "opacity": float(opacity),
            "visible": bool(visible),
        }
        # æ–°å›¾å±‚é»˜è®¤ç½®é¡¶ï¼ˆæ›´ç¬¦åˆPSä¹ æƒ¯ï¼‰
        layers.insert(0, layer)
        return layer

    def update_layer_listbox(self):
        if not hasattr(self, "layer_tree"):
            return

        page = self.pages[self.current_page_index] if self.pages else None
        layers = page.get("layers", []) if page else []

        # å…¼å®¹ï¼šç¼ºå°‘ id çš„å›¾å±‚è¡¥é½
        for layer in layers:
            if layer is not None and not layer.get("id"):
                layer["id"] = uuid.uuid4().hex[:10]
            if layer is not None:
                layer.setdefault("x", 0)
                layer.setdefault("y", 0)
                layer.setdefault("opacity", 1.0)
                layer.setdefault("visible", True)
                layer.setdefault("scale", 1.0)
                layer.setdefault("crop", None)
                layer.setdefault("locked", False)

        for iid in self.layer_tree.get_children(""):
            self.layer_tree.delete(iid)

        self._layer_thumb_refs = {}

        for i, layer in enumerate(layers):
            if not layer:
                continue
            vis = "âœ“" if layer.get("visible", True) else "Ã—"
            opacity = int(float(layer.get("opacity", 1.0)) * 100)
            name = layer.get("name") or f"å›¾å±‚{i+1}"
            lock_flag = "é”" if layer.get("locked") else ""

            thumb = None
            path = layer.get("path")
            if path and os.path.exists(path):
                try:
                    img = Image.open(path).convert("RGBA")
                    img.thumbnail((48, 32), Image.Resampling.LANCZOS)
                    thumb = ImageTk.PhotoImage(img)
                    self._layer_thumb_refs[layer["id"]] = thumb
                except Exception:
                    thumb = None

            self.layer_tree.insert(
                "",
                "end",
                iid=layer["id"],
                text=name,
                image=thumb,
                values=(vis, f"{opacity}%", lock_flag),
            )

        # æ¢å¤é€‰æ‹© + åŒæ­¥é€æ˜åº¦
        if self.selected_layer_index >= len(layers):
            self.selected_layer_index = -1

        page, layers, layer = self._get_selected_layer()
        if layer is not None:
            try:
                self.layer_tree.selection_set(layer["id"])
                self.layer_tree.focus(layer["id"])
            except Exception:
                pass
            try:
                opacity = float(layer.get("opacity", 1.0))
                self._layer_opacity_syncing = True
                self.layer_opacity_scale.set(int(opacity * 100))
            except Exception:
                self._layer_opacity_syncing = True
                self.layer_opacity_scale.set(100)
            finally:
                self._layer_opacity_syncing = False

            self._sync_layer_transform_controls(layer)
        else:
            self._layer_opacity_syncing = True
            self.layer_opacity_scale.set(100)
            self._layer_opacity_syncing = False
            self._sync_layer_transform_controls(None)

    def on_layer_select(self, event=None):
        if not self.pages:
            return
        page = self.pages[self.current_page_index]
        layers = page.get("layers", [])

        idx = -1
        if hasattr(self, "layer_tree"):
            sel = self.layer_tree.selection()
            if sel:
                selected_iid = sel[0]
                for i, layer in enumerate(layers):
                    if layer and layer.get("id") == selected_iid:
                        idx = i
                        break

        self.selected_layer_index = idx
        if 0 <= idx < len(layers):
            try:
                opacity = float(layers[idx].get("opacity", 1.0))
                self._layer_opacity_syncing = True
                self.layer_opacity_scale.set(int(opacity * 100))
            except Exception:
                self._layer_opacity_syncing = True
                self.layer_opacity_scale.set(100)
            finally:
                self._layer_opacity_syncing = False

            self._sync_layer_transform_controls(layers[idx])
        else:
            self._sync_layer_transform_controls(None)

        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

    def _get_selected_layer(self):
        if not self.pages:
            return None, None, None
        page = self.pages[self.current_page_index]
        layers = page.get("layers", [])

        idx = self.selected_layer_index
        if hasattr(self, "layer_tree"):
            sel = self.layer_tree.selection()
            if sel:
                selected_iid = sel[0]
                for i, layer in enumerate(layers):
                    if layer and layer.get("id") == selected_iid:
                        idx = i
                        break

        if idx is None or idx < 0 or idx >= len(layers):
            return page, layers, None

        self.selected_layer_index = idx
        return page, layers, layers[idx]

    def _layer_bbox(self, layer):
        """è¿”å›å›¾å±‚åœ¨é¡µé¢åæ ‡ç³»ä¸‹çš„åŒ…å›´ç›’ (x0,y0,x1,y1)ï¼Œè€ƒè™‘è£å‰ªä¸ç¼©æ”¾ã€‚"""
        try:
            path = layer.get("path")
            if not path or not os.path.exists(path):
                return None
            w0, h0 = Image.open(path).size

            crop = layer.get("crop")
            if crop:
                try:
                    if isinstance(crop, dict):
                        x0 = int(crop.get("x0", 0))
                        y0 = int(crop.get("y0", 0))
                        x1 = int(crop.get("x1", w0))
                        y1 = int(crop.get("y1", h0))
                    else:
                        x0, y0, x1, y1 = [int(v) for v in crop]
                    x0 = max(0, min(w0, x0))
                    y0 = max(0, min(h0, y0))
                    x1 = max(0, min(w0, x1))
                    y1 = max(0, min(h0, y1))
                    if x1 > x0 and y1 > y0:
                        w0 = x1 - x0
                        h0 = y1 - y0
                except Exception:
                    pass

            try:
                scale = float(layer.get("scale", 1.0))
            except Exception:
                scale = 1.0
            if scale <= 0:
                scale = 1.0

            w = max(1, int(round(w0 * scale)))
            h = max(1, int(round(h0 * scale)))
            x = int(layer.get("x", 0))
            y = int(layer.get("y", 0))
            return x, y, x + w, y + h
        except Exception:
            return None

    def toggle_selected_layer(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        self.save_state("layers")
        layer["visible"] = not layer.get("visible", True)
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def delete_selected_layer(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        self.save_state("layers")
        del layers[self.selected_layer_index]
        self.selected_layer_index = min(self.selected_layer_index, len(layers) - 1)
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def move_layer_up(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        idx = self.selected_layer_index
        if idx <= 0:
            return
        self.save_state("layers")
        layers[idx - 1], layers[idx] = layers[idx], layers[idx - 1]
        self.selected_layer_index = idx - 1
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def move_layer_down(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        idx = self.selected_layer_index
        if idx >= len(layers) - 1:
            return
        self.save_state("layers")
        layers[idx + 1], layers[idx] = layers[idx], layers[idx + 1]
        self.selected_layer_index = idx + 1
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def on_layer_opacity_change(self, value):
        if getattr(self, "_layer_opacity_syncing", False):
            return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        try:
            opacity = float(value) / 100.0
        except Exception:
            opacity = 1.0
        opacity = max(0.0, min(opacity, 1.0))
        prev = float(layer.get("opacity", 1.0))
        if abs(prev - opacity) < 1e-6:
            return
        # æ‹–åŠ¨æ»‘æ†æ—¶é¿å…åˆ·å±å†å²ï¼šåªåœ¨ä¸€æ¬¡æ‹–åŠ¨çš„é¦–æ¬¡å˜æ›´æ—¶ä¿å­˜å¿«ç…§
        if getattr(self, "_layer_opacity_drag_active", False):
            if not getattr(self, "_layer_opacity_saved", False):
                self.save_state("layers")
                self._layer_opacity_saved = True
        else:
            self.save_state("layers")
        layer["opacity"] = opacity
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def _begin_layer_opacity_drag(self, event=None):
        self._layer_opacity_drag_active = True
        self._layer_opacity_saved = False

    def _end_layer_opacity_drag(self, event=None):
        self._layer_opacity_drag_active = False
        self._layer_opacity_saved = False

    def _sync_layer_transform_controls(self, layer):
        """åŒæ­¥å›¾å±‚å˜æ¢ UIï¼ˆX/Y/ç¼©æ”¾/é”å®šï¼‰åˆ°å½“å‰é€‰æ‹©ã€‚"""
        if not hasattr(self, "layer_x_var"):
            return
        try:
            self._layer_transform_syncing = True
            if not layer:
                self.layer_x_var.set(0)
                self.layer_y_var.set(0)
                self.layer_scale_var.set(100)
                self.layer_lock_var.set(False)
                for w in (
                    getattr(self, "layer_x_entry", None),
                    getattr(self, "layer_y_entry", None),
                    getattr(self, "layer_scale_slider", None),
                    getattr(self, "layer_lock_check", None),
                    getattr(self, "layer_crop_btn", None),
                    getattr(self, "layer_reset_crop_btn", None),
                ):
                    try:
                        if w is not None:
                            w.config(state="disabled")
                    except Exception:
                        pass
                return

            layer.setdefault("x", 0)
            layer.setdefault("y", 0)
            layer.setdefault("scale", 1.0)
            layer.setdefault("crop", None)
            layer.setdefault("locked", False)

            self.layer_x_var.set(int(layer.get("x", 0)))
            self.layer_y_var.set(int(layer.get("y", 0)))
            try:
                s = float(layer.get("scale", 1.0))
            except Exception:
                s = 1.0
            self.layer_scale_var.set(int(max(10, min(300, round(s * 100)))))
            self.layer_lock_var.set(bool(layer.get("locked", False)))

            state = "disabled" if layer.get("locked") else "normal"
            for w in (
                getattr(self, "layer_x_entry", None),
                getattr(self, "layer_y_entry", None),
                getattr(self, "layer_scale_slider", None),
                getattr(self, "layer_crop_btn", None),
                getattr(self, "layer_reset_crop_btn", None),
            ):
                try:
                    if w is not None:
                        w.config(state=state)
                except Exception:
                    pass
            try:
                if getattr(self, "layer_lock_check", None) is not None:
                    self.layer_lock_check.config(state="normal")
            except Exception:
                pass
        finally:
            self._layer_transform_syncing = False

    def toggle_selected_layer_lock(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        self.save_state("layers")
        layer["locked"] = not bool(layer.get("locked", False))
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def _on_layer_lock_toggle(self):
        if getattr(self, "_layer_transform_syncing", False):
            return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        self.save_state("layers")
        layer["locked"] = bool(self.layer_lock_var.get())
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def _apply_layer_transform_from_ui(self):
        if getattr(self, "_layer_transform_syncing", False):
            return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        if layer.get("locked"):
            self.update_status("å›¾å±‚å·²é”å®šï¼Œæ— æ³•ä¿®æ”¹ä½ç½®/ç¼©æ”¾/è£å‰ª")
            self._sync_layer_transform_controls(layer)
            return

        try:
            x = int(self.layer_x_var.get())
            y = int(self.layer_y_var.get())
        except Exception:
            return
        old_x = int(layer.get("x", 0))
        old_y = int(layer.get("y", 0))
        if x == old_x and y == old_y:
            return
        self.save_state("layers")
        layer["x"] = x
        layer["y"] = y
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def _begin_layer_scale_drag(self, event=None):
        self._layer_scale_drag_active = True
        self._layer_scale_saved = False

    def _end_layer_scale_drag(self, event=None):
        self._layer_scale_drag_active = False
        self._layer_scale_saved = False

    def _on_layer_scale_change(self, value):
        if getattr(self, "_layer_transform_syncing", False):
            return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        if layer.get("locked"):
            self.update_status("å›¾å±‚å·²é”å®šï¼Œæ— æ³•ä¿®æ”¹ä½ç½®/ç¼©æ”¾/è£å‰ª")
            self._sync_layer_transform_controls(layer)
            return
        try:
            scale_pct = int(float(value))
        except Exception:
            return
        scale_pct = max(10, min(300, scale_pct))
        new_scale = scale_pct / 100.0
        try:
            old_scale = float(layer.get("scale", 1.0))
        except Exception:
            old_scale = 1.0
        if abs(old_scale - new_scale) < 1e-6:
            return
        if getattr(self, "_layer_scale_drag_active", False):
            if not getattr(self, "_layer_scale_saved", False):
                self.save_state("layers")
                self._layer_scale_saved = True
        else:
            self.save_state("layers")
        layer["scale"] = new_scale
        self.refresh_canvas()
        self.mark_unsaved()
        self.update_layer_listbox()

    def solid_color_cutout_selected_layer(self):
        """å¯¹é€‰å®šçš„å›¾å±‚è¿›è¡Œçº¯è‰²æŠ å›¾"""
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            messagebox.showwarning("æç¤º", "è¯·å…ˆé€‰æ‹©ä¸€ä¸ªå›¾å±‚")
            return
        if layer.get("locked"):
            self.update_status("å›¾å±‚å·²é”å®šï¼Œæ— æ³•ç¼–è¾‘")
            return

        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨")
            return

        try:
            img = Image.open(path).convert("RGBA")
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•æ‰“å¼€å›¾å±‚å›¾ç‰‡: {e}")
            return

        # åˆ›å»ºå¯¹è¯æ¡†
        dialog = tk.Toplevel(self.root)
        dialog.title("çº¯è‰²æŠ å›¾")
        dialog.geometry("900x600")
        dialog.transient(self.root)
        dialog.grab_set()

        # ä¸»å®¹å™¨
        main_frame = tk.Frame(dialog)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        # å·¦ä¾§ï¼šå›¾ç‰‡é¢„è§ˆ
        left_frame = tk.Frame(main_frame, bg="#f0f0f0")
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))

        tk.Label(left_frame, text="å›¾ç‰‡é¢„è§ˆ (ç‚¹å‡»å¸å–é¢œè‰²)",
                font=(FONT_FAMILY, 10, "bold"), bg="#f0f0f0").pack(pady=5)

        # è®¡ç®—é¢„è§ˆå›¾å°ºå¯¸
        max_w, max_h = 550, 500
        scale = min(1.0, max_w / img.size[0], max_h / img.size[1])
        disp_w = max(1, int(img.size[0] * scale))
        disp_h = max(1, int(img.size[1] * scale))
        preview_img = img.resize((disp_w, disp_h), Image.Resampling.LANCZOS)
        tk_preview = ImageTk.PhotoImage(preview_img)

        # åˆ›å»ºç”»å¸ƒ
        canvas = tk.Canvas(left_frame, width=disp_w, height=disp_h,
                          bg="#222", highlightthickness=1, highlightbackground="#999")
        canvas.pack(padx=5, pady=5)
        canvas.create_image(0, 0, anchor=tk.NW, image=tk_preview)
        canvas.image = tk_preview

        # å³ä¾§ï¼šæ§åˆ¶é¢æ¿
        right_frame = tk.Frame(main_frame, bg=COLOR_WHITE, width=300)
        right_frame.pack(side=tk.LEFT, fill=tk.BOTH, padx=(10, 0))
        right_frame.pack_propagate(False)

        tk.Label(right_frame, text="çº¯è‰²æŠ å›¾è®¾ç½®", font=(FONT_FAMILY, 12, "bold"),
                bg=COLOR_WHITE).pack(pady=(15, 20))

        # é¢œè‰²é€‰æ‹©
        color_section = tk.LabelFrame(right_frame, text="é¢œè‰²é€‰æ‹©",
                                      font=(FONT_FAMILY, 10, "bold"),
                                      bg=COLOR_WHITE, padx=10, pady=10)
        color_section.pack(fill=tk.X, padx=15, pady=(0, 15))

        selected_color = tk.StringVar(value="#FFFFFF")
        eyedropper_mode = tk.BooleanVar(value=False)

        # é¢œè‰²æ˜¾ç¤º
        color_display_frame = tk.Frame(color_section, bg=COLOR_WHITE)
        color_display_frame.pack(fill=tk.X, pady=5)

        tk.Label(color_display_frame, text="å½“å‰é¢œè‰²:",
                font=(FONT_FAMILY, 9), bg=COLOR_WHITE).pack(side=tk.LEFT)

        color_display = tk.Label(color_display_frame, text="      ",
                                bg="#FFFFFF", relief=tk.RIDGE, borderwidth=2)
        color_display.pack(side=tk.LEFT, padx=10)

        color_label = tk.Label(color_display_frame, text="#FFFFFF",
                              font=(FONT_FAMILY, 9), bg=COLOR_WHITE)
        color_label.pack(side=tk.LEFT)

        def update_color_display(color_hex):
            selected_color.set(color_hex)
            color_display.config(bg=color_hex)
            color_label.config(text=color_hex)

        # æŒ‰é’®è¡Œ
        btn_row = tk.Frame(color_section, bg=COLOR_WHITE)
        btn_row.pack(fill=tk.X, pady=5)

        def choose_color():
            color = colorchooser.askcolor(title="é€‰æ‹©è¦æŠ é™¤çš„é¢œè‰²",
                                         initialcolor=selected_color.get())
            if color[1]:
                update_color_display(color[1])
                eyedropper_mode.set(False)
                eyedropper_btn.config(relief=tk.FLAT, bg="#2196F3")
                canvas.config(cursor="")

        def toggle_eyedropper():
            if eyedropper_mode.get():
                eyedropper_mode.set(False)
                eyedropper_btn.config(relief=tk.FLAT, bg="#2196F3")
                canvas.config(cursor="")
            else:
                eyedropper_mode.set(True)
                eyedropper_btn.config(relief=tk.SUNKEN, bg="#1976D2")
                canvas.config(cursor="crosshair")

        tk.Button(btn_row, text="é€‰æ‹©é¢œè‰²", command=choose_color,
                 bg=COLOR_THEME, fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 9), padx=10, pady=5).pack(side=tk.LEFT, padx=(0, 5))

        eyedropper_btn = tk.Button(btn_row, text="ğŸ¨ å¸ç®¡å·¥å…·", command=toggle_eyedropper,
                                   bg="#2196F3", fg="white", relief=tk.FLAT,
                                   font=(FONT_FAMILY, 9), padx=10, pady=5)
        eyedropper_btn.pack(side=tk.LEFT)

        # å¸ç®¡å·¥å…·ç‚¹å‡»äº‹ä»¶
        def on_canvas_click(event):
            if not eyedropper_mode.get():
                return

            # è·å–ç‚¹å‡»ä½ç½®çš„é¢œè‰²
            try:
                x = int(event.x / scale)
                y = int(event.y / scale)
                if 0 <= x < img.size[0] and 0 <= y < img.size[1]:
                    pixel = img.getpixel((x, y))
                    r, g, b = pixel[0], pixel[1], pixel[2]
                    color_hex = f"#{r:02x}{g:02x}{b:02x}"
                    update_color_display(color_hex)
                    # å–è‰²åè‡ªåŠ¨å…³é—­å¸ç®¡æ¨¡å¼
                    eyedropper_mode.set(False)
                    eyedropper_btn.config(relief=tk.FLAT, bg="#2196F3")
                    canvas.config(cursor="")
            except Exception as e:
                print(f"å¸ç®¡å·¥å…·é”™è¯¯: {e}")

        canvas.bind("<Button-1>", on_canvas_click)

        # å®¹å·®è®¾ç½®
        tolerance_section = tk.LabelFrame(right_frame, text="å®¹å·®è®¾ç½®",
                                         font=(FONT_FAMILY, 10, "bold"),
                                         bg=COLOR_WHITE, padx=10, pady=10)
        tolerance_section.pack(fill=tk.X, padx=15, pady=(0, 15))

        tk.Label(tolerance_section, text="é¢œè‰²å®¹å·® (0-255):",
                font=(FONT_FAMILY, 9), bg=COLOR_WHITE).pack(anchor="w")

        tolerance_var = tk.IntVar(value=30)

        tolerance_frame = tk.Frame(tolerance_section, bg=COLOR_WHITE)
        tolerance_frame.pack(fill=tk.X, pady=5)

        tolerance_scale = tk.Scale(tolerance_frame, from_=0, to=255, orient=tk.HORIZONTAL,
                                   variable=tolerance_var, bg=COLOR_WHITE, length=220)
        tolerance_scale.pack(side=tk.LEFT)

        tolerance_value_label = tk.Label(tolerance_frame, text="30",
                                        font=(FONT_FAMILY, 9, "bold"),
                                        bg=COLOR_WHITE, width=3)
        tolerance_value_label.pack(side=tk.LEFT, padx=5)

        def update_tolerance_label(*args):
            tolerance_value_label.config(text=str(tolerance_var.get()))

        tolerance_var.trace('w', update_tolerance_label)

        tk.Label(tolerance_section, text="å®¹å·®è¶Šå¤§ï¼ŒæŠ é™¤çš„é¢œè‰²èŒƒå›´è¶Šå¹¿",
                font=(FONT_FAMILY, 8), fg="#666", bg=COLOR_WHITE).pack(anchor="w")

        # æŒ‰é’®åŒºåŸŸ
        btn_frame = tk.Frame(right_frame, bg=COLOR_WHITE)
        btn_frame.pack(side=tk.BOTTOM, pady=20)

        def on_apply():
            color_hex = selected_color.get()
            tolerance = tolerance_var.get()

            try:
                # å°†åå…­è¿›åˆ¶é¢œè‰²è½¬æ¢ä¸ºRGB
                color_hex = color_hex.lstrip('#')
                target_r = int(color_hex[0:2], 16)
                target_g = int(color_hex[2:4], 16)
                target_b = int(color_hex[4:6], 16)

                # è½¬æ¢ä¸ºnumpyæ•°ç»„è¿›è¡Œå¤„ç†
                img_array = np.array(img)

                # è®¡ç®—æ¯ä¸ªåƒç´ ä¸ç›®æ ‡é¢œè‰²çš„è·ç¦»
                diff = np.abs(img_array[:, :, 0].astype(int) - target_r) + \
                       np.abs(img_array[:, :, 1].astype(int) - target_g) + \
                       np.abs(img_array[:, :, 2].astype(int) - target_b)

                # åˆ›å»ºmaskï¼šè·ç¦»å°äºå®¹å·®çš„åƒç´ è®¾ä¸ºé€æ˜
                mask = diff <= tolerance * 3  # ä¹˜ä»¥3å› ä¸ºæ˜¯ä¸‰ä¸ªé€šé“çš„æ€»å’Œ
                img_array[mask, 3] = 0  # å°†alphaé€šé“è®¾ä¸º0ï¼ˆé€æ˜ï¼‰

                # è½¬æ¢å›PILå›¾åƒ
                result_img = Image.fromarray(img_array, 'RGBA')

                # ä¿å­˜å¤„ç†åçš„å›¾ç‰‡
                temp_dir = os.path.join(get_base_dir(), "temp_cutout")
                os.makedirs(temp_dir, exist_ok=True)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                new_path = os.path.join(temp_dir, f"cutout_{timestamp}.png")
                result_img.save(new_path)

                # æ›´æ–°å›¾å±‚è·¯å¾„
                self.save_state("layers")
                layer["path"] = new_path

                # åˆ·æ–°æ˜¾ç¤º
                self.update_layer_listbox()
                self.refresh_canvas()
                self.mark_unsaved()

                dialog.destroy()
                messagebox.showinfo("æˆåŠŸ", "çº¯è‰²æŠ å›¾å®Œæˆï¼")

            except Exception as e:
                messagebox.showerror("é”™è¯¯", f"æŠ å›¾å¤±è´¥: {e}")

        def on_cancel():
            dialog.destroy()

        tk.Button(btn_frame, text="åº”ç”¨", command=on_apply,
                 bg="#00897B", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10, "bold"), padx=25, pady=10).pack(pady=5)

        tk.Button(btn_frame, text="å–æ¶ˆ", command=on_cancel,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10), padx=25, pady=10).pack(pady=5)

    def detect_text_in_selected_layers(self):
        """å¯¹é€‰ä¸­çš„å›¾å±‚è¿›è¡Œæ–‡å­—æ£€æµ‹ï¼Œåœ¨å›¾å±‚ä¸Šç›´æ¥æ£€æµ‹æ–‡æœ¬ä½ç½®å¹¶æ·»åŠ æ–‡æœ¬æ¡†åˆ°å½“å‰é¡µé¢"""
        if not self.ocr:
            messagebox.showwarning("æç¤º", "OCRæœªåˆå§‹åŒ–")
            return

        page, layers, layer = self._get_selected_layer()
        if layer is None:
            messagebox.showwarning("æç¤º", "è¯·å…ˆé€‰æ‹©ä¸€ä¸ªå›¾å±‚")
            return

        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨")
            return

        current_page = self.pages[self.current_page_index]

        # è¯¢é—®æ˜¯å¦æ¸…ç©ºç°æœ‰æ¡†
        if current_page.get("text_boxes"):
            result = messagebox.askyesnocancel(
                "æç¤º", "æ˜¯å¦æ¸…ç©ºç°æœ‰æ–‡æœ¬æ¡†ï¼Ÿ\n\næ˜¯ - æ¸…ç©ºåæ£€æµ‹\nå¦ - è¿½åŠ æ£€æµ‹\nå–æ¶ˆ - å–æ¶ˆ"
            )
            if result is None:
                return
            elif result:
                current_page["text_boxes"] = []

        # åˆ›å»ºè¿›åº¦å¯¹è¯æ¡†
        progress_dialog = tk.Toplevel(self.root)
        progress_dialog.title("OCRæ£€æµ‹ä¸­")
        progress_dialog.geometry("400x150")
        progress_dialog.transient(self.root)
        progress_dialog.grab_set()

        tk.Label(progress_dialog, text="æ­£åœ¨æ£€æµ‹å›¾å±‚ä¸­çš„æ–‡å­—åŒºåŸŸ...",
                font=(FONT_FAMILY, 11, "bold")).pack(pady=20)

        progress_label = tk.Label(progress_dialog, text="è¯·ç¨å€™...",
                                 font=(FONT_FAMILY, 9), fg="#666")
        progress_label.pack(pady=10)

        def worker():
            try:
                # è¯»å–å›¾å±‚å›¾ç‰‡
                layer_img = Image.open(path).convert("RGB")

                # è·å–å›¾å±‚åœ¨é¡µé¢ä¸­çš„ä½ç½®
                layer_x = layer.get("x", 0)
                layer_y = layer.get("y", 0)
                layer_scale = layer.get("scale", 1.0)

                # è½¬æ¢ä¸ºOpenCVæ ¼å¼
                img_array = np.array(layer_img)
                img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)

                img_h, img_w = img_bgr.shape[:2]

                # ä¿å­˜ä¸´æ—¶æ–‡ä»¶ä¾›OCRä½¿ç”¨
                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                cv2.imwrite(temp_path, img_bgr)

                # OCRæ£€æµ‹ - ä½¿ç”¨ä¸é¡¶éƒ¨èœå•ç›¸åŒçš„é€»è¾‘
                result = self.ocr.predict(temp_path)
                os.remove(temp_path)

                if not result or len(result) == 0:
                    self.root.after(0, progress_dialog.destroy)
                    self.root.after(0, lambda: messagebox.showinfo("æç¤º", "æœªæ£€æµ‹åˆ°æ–‡å­—"))
                    return

                ocr_result = result[0]
                dt_polys = ocr_result.get('dt_polys', [])
                rec_texts = ocr_result.get('rec_texts', [])

                if not dt_polys:
                    self.root.after(0, progress_dialog.destroy)
                    self.root.after(0, lambda: messagebox.showinfo("æç¤º", "æœªæ£€æµ‹åˆ°æ–‡å­—"))
                    return

                # å°†æ£€æµ‹ç»“æœæ·»åŠ åˆ°å½“å‰é¡µé¢çš„text_boxes
                text_boxes = current_page.get("text_boxes", [])
                added_count = 0

                for i, poly in enumerate(dt_polys):
                    # è®¡ç®—æ–‡æœ¬æ¡†ä½ç½®ï¼ˆç›¸å¯¹äºå›¾å±‚ï¼‰
                    x_coords = [p[0] for p in poly]
                    y_coords = [p[1] for p in poly]

                    x_min = int(min(x_coords))
                    y_min = int(min(y_coords))
                    x_max = int(max(x_coords))
                    y_max = int(max(y_coords))

                    box_w = x_max - x_min
                    box_h = y_max - y_min

                    if box_w < 10 or box_h < 10:
                        continue

                    # è½¬æ¢åˆ°é¡µé¢åæ ‡ï¼ˆè€ƒè™‘å›¾å±‚ä½ç½®å’Œç¼©æ”¾ï¼‰
                    page_x = int(layer_x + x_min * layer_scale)
                    page_y = int(layer_y + y_min * layer_scale)
                    page_w = int(box_w * layer_scale)
                    page_h = int(box_h * layer_scale)

                    # è·å–è¯†åˆ«çš„æ–‡å­—
                    text = rec_texts[i] if i < len(rec_texts) else ""

                    # åˆ›å»ºæ–‡æœ¬æ¡† - ä½¿ç”¨æ ‡å‡†æ ¼å¼
                    font_size = 12
                    if text:
                        font_size = fit_font_size_pt(text, page_w, page_h, editor=self)

                    text_box = {
                        "x": page_x,
                        "y": page_y,
                        "width": page_w,
                        "height": page_h,
                        "text": text,  # å·²ç»æœ‰æ–‡å­—äº†
                        "font_name": "å¾®è½¯é›…é»‘",
                        "font_size": font_size,
                        "font_color": "#000000",
                        "bold": False,
                        "italic": False,
                        "align": "left"
                    }
                    text_boxes.append(text_box)
                    added_count += 1

                current_page["text_boxes"] = text_boxes

                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, self.load_current_page)  # é‡æ–°åŠ è½½é¡µé¢ï¼Œå°†å­—å…¸è½¬æ¢ä¸ºTextBoxå¯¹è±¡
                self.root.after(0, self.mark_unsaved)
                self.root.after(0, lambda c=added_count: messagebox.showinfo(
                    "æˆåŠŸ", f"æ£€æµ‹å¹¶è¯†åˆ«å®Œæˆï¼\nå…±æ£€æµ‹åˆ° {c} ä¸ªæ–‡æœ¬æ¡†\næ–‡å­—å·²è‡ªåŠ¨è¯†åˆ«"))

            except Exception as e:
                import traceback
                traceback.print_exc()
                err_text = str(e)
                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"OCRæ£€æµ‹å¤±è´¥:\n{t}"))

        threading.Thread(target=worker, daemon=True).start()

    def remove_text_background_from_layer(self):
        """å¯¹é€‰ä¸­çš„å›¾å±‚å»é™¤æ–‡æœ¬èƒŒæ™¯"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        page, layers, layer = self._get_selected_layer()
        if layer is None:
            messagebox.showwarning("æç¤º", "è¯·å…ˆé€‰æ‹©ä¸€ä¸ªå›¾å±‚")
            return

        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨")
            return

        current_page = self.pages[self.current_page_index]
        text_boxes = current_page.get("text_boxes", [])

        if not text_boxes:
            messagebox.showwarning("æç¤º", "å½“å‰é¡µæ²¡æœ‰æ–‡æœ¬æ¡†\n\nè¯·å…ˆä½¿ç”¨ã€Œæ£€æµ‹ã€åŠŸèƒ½è¯†åˆ«æ–‡æœ¬åŒºåŸŸ")
            return

        if not self.config.get("inpaint_enabled", True):
            messagebox.showwarning("æç¤º", "èƒŒæ™¯ç”ŸæˆåŠŸèƒ½å·²ç¦ç”¨\n\nè¯·åœ¨è®¾ç½®ä¸­å¯ç”¨")
            return

        # ç¡®è®¤å¯¹è¯æ¡†
        result = messagebox.askyesno(
            "ç¡®è®¤",
            f"å³å°†å¯¹é€‰ä¸­å›¾å±‚è¿›è¡Œå»å­—å¤„ç†\n\n"
            f"å½“å‰é¡µæœ‰ {len(text_boxes)} ä¸ªæ–‡æœ¬æ¡†\n"
            "ç³»ç»Ÿå°†è‡ªåŠ¨å¯¹è¿™äº›æ–‡å­—åŒºåŸŸè¿›è¡Œä¿®å¤\n\n"
            "æç¤ºï¼šç»“æœä¼šä½œä¸ºæ–°å›¾å±‚å åŠ ï¼ŒåŸå›¾å±‚ä¸ä¼šè¢«ä¿®æ”¹\n\n"
            "æ­¤æ“ä½œéœ€è¦è°ƒç”¨ IOPaint API æœåŠ¡\n"
            "å¤„ç†æ—¶é—´çº¦ 5-30 ç§’\n\n"
            "æ˜¯å¦ç»§ç»­ï¼Ÿ",
        )

        if not result:
            return

        # ä¿å­˜å›¾å±‚å¿«ç…§ä¾¿äºæ’¤é”€
        self.save_state("layers")

        self.update_status("æ­£åœ¨å¯¹å›¾å±‚è¿›è¡Œå»å­—å¤„ç†...")

        def generate_bg():
            try:
                # è¯»å–å›¾å±‚å›¾ç‰‡
                layer_img = Image.open(path).convert("RGB")
                layer_x = layer.get("x", 0)
                layer_y = layer.get("y", 0)
                layer_scale = layer.get("scale", 1.0)

                # åˆ›å»ºè’™ç‰ˆ - éœ€è¦å°†é¡µé¢åæ ‡çš„text_boxesè½¬æ¢åˆ°å›¾å±‚åæ ‡
                self.root.after(0, lambda: self.update_status("æ­£åœ¨åˆ›å»ºè’™ç‰ˆ..."))

                # åˆ›å»ºå›¾å±‚å¤§å°çš„è’™ç‰ˆ
                mask = Image.new("L", layer_img.size, 0)  # å…¨é»‘èƒŒæ™¯
                draw = ImageDraw.Draw(mask)

                for box in text_boxes:
                    # å°†é¡µé¢åæ ‡è½¬æ¢å›å›¾å±‚åæ ‡
                    # page_x = layer_x + layer_img_x * layer_scale
                    # => layer_img_x = (page_x - layer_x) / layer_scale

                    box_x_on_layer = (box["x"] - layer_x) / layer_scale
                    box_y_on_layer = (box["y"] - layer_y) / layer_scale
                    box_w_on_layer = box["width"] / layer_scale
                    box_h_on_layer = box["height"] / layer_scale

                    # æ£€æŸ¥æ–‡æœ¬æ¡†æ˜¯å¦åœ¨å›¾å±‚èŒƒå›´å†…
                    if (box_x_on_layer + box_w_on_layer < 0 or box_x_on_layer > layer_img.size[0] or
                        box_y_on_layer + box_h_on_layer < 0 or box_y_on_layer > layer_img.size[1]):
                        continue

                    # ç¨å¾®æ‰©å¤§æ–‡æœ¬æ¡†åŒºåŸŸ
                    padding = 5
                    x1 = max(0, int(box_x_on_layer - padding))
                    y1 = max(0, int(box_y_on_layer - padding))
                    x2 = min(layer_img.size[0], int(box_x_on_layer + box_w_on_layer + padding))
                    y2 = min(layer_img.size[1], int(box_y_on_layer + box_h_on_layer + padding))

                    # æ ‡è®°ä¸ºç™½è‰²ï¼ˆéœ€è¦ä¿®å¤ï¼‰
                    draw.rectangle([x1, y1, x2, y2], fill=255)

                # è°ƒç”¨APIä¿®å¤
                self.root.after(0, lambda: self.update_status("æ­£åœ¨è°ƒç”¨IOPaint APIä¿®å¤..."))
                result_img = self.call_inpaint_api(layer_img, mask)

                if result_img:
                    # å°†ä¿®å¤åçš„å›¾ç‰‡ä½œä¸ºæ–°å›¾å±‚æ·»åŠ 
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    layer_name = f"å»å­—_{os.path.splitext(os.path.basename(path))[0]}_{timestamp}"

                    # ä¿å­˜ä¿®å¤åçš„å›¾ç‰‡
                    temp_dir = os.path.join(get_base_dir(), "temp_inpaint")
                    os.makedirs(temp_dir, exist_ok=True)
                    new_path = os.path.join(temp_dir, f"{layer_name}.png")
                    result_img.save(new_path)

                    # æ·»åŠ ä¸ºæ–°å›¾å±‚ï¼Œç»§æ‰¿åŸå›¾å±‚çš„ä½ç½®å’Œç¼©æ”¾
                    new_layer = self.add_image_layer(
                        current_page,
                        result_img.convert("RGBA"),
                        name=layer_name,
                        x=layer_x,
                        y=layer_y,
                        opacity=1.0,
                        visible=True
                    )
                    if new_layer:
                        new_layer["scale"] = layer_scale
                        new_layer["locked"] = False

                    self.root.after(0, self.update_thumbnails)
                    self.root.after(0, self.update_layer_listbox)
                    self.root.after(0, self.scroll_to_layers)
                    if new_layer and new_layer.get("id"):
                        self.root.after(0, lambda lid=new_layer["id"]: self.select_layer_by_id(lid))
                    self.root.after(0, self.refresh_canvas)
                    self.root.after(0, self.mark_unsaved)

                    self.root.after(0, lambda: self.update_status(f"å·²ç”Ÿæˆå»å­—å›¾å±‚ï¼š{layer_name}"))
                    self.root.after(
                        0,
                        lambda: messagebox.showinfo(
                            "å®Œæˆ",
                            "å»å­—å¤„ç†å®Œæˆï¼\n\n"
                            f"å·²å»é™¤ {len(text_boxes)} ä¸ªæ–‡å­—åŒºåŸŸ\n"
                            "ç»“æœå·²ä½œä¸ºæ–°å›¾å±‚å åŠ ï¼ˆå³ä¾§å›¾å±‚é¢æ¿å¯è§ï¼‰\n\n"
                            "æç¤ºï¼šCtrl+Z å¯ä»¥æ’¤é”€",
                        ),
                    )
                else:
                    self.root.after(0, lambda: self.update_status("å»å­—å¤„ç†å¤±è´¥"))
                    self.root.after(0, lambda: messagebox.showerror("é”™è¯¯", "IOPaint API è°ƒç”¨å¤±è´¥"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(error_msg)
                self.root.after(0, lambda: self.update_status("å»å­—å¤„ç†å‡ºé”™"))
                self.root.after(0, lambda: messagebox.showerror("é”™è¯¯", f"å»å­—å¤„ç†å¤±è´¥:\n{str(e)}"))

        threading.Thread(target=generate_bg, daemon=True).start()

    def recognize_text_in_selected_layers(self):
        """å¯¹å½“å‰é¡µé¢ä¸­çš„ç©ºæ–‡æœ¬æ¡†è¿›è¡ŒOCRè¯†åˆ«"""
        if not self.ocr:
            messagebox.showwarning("æç¤º", "OCRæœªåˆå§‹åŒ–")
            return

        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        current_page = self.pages[self.current_page_index]
        text_boxes = current_page.get("text_boxes", [])

        if not text_boxes:
            messagebox.showwarning("æç¤º", "è¯·å…ˆæ£€æµ‹æ–‡æœ¬æ¡†")
            return

        # ç»Ÿè®¡ç©ºæ–‡æœ¬æ¡†
        empty_boxes = [box for box in text_boxes if not box.get("text")]

        if not empty_boxes:
            messagebox.showinfo("æç¤º", "æ‰€æœ‰æ–‡æœ¬æ¡†éƒ½å·²æœ‰æ–‡å­—ï¼Œæ— éœ€è¯†åˆ«")
            return

        # åˆ›å»ºè¿›åº¦å¯¹è¯æ¡†
        progress_dialog = tk.Toplevel(self.root)
        progress_dialog.title("OCRè¯†åˆ«ä¸­")
        progress_dialog.geometry("400x150")
        progress_dialog.transient(self.root)
        progress_dialog.grab_set()

        tk.Label(progress_dialog, text=f"æ­£åœ¨è¯†åˆ« {len(empty_boxes)} ä¸ªæ–‡æœ¬æ¡†...",
                font=(FONT_FAMILY, 11, "bold")).pack(pady=20)

        progress_label = tk.Label(progress_dialog, text="è¯·ç¨å€™...",
                                 font=(FONT_FAMILY, 9), fg="#666")
        progress_label.pack(pady=10)

        def worker():
            try:
                # ä½¿ç”¨å½“å‰é¡µçš„ç¼–è¾‘å›¾ç‰‡
                page_img = current_page["image"]
                img = np.array(page_img)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]
                recognized_count = 0

                for box_data in empty_boxes:
                    if box_data.get("text"):
                        continue

                    x, y, w, h = box_data["x"], box_data["y"], box_data["width"], box_data["height"]
                    expand_h, expand_w = int(h * 0.3), int(w * 0.1)

                    crop_x = max(0, x - expand_w)
                    crop_y = max(0, y - expand_h)
                    crop_x2 = min(x + w + expand_w, img_w)
                    crop_y2 = min(y + h + expand_h, img_h)

                    cropped = img[crop_y:crop_y2, crop_x:crop_x2]

                    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_path = temp_file.name
                    temp_file.close()
                    cv2.imwrite(temp_path, cropped)

                    try:
                        result = self.ocr.predict(temp_path)
                        os.remove(temp_path)

                        if result and len(result) > 0:
                            ocr_result = result[0]
                            rec_texts = ocr_result.get('rec_texts', [])
                            if rec_texts:
                                box_data["text"] = ''.join(rec_texts)
                                if box_data["text"]:
                                    box_data["font_size"] = fit_font_size_pt(
                                        box_data["text"],
                                        w,
                                        h,
                                        editor=self,
                                        font_name=box_data.get("font_name"),
                                    )
                                    recognized_count += 1
                    except:
                        try:
                            os.remove(temp_path)
                        except:
                            pass

                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, self.load_current_page)  # é‡æ–°åŠ è½½é¡µé¢ï¼Œæ›´æ–°TextBoxå¯¹è±¡
                self.root.after(0, self.mark_unsaved)
                self.root.after(0, lambda c=recognized_count: messagebox.showinfo(
                    "æˆåŠŸ", f"è¯†åˆ«å®Œæˆï¼\næˆåŠŸè¯†åˆ« {c} ä¸ªæ–‡æœ¬æ¡†"))

            except Exception as e:
                import traceback
                traceback.print_exc()
                err_text = str(e)
                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"OCRè¯†åˆ«å¤±è´¥:\n{t}"))

        threading.Thread(target=worker, daemon=True).start()

    def ocr_selected_layer(self):
        """å¯¹é€‰ä¸­çš„å›¾å±‚è¿›è¡ŒOCRè¯†åˆ«ï¼Œå°†è¯†åˆ«ç»“æœæ·»åŠ åˆ°å½“å‰é¡µé¢ï¼ˆå·²å¼ƒç”¨ï¼Œä¿ç•™å…¼å®¹æ€§ï¼‰"""
        self.detect_text_in_selected_layers()

    def reset_selected_layer_crop(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        if layer.get("locked"):
            self.update_status("å›¾å±‚å·²é”å®šï¼Œæ— æ³•ä¿®æ”¹ä½ç½®/ç¼©æ”¾/è£å‰ª")
            return
        if not layer.get("crop"):
            return
        self.save_state("layers")
        layer["crop"] = None
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def crop_selected_layer(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return
        if layer.get("locked"):
            self.update_status("å›¾å±‚å·²é”å®šï¼Œæ— æ³•ä¿®æ”¹ä½ç½®/ç¼©æ”¾/è£å‰ª")
            return
        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨")
            return

        try:
            src = Image.open(path).convert("RGBA")
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•æ‰“å¼€å›¾å±‚å›¾ç‰‡: {e}")
            return

        win = tk.Toplevel(self.root)
        win.title("è£å‰ªå›¾å±‚")
        win.geometry("900x700")
        win.transient(self.root)

        max_w, max_h = 860, 560
        scale = min(1.0, max_w / src.size[0], max_h / src.size[1])
        disp_w = max(1, int(src.size[0] * scale))
        disp_h = max(1, int(src.size[1] * scale))
        disp = src.resize((disp_w, disp_h), Image.Resampling.LANCZOS)
        tk_img = ImageTk.PhotoImage(disp)

        canvas = tk.Canvas(win, width=disp_w, height=disp_h, bg="#222", highlightthickness=0)
        canvas.pack(padx=10, pady=10)
        canvas.create_image(0, 0, anchor=tk.NW, image=tk_img)
        canvas.image = tk_img

        rect_id = None
        start = {"x": 0, "y": 0}
        current = {"x0": 0, "y0": 0, "x1": disp_w, "y1": disp_h}

        # åˆå§‹åŒ–ä¸ºå·²æœ‰ crop æˆ–å…¨å›¾
        crop = layer.get("crop")
        if crop:
            try:
                if isinstance(crop, dict):
                    x0 = int(crop.get("x0", 0))
                    y0 = int(crop.get("y0", 0))
                    x1 = int(crop.get("x1", src.size[0]))
                    y1 = int(crop.get("y1", src.size[1]))
                else:
                    x0, y0, x1, y1 = [int(v) for v in crop]
                current["x0"] = int(x0 * scale)
                current["y0"] = int(y0 * scale)
                current["x1"] = int(x1 * scale)
                current["y1"] = int(y1 * scale)
            except Exception:
                pass

        rect_id = canvas.create_rectangle(
            current["x0"], current["y0"], current["x1"], current["y1"], outline="#00E5FF", width=2
        )

        def on_press(ev):
            start["x"], start["y"] = ev.x, ev.y
            current["x0"], current["y0"] = ev.x, ev.y
            current["x1"], current["y1"] = ev.x, ev.y
            canvas.coords(rect_id, ev.x, ev.y, ev.x, ev.y)

        def on_drag(ev):
            x0 = min(start["x"], ev.x)
            y0 = min(start["y"], ev.y)
            x1 = max(start["x"], ev.x)
            y1 = max(start["y"], ev.y)
            x0 = max(0, min(disp_w, x0))
            y0 = max(0, min(disp_h, y0))
            x1 = max(0, min(disp_w, x1))
            y1 = max(0, min(disp_h, y1))
            current["x0"], current["y0"], current["x1"], current["y1"] = x0, y0, x1, y1
            canvas.coords(rect_id, x0, y0, x1, y1)

        canvas.bind("<ButtonPress-1>", on_press)
        canvas.bind("<B1-Motion>", on_drag)

        btn_row = tk.Frame(win, bg=COLOR_WHITE)
        btn_row.pack(fill=tk.X, padx=10, pady=(0, 10))

        def apply_crop():
            x0, y0, x1, y1 = current["x0"], current["y0"], current["x1"], current["y1"]
            if x1 - x0 < 2 or y1 - y0 < 2:
                win.destroy()
                return
            rx0 = int(round(x0 / scale))
            ry0 = int(round(y0 / scale))
            rx1 = int(round(x1 / scale))
            ry1 = int(round(y1 / scale))
            rx0 = max(0, min(src.size[0], rx0))
            ry0 = max(0, min(src.size[1], ry0))
            rx1 = max(0, min(src.size[0], rx1))
            ry1 = max(0, min(src.size[1], ry1))
            if rx1 <= rx0 or ry1 <= ry0:
                win.destroy()
                return
            self.save_state("layers")
            layer["crop"] = {"x0": rx0, "y0": ry0, "x1": rx1, "y1": ry1}
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()
            win.destroy()

        tk.Button(btn_row, text="åº”ç”¨è£å‰ª", command=apply_crop, bg=COLOR_GREEN, fg="white",
                  font=(FONT_FAMILY, 9), cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=4)
        tk.Button(btn_row, text="å–æ¶ˆ", command=win.destroy, bg="#757575", fg="white",
                  font=(FONT_FAMILY, 9), cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=4)

    def import_layer_from_file(self):
        """ä»æœ¬åœ°å¯¼å…¥å›¾ç‰‡ï¼ˆåŒ…å« SVGï¼‰ä½œä¸ºæ–°å›¾å±‚ã€‚"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        file_path = filedialog.askopenfilename(
            title="é€‰æ‹©è¦å¯¼å…¥åˆ°å›¾å±‚çš„å›¾ç‰‡",
            filetypes=[
                ("å›¾ç‰‡æ–‡ä»¶", "*.png *.jpg *.jpeg *.bmp *.webp *.gif *.tif *.tiff *.svg"),
                ("æ‰€æœ‰æ–‡ä»¶", "*.*"),
            ],
        )
        if not file_path:
            return

        ext = os.path.splitext(file_path)[1].lower()
        img = None
        try:
            if ext == ".svg":
                try:
                    import cairosvg  # type: ignore
                except Exception:
                    messagebox.showerror(
                        "ç¼ºå°‘ä¾èµ–",
                        "å½“å‰ç¯å¢ƒæœªå®‰è£… SVG æ¸²æŸ“ä¾èµ–ï¼Œæ— æ³•ç›´æ¥å¯¼å…¥ SVGã€‚\n\n"
                        "å¯é€‰æ–¹æ¡ˆï¼š\n"
                        "1) pip install cairosvg\n"
                        "2) å…ˆæŠŠ SVG å¯¼å‡ºä¸º PNG å†å¯¼å…¥\n",
                    )
                    return
                from io import BytesIO

                png_bytes = cairosvg.svg2png(url=file_path)
                img = Image.open(BytesIO(png_bytes)).convert("RGBA")
            else:
                img = Image.open(file_path)
                if "A" in img.getbands():
                    img = img.convert("RGBA")
                else:
                    img = img.convert("RGB").convert("RGBA")
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•å¯¼å…¥å›¾ç‰‡: {e}")
            return

        page = self.pages[self.current_page_index]
        base = page.get("image") or self.original_image
        if base is None:
            return
        base_w, base_h = base.size

        # é»˜è®¤ç¼©æ”¾åˆ°ç”»é¢å†…ï¼ˆä¸æ”¾å¤§ï¼‰
        s = 1.0
        try:
            s = min(1.0, base_w / max(1, img.size[0]), base_h / max(1, img.size[1])) * 0.9
            s = max(0.05, min(1.0, s))
        except Exception:
            s = 1.0

        x = int((base_w - img.size[0] * s) / 2)
        y = int((base_h - img.size[1] * s) / 2)

        name = os.path.splitext(os.path.basename(file_path))[0] or "å¯¼å…¥å›¾å±‚"

        self.save_state("layers")
        layer = self.add_image_layer(page, img, name=name, x=x, y=y, opacity=1.0, visible=True)
        if layer is not None:
            layer["scale"] = float(s)
            layer["crop"] = None
            layer["locked"] = True  # é»˜è®¤é”å®šå›¾å±‚ï¼Œé˜²æ­¢è¯¯æ“ä½œ

        self.update_layer_listbox()
        self.scroll_to_layers()
        if layer and layer.get("id"):
            self.select_layer_by_id(layer["id"])
        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"å·²å¯¼å…¥å›¾å±‚: {name}")

    def on_layer_tree_click(self, event):
        # ç‚¹å‡»â€œæ˜¾â€åˆ—ï¼šå¿«é€Ÿæ˜¾ç¤º/éšè—
        if not hasattr(self, "layer_tree"):
            return
        region = self.layer_tree.identify("region", event.x, event.y)
        if region != "cell":
            return
        row = self.layer_tree.identify_row(event.y)
        col = self.layer_tree.identify_column(event.x)
        if not row:
            return
        if col == "#1":  # vis åˆ—
            try:
                self.layer_tree.selection_set(row)
                self.layer_tree.focus(row)
            except Exception:
                pass
            self.on_layer_select()
            self.toggle_selected_layer()
            return "break"
        if col == "#3":  # lock åˆ—
            try:
                self.layer_tree.selection_set(row)
                self.layer_tree.focus(row)
            except Exception:
                pass
            self.on_layer_select()
            self.toggle_selected_layer_lock()
            return "break"

    def on_layer_drag_start(self, event):
        if not hasattr(self, "layer_tree"):
            return
        region = self.layer_tree.identify("region", event.x, event.y)
        if region != "cell":
            return
        col = self.layer_tree.identify_column(event.x)
        if col in ("#1", "#3"):  # vis/lock åˆ—ç‚¹å‡»ä¸è§¦å‘æ‹–æ‹½
            return
        self._layer_drag_iid = self.layer_tree.identify_row(event.y)
        self._layer_drag_moved = False

    def on_layer_drag_motion(self, event):
        if not hasattr(self, "layer_tree"):
            return
        dragged = getattr(self, "_layer_drag_iid", None)
        if not dragged:
            return
        target = self.layer_tree.identify_row(event.y)
        if not target or target == dragged:
            return
        try:
            target_index = self.layer_tree.index(target)
            self.layer_tree.move(dragged, "", target_index)
            self._layer_drag_moved = True
        except Exception:
            pass

    def on_layer_drag_release(self, event):
        if not hasattr(self, "layer_tree") or not self.pages:
            self._layer_drag_iid = None
            self._layer_drag_moved = False
            return
        if not getattr(self, "_layer_drag_moved", False):
            self._layer_drag_iid = None
            return

        page = self.pages[self.current_page_index]
        layers = page.get("layers", [])
        if not layers:
            self._layer_drag_iid = None
            self._layer_drag_moved = False
            return

        new_order = list(self.layer_tree.get_children(""))
        old_order = [layer.get("id") for layer in layers if layer]
        if new_order and old_order and new_order != old_order:
            self.save_state("layers")
            layer_map = {layer.get("id"): layer for layer in layers if layer and layer.get("id")}
            rebuilt = [layer_map[iid] for iid in new_order if iid in layer_map]
            page["layers"] = rebuilt
            self.layers = page.get("layers", [])
            # ä¿æŒé€‰æ‹©åŒæ­¥
            try:
                sel = self.layer_tree.selection()
                if sel:
                    selected_iid = sel[0]
                    for i, layer in enumerate(rebuilt):
                        if layer and layer.get("id") == selected_iid:
                            self.selected_layer_index = i
                            break
            except Exception:
                pass
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()

        self._layer_drag_iid = None
        self._layer_drag_moved = False

    def preview_selected_layer(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return

        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨")
            return

        try:
            overlay = Image.open(path).convert("RGBA")
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•æ‰“å¼€å›¾å±‚å›¾ç‰‡: {e}")
            return

        # é¢„è§ˆæ—¶åº”ç”¨é€æ˜åº¦/è£å‰ª/ç¼©æ”¾æ•ˆæœï¼ˆæ›´ç›´è§‚ï¼‰
        crop = layer.get("crop")
        if crop:
            try:
                if isinstance(crop, dict):
                    x0 = int(crop.get("x0", 0))
                    y0 = int(crop.get("y0", 0))
                    x1 = int(crop.get("x1", overlay.size[0]))
                    y1 = int(crop.get("y1", overlay.size[1]))
                else:
                    x0, y0, x1, y1 = [int(v) for v in crop]
                x0 = max(0, min(overlay.size[0], x0))
                y0 = max(0, min(overlay.size[1], y0))
                x1 = max(0, min(overlay.size[0], x1))
                y1 = max(0, min(overlay.size[1], y1))
                if x1 > x0 and y1 > y0:
                    overlay = overlay.crop((x0, y0, x1, y1))
            except Exception:
                pass

        try:
            scale = float(layer.get("scale", 1.0))
        except Exception:
            scale = 1.0
        if scale <= 0:
            scale = 1.0
        if abs(scale - 1.0) > 1e-6:
            try:
                overlay = overlay.resize(
                    (max(1, int(round(overlay.size[0] * scale))), max(1, int(round(overlay.size[1] * scale)))),
                    Image.Resampling.LANCZOS,
                )
            except Exception:
                pass

        opacity = float(layer.get("opacity", 1.0))
        opacity = max(0.0, min(opacity, 1.0))
        if opacity < 1.0:
            r, g, b, a = overlay.split()
            a = a.point(lambda v: int(v * opacity))
            overlay = Image.merge("RGBA", (r, g, b, a))

        win = tk.Toplevel(self.root)
        win.title(layer.get("name") or "å›¾å±‚é¢„è§ˆ")
        win.transient(self.root)

        img = overlay.copy()
        img.thumbnail((820, 600), Image.Resampling.LANCZOS)
        tk_img = ImageTk.PhotoImage(img)

        label = tk.Label(win, image=tk_img, bg="white")
        label.image = tk_img
        label.pack(padx=10, pady=10)

    def rename_selected_layer(self):
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return

        current_name = layer.get("name") or "å›¾å±‚"
        new_name = simpledialog.askstring("é‡å‘½åå›¾å±‚", "è¯·è¾“å…¥æ–°åç§°ï¼š", initialvalue=current_name, parent=self.root)
        if not new_name:
            return
        self.save_state("layers")
        layer["name"] = new_name.strip()
        self.update_layer_listbox()
        self.mark_unsaved()

    def set_selected_layer_mask_from_file(self):
        # å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        messagebox.showinfo("æç¤º", "å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆé¿å…å¡æ­»ï¼‰ã€‚")
        return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return

        file_path = filedialog.askopenfilename(
            title="é€‰æ‹©è’™ç‰ˆå›¾ç‰‡ï¼ˆç™½=æ˜¾ç¤ºï¼Œé»‘=éšè—ï¼‰",
            filetypes=[("å›¾ç‰‡æ–‡ä»¶", "*.png *.jpg *.jpeg *.bmp")],
        )
        if not file_path:
            return

        try:
            src = Image.open(file_path)
            # ä¼˜å…ˆä½¿ç”¨ alphaï¼ˆå¾ˆå¤šè’™ç‰ˆ PNG ç”¨é€æ˜åº¦è¡¨è¾¾ï¼‰
            if "A" in src.getbands():
                mask_img = src.convert("RGBA").split()[-1]
            else:
                mask_img = src.convert("L")
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•æ‰“å¼€è’™ç‰ˆå›¾ç‰‡: {e}")
            return

        temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
        os.makedirs(temp_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        layer_id = (layer.get("id") or uuid.uuid4().hex[:10])[:10]
        mask_path = os.path.join(temp_dir, f"mask_{layer_id}_{timestamp}.png")

        try:
            mask_img.save(mask_path)
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"ä¿å­˜è’™ç‰ˆå¤±è´¥: {e}")
            return

        self.save_state("layers")
        layer["mask_path"] = mask_path
        layer.setdefault("mask_invert", False)
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def clear_selected_layer_mask(self):
        # å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        messagebox.showinfo("æç¤º", "å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆé¿å…å¡æ­»ï¼‰ã€‚")
        return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return

        if "mask_path" in layer:
            self.save_state("layers")
            layer.pop("mask_path", None)
            layer.pop("mask_invert", None)
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()

    def invert_selected_layer_mask(self):
        # å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        messagebox.showinfo("æç¤º", "å›¾å±‚è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤ï¼ˆé¿å…å¡æ­»ï¼‰ã€‚")
        return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            return

        if not layer.get("mask_path") or not os.path.exists(layer.get("mask_path")):
            messagebox.showwarning("æç¤º", "å½“å‰å›¾å±‚æ²¡æœ‰è’™ç‰ˆ")
            return

        self.save_state("layers")
        layer["mask_invert"] = not bool(layer.get("mask_invert", False))
        self.update_layer_listbox()
        self.refresh_canvas()
        self.mark_unsaved()

    def enter_mask_edit_mode(self):
        # å›¾å±‚è’™ç‰ˆç¼–è¾‘åŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        try:
            self.layer_mask_edit_mode = False
            self.canvas.delete("mask_overlay")
            self.canvas.delete("mask_edit_rect")
        except Exception:
            pass
        messagebox.showinfo("æç¤º", "å›¾å±‚è’™ç‰ˆç¼–è¾‘å·²ç§»é™¤ï¼ˆé¿å…å¡æ­»ï¼‰ã€‚")
        return
        page, layers, layer = self._get_selected_layer()
        if layer is None:
            messagebox.showwarning("æç¤º", "è¯·å…ˆé€‰æ‹©ä¸€ä¸ªå›¾å±‚")
            return

        # åç»­å¯èƒ½ä¼šåˆ›å»º mask_path / ä¿®æ­£å°ºå¯¸ï¼Œå…ˆè®°å½•å›¾å±‚å¿«ç…§ç”¨äºæ’¤é”€
        self.save_state("layers")

        path = layer.get("path")
        if not path or not os.path.exists(path):
            messagebox.showwarning("æç¤º", "å›¾å±‚æ–‡ä»¶ä¸å­˜åœ¨ï¼Œæ— æ³•ç¼–è¾‘è’™ç‰ˆ")
            return

        try:
            # åªå–å°ºå¯¸ï¼Œé¿å…å¤§å›¾ convert å¯¼è‡´ç•Œé¢å‡æ­»
            overlay_size = Image.open(path).size
        except Exception as e:
            messagebox.showerror("é”™è¯¯", f"æ— æ³•æ‰“å¼€å›¾å±‚å›¾ç‰‡: {e}")
            return

        # å…³é—­å¯èƒ½å†²çªçš„æ¨¡å¼
        self.inpaint_mode = False
        self.ai_replace_mode = False
        self.draw_mode = False
        self.is_drawing = False
        self.is_dragging = False
        self.is_resizing = False
        self.is_selecting = False
        self.canvas.delete("temp_rect")
        self.canvas.delete("selection_rect")

        # å‡†å¤‡è’™ç‰ˆï¼ˆé»˜è®¤å…¨ç™½ = å…¨æ˜¾ç¤ºï¼‰
        temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
        os.makedirs(temp_dir, exist_ok=True)

        if not layer.get("id"):
            layer["id"] = uuid.uuid4().hex[:10]

        mask_path = layer.get("mask_path")
        if not mask_path or not os.path.exists(mask_path):
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            mask_path = os.path.join(temp_dir, f"mask_{layer['id']}_{timestamp}.png")
            Image.new("L", overlay_size, 255).save(mask_path)
            layer["mask_path"] = mask_path
            layer.setdefault("mask_invert", False)
        else:
            try:
                existing = Image.open(mask_path).convert("L")
                if existing.size != overlay_size:
                    existing = existing.resize(overlay_size, Image.Resampling.NEAREST)
                    existing.save(mask_path)
            except Exception:
                Image.new("L", overlay_size, 255).save(mask_path)
                layer["mask_path"] = mask_path
                layer.setdefault("mask_invert", False)

        try:
            self._mask_edit_mask = Image.open(mask_path).convert("L")
        except Exception:
            self._mask_edit_mask = Image.new("L", overlay_size, 255)

        self._mask_edit_draw = ImageDraw.Draw(self._mask_edit_mask)
        self._mask_edit_layer_id = layer.get("id")
        self._mask_edit_overlay_size = overlay_size
        self._mask_last_point = None
        self._mask_rect_id = None
        self._mask_rect_start = None
        self._mask_dirty = False
        self._mask_refresh_after = None

        self.layer_mask_edit_mode = True
        self.mask_edit_mode_var.set(True)
        self.update_status("è’™ç‰ˆç¼–è¾‘ï¼šå·¦é”®åˆ·/æ¡†é€‰ï¼ˆå›¾å±‚é¢æ¿å¯åˆ‡æ¢å·¥å…·/æ¨¡å¼ï¼‰")
        self.update_layer_listbox()
        self.refresh_canvas()

    def exit_mask_edit_mode(self):
        # å›¾å±‚è’™ç‰ˆç¼–è¾‘åŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        try:
            self.layer_mask_edit_mode = False
            self.canvas.delete("mask_overlay")
            self.canvas.delete("mask_edit_rect")
        except Exception:
            pass
        return
        if not getattr(self, "layer_mask_edit_mode", False):
            return

        self._mask_save_current()
        self.layer_mask_edit_mode = False
        try:
            self.mask_edit_mode_var.set(False)
        except Exception:
            pass
        self._mask_last_point = None
        self._mask_rect_start = None
        if self._mask_rect_id:
            try:
                self.canvas.delete(self._mask_rect_id)
            except Exception:
                pass
        self._mask_rect_id = None
        self.update_layer_listbox()
        self.refresh_canvas()
        if getattr(self, "_mask_dirty", False):
            self.mark_unsaved()
        self.update_status("å·²é€€å‡ºè’™ç‰ˆç¼–è¾‘")

    def _mask_save_current(self):
        try:
            if not self.pages:
                return
            page = self.pages[self.current_page_index]
            layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
            if layer is None:
                return
            mask_path = layer.get("mask_path")
            if not mask_path:
                return
            if hasattr(self, "_mask_edit_mask") and self._mask_edit_mask is not None:
                self._mask_edit_mask.save(mask_path)
        except Exception:
            pass

    def _mask_get_layer_by_id(self, page, layer_id):
        layers = page.get("layers", [])
        for layer in layers:
            if layer and layer.get("id") == layer_id:
                return layer
        return None

    def _mask_paint_value(self):
        try:
            mode = self.mask_paint_var.get()
        except Exception:
            mode = "hide"
        return 255 if mode == "show" else 0

    def _mask_on_press(self, canvas_x, canvas_y, img_x, img_y):
        if not self.pages:
            return
        page = self.pages[self.current_page_index]
        layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
        if layer is None:
            return

        tool = self.mask_tool_var.get() if hasattr(self, "mask_tool_var") else "brush"

        if tool == "rect":
            self._mask_rect_start = (img_x, img_y)
            if self._mask_rect_id:
                self.canvas.delete(self._mask_rect_id)
            x0 = int(canvas_x)
            y0 = int(canvas_y)
            self._mask_rect_id = self.canvas.create_rectangle(
                x0, y0, x0, y0,
                outline="#FF1744",
                width=2,
                dash=(4, 2),
                tags="mask_edit_rect",
            )
            return

        # brush
        self._mask_last_point = (img_x, img_y)
        self._mask_brush_paint_segment(img_x, img_y, img_x, img_y)
        self._mask_schedule_refresh()

    def _mask_on_drag(self, canvas_x, canvas_y, img_x, img_y):
        if not self.pages:
            return
        page = self.pages[self.current_page_index]
        layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
        if layer is None:
            return

        tool = self.mask_tool_var.get() if hasattr(self, "mask_tool_var") else "brush"

        if tool == "rect":
            if not self._mask_rect_id or not self._mask_rect_start:
                return
            start_canvas_x = (self._mask_rect_start[0] * self.scale) + getattr(self, 'canvas_offset_x', 0)
            start_canvas_y = (self._mask_rect_start[1] * self.scale) + getattr(self, 'canvas_offset_y', 0)
            self.canvas.coords(self._mask_rect_id, start_canvas_x, start_canvas_y, canvas_x, canvas_y)
            return

        # brush
        if not self._mask_last_point:
            self._mask_last_point = (img_x, img_y)
        last_x, last_y = self._mask_last_point
        self._mask_brush_paint_segment(last_x, last_y, img_x, img_y)
        self._mask_last_point = (img_x, img_y)
        self._mask_schedule_refresh()

    def _mask_on_release(self, canvas_x, canvas_y, img_x, img_y):
        if not self.pages:
            return
        page = self.pages[self.current_page_index]
        layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
        if layer is None:
            return

        tool = self.mask_tool_var.get() if hasattr(self, "mask_tool_var") else "brush"

        if tool == "rect" and self._mask_rect_start:
            x0, y0 = self._mask_rect_start
            x1, y1 = img_x, img_y
            self._mask_apply_rect(page, layer, x0, y0, x1, y1)
            if self._mask_rect_id:
                try:
                    self.canvas.delete(self._mask_rect_id)
                except Exception:
                    pass
            self._mask_rect_id = None
            self._mask_rect_start = None
            self._mask_save_current()
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()
            return

        self._mask_last_point = None
        if getattr(self, "_mask_dirty", False):
            self._mask_save_current()
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()

    def _mask_brush_paint_segment(self, img_x0, img_y0, img_x1, img_y1):
        if not hasattr(self, "_mask_edit_mask") or self._mask_edit_mask is None:
            return
        if not self.pages:
            return
        page = self.pages[self.current_page_index]
        layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
        if layer is None:
            return

        try:
            brush_size = int(self.mask_brush_size_var.get())
        except Exception:
            brush_size = 40
        brush_size = max(1, brush_size)

        paint_value = self._mask_paint_value()

        layer_x = int(layer.get("x", 0))
        layer_y = int(layer.get("y", 0))
        w, h = self._mask_edit_mask.size

        lx0 = int(round(img_x0)) - layer_x
        ly0 = int(round(img_y0)) - layer_y
        lx1 = int(round(img_x1)) - layer_x
        ly1 = int(round(img_y1)) - layer_y

        # çº¿æ®µè£å‰ªï¼šåªè¦ç«¯ç‚¹éƒ½åœ¨å¤–é¢ä¹Ÿå¯èƒ½ç©¿è¿‡ï¼Œç®€å•èµ·è§ä¸åšå¤æ‚è£å‰ª
        if (lx0 < -brush_size and lx1 < -brush_size) or (ly0 < -brush_size and ly1 < -brush_size):
            return
        if (lx0 > w + brush_size and lx1 > w + brush_size) or (ly0 > h + brush_size and ly1 > h + brush_size):
            return

        try:
            self._mask_edit_draw.line([(lx0, ly0), (lx1, ly1)], fill=paint_value, width=brush_size, joint="curve")
        except Exception:
            self._mask_edit_draw.line([(lx0, ly0), (lx1, ly1)], fill=paint_value, width=brush_size)
        self._mask_dirty = True

    def _mask_apply_rect(self, page, layer, img_x0, img_y0, img_x1, img_y1):
        if not hasattr(self, "_mask_edit_mask") or self._mask_edit_mask is None:
            return

        layer_x = int(layer.get("x", 0))
        layer_y = int(layer.get("y", 0))
        w, h = self._mask_edit_mask.size

        x0 = int(min(img_x0, img_x1)) - layer_x
        y0 = int(min(img_y0, img_y1)) - layer_y
        x1 = int(max(img_x0, img_x1)) - layer_x
        y1 = int(max(img_y0, img_y1)) - layer_y

        x0 = max(0, min(w, x0))
        y0 = max(0, min(h, y0))
        x1 = max(0, min(w, x1))
        y1 = max(0, min(h, y1))
        if x1 <= x0 or y1 <= y0:
            return

        rect_mode = "å±€éƒ¨ç»˜åˆ¶"
        try:
            rect_mode = self.mask_rect_mode_var.get()
        except Exception:
            rect_mode = "å±€éƒ¨ç»˜åˆ¶"

        if rect_mode == "åªæ˜¾ç¤ºé€‰åŒº(é‡å»º)":
            self._mask_edit_mask.paste(0, (0, 0, w, h))
            self._mask_edit_mask.paste(255, (x0, y0, x1, y1))
        elif rect_mode == "åªéšè—é€‰åŒº(é‡å»º)":
            self._mask_edit_mask.paste(255, (0, 0, w, h))
            self._mask_edit_mask.paste(0, (x0, y0, x1, y1))
        else:
            paint_value = self._mask_paint_value()
            self._mask_edit_mask.paste(paint_value, (x0, y0, x1, y1))

        self._mask_edit_draw = ImageDraw.Draw(self._mask_edit_mask)
        self._mask_dirty = True

    def _mask_schedule_refresh(self):
        """è’™ç‰ˆç¼–è¾‘çš„åˆ·æ–°èŠ‚æµï¼Œé¿å…æ¯ä¸ªé¼ æ ‡äº‹ä»¶éƒ½è§¦å‘ç£ç›˜/æ¸²æŸ“å¯¼è‡´å¡æ­»ã€‚"""
        try:
            if self._mask_refresh_after is not None:
                self.root.after_cancel(self._mask_refresh_after)
        except Exception:
            pass
        try:
            self._mask_refresh_after = self.root.after(30, self.refresh_canvas)
        except Exception:
            self._mask_refresh_after = None

    def _draw_mask_edit_overlay(self, base_size, offset_x, offset_y, display_w, display_h):
        # å›¾å±‚è’™ç‰ˆç¼–è¾‘åŠŸèƒ½å·²ç§»é™¤ï¼ˆæ­¤å‰ä¼šå¯¼è‡´å¡æ­»/é—ªé€€ï¼‰ã€‚
        return
        if not getattr(self, "layer_mask_edit_mode", False):
            return
        if not self.pages:
            return

        page = self.pages[self.current_page_index]
        layer = self._mask_get_layer_by_id(page, getattr(self, "_mask_edit_layer_id", None))
        if layer is None:
            return

        # ä¼˜å…ˆä½¿ç”¨å†…å­˜ä¸­çš„è’™ç‰ˆï¼ˆç¼–è¾‘ä¸­ï¼‰ï¼Œé¿å…æ¯æ¬¡åˆ·æ–°éƒ½è¯»å†™ç£ç›˜å¯¼è‡´å¡é¡¿
        mask = None
        try:
            if (
                getattr(self, "layer_mask_edit_mode", False)
                and getattr(self, "_mask_edit_layer_id", None) == layer.get("id")
                and getattr(self, "_mask_edit_mask", None) is not None
            ):
                mask = self._mask_edit_mask
        except Exception:
            mask = None

        if mask is None:
            mask_path = layer.get("mask_path")
            if not mask_path or not os.path.exists(mask_path):
                return
            try:
                mask = Image.open(mask_path).convert("L")
            except Exception:
                return

        # ç»Ÿä¸€å°ºå¯¸åˆ°å›¾å±‚å°ºå¯¸ï¼ˆç”¨ nearestï¼Œé¿å…ç°è¾¹/æ€§èƒ½é—®é¢˜ï¼‰
        try:
            overlay_path = layer.get("path")
            if not overlay_path or not os.path.exists(overlay_path):
                return
            overlay_size = Image.open(overlay_path).size
            if mask.size != overlay_size:
                mask = mask.resize(overlay_size, Image.Resampling.NEAREST)
        except Exception:
            return

        if layer.get("mask_invert"):
            try:
                mask = ImageOps.invert(mask)
            except Exception:
                pass

        # éšè—åŒºåŸŸæ˜¾ç¤ºä¸ºçº¢è‰²åŠé€æ˜ï¼šalpha = (255 - mask) * (opacity/100)
        try:
            opacity = int(getattr(self, "mask_overlay_opacity_var").get())
        except Exception:
            opacity = 55
        opacity = max(0, min(opacity, 90)) / 100.0
        hidden = ImageOps.invert(mask)
        hidden = hidden.point(lambda v, a=opacity: int(v * a))
        red = Image.new("RGBA", overlay_img.size, (255, 0, 0, 0))
        red.putalpha(hidden)

        base_w, base_h = base_size
        full = Image.new("RGBA", (base_w, base_h), (255, 0, 0, 0))
        x = int(layer.get("x", 0))
        y = int(layer.get("y", 0))
        full.paste(red, (x, y), red)

        full_disp = full.resize((display_w, display_h), Image.Resampling.NEAREST)
        self._mask_overlay_tk = ImageTk.PhotoImage(full_disp)
        self.canvas.create_image(offset_x, offset_y, anchor=tk.NW, image=self._mask_overlay_tk, tags="mask_overlay")

    def _get_font_path(self, font_name):
        """è·å–å­—ä½“è·¯å¾„"""
        font_map = {
            "å¾®è½¯é›…é»‘": "C:/Windows/Fonts/msyh.ttc",
            "å®‹ä½“": "C:/Windows/Fonts/simsun.ttc",
            "é»‘ä½“": "C:/Windows/Fonts/simhei.ttf",
            "æ¥·ä½“": "C:/Windows/Fonts/simkai.ttf",
            "ä»¿å®‹": "C:/Windows/Fonts/simfang.ttf",
            "Arial": "C:/Windows/Fonts/arial.ttf"
        }
        path = font_map.get(font_name)
        if path and os.path.exists(path):
            return path
        return font_map.get("å¾®è½¯é›…é»‘")

    def draw_box(self, idx, box, offset_x, offset_y):
        """ç»˜åˆ¶æ–‡æœ¬æ¡†"""
        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        is_primary = (idx == self.selected_box_index)
        is_multi = (idx in self.selected_boxes)

        if is_primary:
            color, width = "#1976D2", 3
        elif is_multi:
            color, width = "#4CAF50", 2
        else:
            color, width = "#f44336", 2

        self.canvas.create_rectangle(x1, y1, x2, y2, outline=color, width=width, tags=f"box_{idx}")

        # åºå·
        self.canvas.create_oval(x1 + 5, y1 + 5, x1 + 22, y1 + 22, fill="#FF9800", outline="")
        self.canvas.create_text(x1 + 13, y1 + 13, text=str(idx + 1), fill="white", font=("Arial", 8, "bold"))

        # æ–‡æœ¬é¢„è§ˆ
        if box.text and y2 - y1 > 30:
            preview = box.text[:15] + "..." if len(box.text) > 15 else box.text
            self.canvas.create_text(x1 + 5, y2 - 12, text=preview, fill="#333333",
                                   anchor=tk.NW, font=("å¾®è½¯é›…é»‘", 8))

        # é€‰ä¸­æ‰‹æŸ„
        if is_primary:
            handle_size = 8
            handles = [(x1, y1), (x2, y1), (x1, y2), (x2, y2),
                      ((x1+x2)//2, y1), ((x1+x2)//2, y2), (x1, (y1+y2)//2), (x2, (y1+y2)//2)]
            for hx, hy in handles:
                self.canvas.create_rectangle(hx - handle_size//2, hy - handle_size//2,
                                            hx + handle_size//2, hy + handle_size//2,
                                            fill="#1976D2", outline="white")

    def _draw_ppt_edit_box(self, idx, box, offset_x, offset_y):
        """PPTé¢„è§ˆæ¨¡å¼ä¸‹çš„ç¼–è¾‘æ¡†"""
        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        is_primary = (idx == self.selected_box_index)
        is_multi = (idx in self.selected_boxes)

        if is_primary:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#1976D2", width=2, dash=(4, 4))
            handle_size = 8
            handles = [(x1, y1), (x2, y1), (x1, y2), (x2, y2),
                      ((x1+x2)//2, y1), ((x1+x2)//2, y2), (x1, (y1+y2)//2), (x2, (y1+y2)//2)]
            for hx, hy in handles:
                self.canvas.create_rectangle(hx - handle_size//2, hy - handle_size//2,
                                            hx + handle_size//2, hy + handle_size//2,
                                            fill="#1976D2", outline="white")
        elif is_multi:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#4CAF50", width=2, dash=(4, 4))
        else:
            self.canvas.create_rectangle(x1, y1, x2, y2, outline="#999999", width=1, dash=(2, 4))

    # ==================== é¼ æ ‡äº‹ä»¶ ====================

    def on_canvas_press(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        # è½¬æ¢ä¸ºå›¾ç‰‡åæ ‡
        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        # å›¾å±‚æ‹–åŠ¨ï¼ˆä»…æ‹–åŠ¨â€œå½“å‰é€‰ä¸­å›¾å±‚â€ï¼Œé¿å…ä¸æ–‡æœ¬æ¡†æ“ä½œå†²çªï¼‰
        page, layers, layer = self._get_selected_layer()
        if (
            layer is not None
            and self.current_preview_mode in ("edit", "ppt")
            and not self.inpaint_mode
            and not self.ai_replace_mode
            and not layer.get("locked")
            and layer.get("visible", True)
        ):
            bbox = self._layer_bbox(layer)
            if bbox is not None:
                x0, y0, x1, y1 = bbox
                if x0 <= img_x <= x1 and y0 <= img_y <= y1:
                    self.is_layer_dragging = True
                    self._layer_drag_start_canvas = (canvas_x, canvas_y)
                    self._layer_drag_origin_xy = (int(layer.get("x", 0)), int(layer.get("y", 0)))
                    self.save_state("layers")
                    return

        # æ¶‚æŠ¹æ¨¡å¼å¤„ç†
        if self.inpaint_mode:
            self.handle_inpaint_press(img_x, img_y)
            return

        # AIæ›¿æ¢æ¨¡å¼å¤„ç†
        if self.ai_replace_mode:
            self.handle_ai_replace_press(img_x, img_y)
            return

        if self.selected_box_index >= 0:
            handle = self.check_resize_handle(canvas_x, canvas_y)
            if handle:
                self.is_resizing = True
                self.resize_handle = handle
                self.drag_start_x = canvas_x
                self.drag_start_y = canvas_y
                return

        clicked_idx = self.find_box_at(img_x, img_y)

        if clicked_idx >= 0:
            self.select_box(clicked_idx)
            self.is_dragging = True
            self.drag_start_x = canvas_x
            self.drag_start_y = canvas_y
        elif self.draw_mode:
            # ç”»æ¡†æ¨¡å¼
            self.is_drawing = True
            self.draw_start_x = img_x
            self.draw_start_y = img_y
        else:
            # é€‰æ‹©æ¨¡å¼ï¼šå¼€å§‹æ¡†é€‰
            self.is_selecting = True
            self.select_start_x = canvas_x
            self.select_start_y = canvas_y

    def on_canvas_ctrl_click(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        clicked_idx = self.find_box_at(img_x, img_y)

        if clicked_idx >= 0:
            if clicked_idx in self.selected_boxes:
                self.selected_boxes.remove(clicked_idx)
            else:
                self.selected_boxes.append(clicked_idx)

            if self.selected_boxes:
                self.selected_box_index = self.selected_boxes[-1]
            else:
                self.selected_box_index = -1

            self.refresh_canvas()
            self.update_property_panel()
            self.update_status(f"å·²é€‰ä¸­ {len(self.selected_boxes)} ä¸ªæ¡†")

    def on_canvas_drag(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)
        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        if getattr(self, "is_layer_dragging", False):
            page, layers, layer = self._get_selected_layer()
            if layer is None or layer.get("locked"):
                self.is_layer_dragging = False
                return
            start_cx, start_cy = getattr(self, "_layer_drag_start_canvas", (canvas_x, canvas_y))
            ox, oy = getattr(self, "_layer_drag_origin_xy", (int(layer.get("x", 0)), int(layer.get("y", 0))))
            dx = (canvas_x - start_cx) / self.scale
            dy = (canvas_y - start_cy) / self.scale
            layer["x"] = int(round(ox + dx))
            layer["y"] = int(round(oy + dy))
            try:
                self._layer_transform_syncing = True
                if hasattr(self, "layer_x_var"):
                    self.layer_x_var.set(int(layer["x"]))
                if hasattr(self, "layer_y_var"):
                    self.layer_y_var.set(int(layer["y"]))
            finally:
                self._layer_transform_syncing = False
            self.refresh_canvas()
            return

        # æ¶‚æŠ¹æ¨¡å¼å¤„ç†
        if self.inpaint_mode:
            img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
            img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale
            self.handle_inpaint_drag(img_x, img_y)
            return

        # AIæ›¿æ¢æ¨¡å¼å¤„ç†
        if self.ai_replace_mode:
            self.handle_ai_replace_drag(canvas_x, canvas_y)
            return

        if self.is_resizing and self.selected_box_index >= 0:
            self.resize_selected_box(canvas_x, canvas_y)
        elif self.is_dragging and self.selected_box_index >= 0:
            self.drag_selected_box(canvas_x, canvas_y)
        elif self.is_drawing:
            self.draw_temp_rect(canvas_x, canvas_y)
        elif self.is_selecting:
            self.draw_selection_rect(canvas_x, canvas_y)

    def on_canvas_release(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)
        # è’™ç‰ˆåŠŸèƒ½å·²ç§»é™¤

        if getattr(self, "is_layer_dragging", False):
            self.is_layer_dragging = False
            self.update_layer_listbox()
            self.refresh_canvas()
            self.mark_unsaved()
            return

        # æ¶‚æŠ¹æ¨¡å¼å¤„ç†
        if self.inpaint_mode:
            img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
            img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale
            self.handle_inpaint_release(img_x, img_y)
            return

        # AIæ›¿æ¢æ¨¡å¼å¤„ç†
        if self.ai_replace_mode:
            self.handle_ai_replace_release(canvas_x, canvas_y)
            return

        if self.is_drawing:
            self.finish_drawing(canvas_x, canvas_y)
        elif self.is_selecting:
            self.finish_selection(canvas_x, canvas_y)

        self.is_drawing = False
        self.is_dragging = False
        self.is_resizing = False
        self.is_selecting = False
        self.resize_handle = None
        self.canvas.delete("temp_rect")
        self.canvas.delete("selection_rect")

    def on_canvas_double_click(self, event):
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        clicked_idx = self.find_box_at(img_x, img_y)
        if clicked_idx >= 0:
            self.select_box(clicked_idx)
            self.show_inline_text_editor(clicked_idx)

    def on_canvas_right_click(self, event):
        """å³é”®èœå•"""
        canvas_x = self.canvas.canvasx(event.x)
        canvas_y = self.canvas.canvasy(event.y)

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        # æŸ¥æ‰¾ç‚¹å‡»çš„æ–‡æœ¬æ¡†
        clicked_idx = self.find_box_at(img_x, img_y)

        # åˆ›å»ºå³é”®èœå•
        menu = tk.Menu(self.root, tearoff=0, font=(FONT_FAMILY, 9))

        if clicked_idx >= 0:
            # ç‚¹å‡»åœ¨æ–‡æœ¬æ¡†ä¸Š
            self.select_box(clicked_idx)

            menu.add_command(label="ğŸ” OCRè¯†åˆ«æ­¤æ¡†", command=self.ocr_single_box,
                           font=(FONT_FAMILY, 9, "bold"))
            menu.add_separator()
            menu.add_command(label="âœï¸ ç¼–è¾‘æ–‡å­—", command=lambda: self.show_inline_text_editor(clicked_idx))
            menu.add_separator()
            menu.add_command(label="ğŸ“‹ å¤åˆ¶ (Ctrl+C)", command=self.copy_boxes)
            menu.add_command(label="ğŸ“„ ç²˜è´´ (Ctrl+V)", command=self.paste_boxes)
            menu.add_separator()
            menu.add_command(label="ğŸ—‘ï¸ åˆ é™¤ (Del)", command=self.delete_selected_box,
                           foreground=COLOR_RED)
        else:
            # ç‚¹å‡»åœ¨ç©ºç™½å¤„
            if self.clipboard_boxes:
                menu.add_command(label="ğŸ“„ ç²˜è´´ (Ctrl+V)", command=self.paste_boxes)
                menu.add_separator()

            menu.add_command(label="ğŸ“ å¼€å§‹ç”»æ¡†", command=self.toggle_draw_mode_btn)

            if self.text_boxes:
                menu.add_separator()
                menu.add_command(label="ğŸ” OCRè¯†åˆ«å…¨éƒ¨", command=self.ocr_all_boxes)

        # æ˜¾ç¤ºèœå•
        try:
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()

    def find_box_at(self, x, y):
        for idx in range(len(self.text_boxes) - 1, -1, -1):
            box = self.text_boxes[idx]
            if box.x <= x <= box.x + box.width and box.y <= y <= box.y + box.height:
                return idx
        return -1

    def check_resize_handle(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return None

        box = self.text_boxes[self.selected_box_index]
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = int(box.x * self.scale) + offset_x
        y1 = int(box.y * self.scale) + offset_y
        x2 = int((box.x + box.width) * self.scale) + offset_x
        y2 = int((box.y + box.height) * self.scale) + offset_y

        handle_size = 10
        handles = {
            "nw": (x1, y1), "ne": (x2, y1), "sw": (x1, y2), "se": (x2, y2),
            "n": ((x1+x2)//2, y1), "s": ((x1+x2)//2, y2),
            "w": (x1, (y1+y2)//2), "e": (x2, (y1+y2)//2)
        }

        for handle_type, (hx, hy) in handles.items():
            if abs(canvas_x - hx) < handle_size and abs(canvas_y - hy) < handle_size:
                return handle_type
        return None

    def draw_temp_rect(self, canvas_x, canvas_y):
        self.canvas.delete("temp_rect")

        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = int(self.draw_start_x * self.scale) + offset_x
        y1 = int(self.draw_start_y * self.scale) + offset_y
        x2 = int(canvas_x)
        y2 = int(canvas_y)

        self.canvas.create_rectangle(x1, y1, x2, y2, outline="#1976D2", width=2,
                                    dash=(5, 5), tags="temp_rect")

    def finish_drawing(self, canvas_x, canvas_y):
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        x1 = self.draw_start_x
        y1 = self.draw_start_y
        x2 = (canvas_x - offset_x) / self.scale
        y2 = (canvas_y - offset_y) / self.scale

        if x1 > x2: x1, x2 = x2, x1
        if y1 > y2: y1, y2 = y2, y1

        width = x2 - x1
        height = y2 - y1

        if width < 10 or height < 10:
            return

        self.save_state()

        new_box = TextBox(int(x1), int(y1), int(width), int(height))
        self.text_boxes.append(new_box)
        self.select_box(len(self.text_boxes) - 1)
        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()

    def draw_selection_rect(self, canvas_x, canvas_y):
        """ç»˜åˆ¶æ¡†é€‰åŒºåŸŸ"""
        self.canvas.delete("selection_rect")

        x1 = int(self.select_start_x)
        y1 = int(self.select_start_y)
        x2 = int(canvas_x)
        y2 = int(canvas_y)

        # ç»˜åˆ¶åŠé€æ˜çš„è“è‰²é€‰åŒºçŸ©å½¢
        self.canvas.create_rectangle(x1, y1, x2, y2,
                                     outline="#2196F3", width=2,
                                     dash=(3, 3), tags="selection_rect")

    def finish_selection(self, canvas_x, canvas_y):
        """å®Œæˆæ¡†é€‰ï¼Œé€‰ä¸­é€‰åŒºå†…çš„æ‰€æœ‰æ¡†"""
        offset_x = getattr(self, 'canvas_offset_x', 0)
        offset_y = getattr(self, 'canvas_offset_y', 0)

        # è®¡ç®—é€‰åŒºçš„å›¾ç‰‡åæ ‡
        x1 = (self.select_start_x - offset_x) / self.scale
        y1 = (self.select_start_y - offset_y) / self.scale
        x2 = (canvas_x - offset_x) / self.scale
        y2 = (canvas_y - offset_y) / self.scale

        # ç¡®ä¿x1 < x2, y1 < y2
        if x1 > x2: x1, x2 = x2, x1
        if y1 > y2: y1, y2 = y2, y1

        # é€‰åŒºå¤ªå°åˆ™å¿½ç•¥
        if abs(x2 - x1) < 5 or abs(y2 - y1) < 5:
            return

        # æŸ¥æ‰¾é€‰åŒºå†…çš„æ‰€æœ‰æ–‡æœ¬æ¡†
        selected_indices = []
        for idx, box in enumerate(self.text_boxes):
            # æ£€æŸ¥æ–‡æœ¬æ¡†æ˜¯å¦ä¸é€‰åŒºç›¸äº¤æˆ–åŒ…å«åœ¨é€‰åŒºå†…
            box_left = box.x
            box_right = box.x + box.width
            box_top = box.y
            box_bottom = box.y + box.height

            # åˆ¤æ–­ç›¸äº¤ï¼šé€‰åŒºçš„ä»»æ„éƒ¨åˆ†ä¸æ¡†é‡å 
            if (box_left < x2 and box_right > x1 and
                box_top < y2 and box_bottom > y1):
                selected_indices.append(idx)

        # é€‰ä¸­æ‰¾åˆ°çš„æ¡†
        if selected_indices:
            self.selected_boxes = selected_indices
            self.selected_box_index = selected_indices[0] if selected_indices else -1

            # æ›´æ–°ç•Œé¢
            self.refresh_canvas()
            self.update_property_panel()

            # æ›´æ–°åˆ—è¡¨æ¡†é€‰æ‹©
            self.box_listbox.selection_clear(0, tk.END)
            for idx in self.selected_boxes:
                self.box_listbox.selection_set(idx)

            self.update_status(f"æ¡†é€‰é€‰ä¸­ {len(selected_indices)} ä¸ªæ–‡æœ¬æ¡† âœ“")
        else:
            # æ²¡æœ‰é€‰ä¸­ä»»ä½•æ¡†ï¼Œæ¸…ç©ºé€‰æ‹©
            self.selected_boxes = []
            self.selected_box_index = -1
            self.refresh_canvas()
            self.update_status("æ¡†é€‰åŒºåŸŸå†…æ²¡æœ‰æ–‡æœ¬æ¡†")

    def resize_selected_box(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        dx = (canvas_x - self.drag_start_x) / self.scale
        dy = (canvas_y - self.drag_start_y) / self.scale

        if "w" in self.resize_handle:
            new_x = box.x + dx
            new_w = box.width - dx
            if new_w > 10:
                box.x = int(new_x)
                box.width = int(new_w)
        if "e" in self.resize_handle:
            new_w = box.width + dx
            if new_w > 10:
                box.width = int(new_w)
        if "n" in self.resize_handle:
            new_y = box.y + dy
            new_h = box.height - dy
            if new_h > 10:
                box.y = int(new_y)
                box.height = int(new_h)
        if "s" in self.resize_handle:
            new_h = box.height + dy
            if new_h > 10:
                box.height = int(new_h)

        self.drag_start_x = canvas_x
        self.drag_start_y = canvas_y
        self.refresh_canvas()
        self.update_property_panel()

    def drag_selected_box(self, canvas_x, canvas_y):
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        dx = (canvas_x - self.drag_start_x) / self.scale
        dy = (canvas_y - self.drag_start_y) / self.scale

        box.x = int(box.x + dx)
        box.y = int(box.y + dy)

        self.drag_start_x = canvas_x
        self.drag_start_y = canvas_y
        self.refresh_canvas()
        self.update_property_panel()

    # ==================== é€‰æ‹©ä¸å±æ€§ ====================

    def select_box(self, idx):
        self.selected_box_index = idx
        self.selected_boxes = [idx] if idx >= 0 else []
        self.refresh_canvas()
        self.update_property_panel()

        self.box_listbox.selection_clear(0, tk.END)
        if idx >= 0:
            self.box_listbox.selection_set(idx)
            self.box_listbox.see(idx)

    def update_listbox(self):
        self.box_listbox.delete(0, tk.END)
        for idx, box in enumerate(self.text_boxes):
            text_preview = box.text[:15] + "..." if len(box.text) > 15 else box.text
            if not text_preview:
                text_preview = "(ç©º)"
            self.box_listbox.insert(tk.END, f"{idx+1}. {text_preview}")

    def on_listbox_select(self, event):
        selection = self.box_listbox.curselection()
        if selection:
            self.select_box(selection[0])

    def update_property_panel(self):
        if self.selected_box_index < 0 or self.selected_box_index >= len(self.text_boxes):
            return

        box = self.text_boxes[self.selected_box_index]

        self.text_entry.delete("1.0", tk.END)
        self.text_entry.insert("1.0", box.text)

        self.x_entry.delete(0, tk.END)
        self.x_entry.insert(0, str(box.x))
        self.y_entry.delete(0, tk.END)
        self.y_entry.insert(0, str(box.y))
        self.w_entry.delete(0, tk.END)
        self.w_entry.insert(0, str(box.width))
        self.h_entry.delete(0, tk.END)
        self.h_entry.insert(0, str(box.height))

        self.fontsize_var.set(str(box.font_size))
        self.fontname_var.set(box.font_name)
        self.bold_var.set(box.bold)
        self.italic_var.set(box.italic)
        self.align_var.set(box.align)
        self.color_btn.config(bg=box.font_color)

        self.update_style_buttons()
        self.update_align_buttons()

    def update_style_buttons(self):
        if self.bold_var.get():
            self.bold_btn.config(bg="#1976D2", fg="white")
        else:
            self.bold_btn.config(bg="#e0e0e0", fg="black")

        if self.italic_var.get():
            self.italic_btn.config(bg="#1976D2", fg="white")
        else:
            self.italic_btn.config(bg="#e0e0e0", fg="black")

    def on_text_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        box.text = self.text_entry.get("1.0", tk.END).strip()
        self.update_listbox()
        self.refresh_canvas()

    def on_position_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        try:
            box.x = int(self.x_entry.get())
            box.y = int(self.y_entry.get())
            box.width = int(self.w_entry.get())
            box.height = int(self.h_entry.get())
            self.refresh_canvas()
        except ValueError:
            pass

    def on_font_change(self, event=None):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        try:
            box.font_size = int(self.fontsize_var.get())
        except:
            pass
        box.font_name = self.fontname_var.get()
        self.refresh_canvas()

    def set_align(self, align):
        """è®¾ç½®å¯¹é½æ–¹å¼"""
        self.align_var.set(align)
        self.update_align_buttons()
        self.on_style_change()

    def update_align_buttons(self):
        """æ›´æ–°å¯¹é½æŒ‰é’®çŠ¶æ€"""
        align = self.align_var.get()
        # å·¦å¯¹é½
        if align == "left":
            self.align_left_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_left_btn.config(bg="#e0e0e0", fg="#333")
        # å±…ä¸­
        if align == "center":
            self.align_center_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_center_btn.config(bg="#e0e0e0", fg="#333")
        # å³å¯¹é½
        if align == "right":
            self.align_right_btn.config(bg="#1976D2", fg="white")
        else:
            self.align_right_btn.config(bg="#e0e0e0", fg="#333")

    def on_style_change(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        box.bold = self.bold_var.get()
        box.italic = self.italic_var.get()
        box.align = self.align_var.get()
        self.refresh_canvas()

    def toggle_bold(self):
        self.bold_var.set(not self.bold_var.get())
        self.update_style_buttons()
        self.on_style_change()

    def toggle_italic(self):
        self.italic_var.set(not self.italic_var.get())
        self.update_style_buttons()
        self.on_style_change()

    def choose_color(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        color = colorchooser.askcolor(color=box.font_color, title="é€‰æ‹©æ–‡å­—é¢œè‰²")
        if color[1]:
            box.font_color = color[1]
            self.color_btn.config(bg=color[1])
            self.refresh_canvas()

    # ==================== å…¶ä»–æ“ä½œ ====================

    def toggle_draw_mode(self):
        self.draw_mode = self.draw_mode_var.get()
        if self.draw_mode:
            self.canvas.config(cursor="crosshair")
        else:
            self.canvas.config(cursor="")

    def switch_preview_mode(self):
        self.current_preview_mode = self.preview_mode_var.get()
        self.refresh_canvas()

    def refresh_ppt_preview(self):
        self.preview_mode_var.set("ppt")
        self.current_preview_mode = "ppt"
        self.refresh_canvas()
        self.update_status("PPTé¢„è§ˆå·²åˆ·æ–° âœ“")

    def show_inline_text_editor(self, box_idx):
        """å†…è”æ–‡å­—ç¼–è¾‘å™¨"""
        if box_idx < 0 or box_idx >= len(self.text_boxes):
            return

        box = self.text_boxes[box_idx]

        edit_window = tk.Toplevel(self.root)
        edit_window.title(f"ç¼–è¾‘æ–‡æœ¬æ¡† {box_idx + 1}")
        edit_window.geometry("420x300")
        edit_window.configure(bg="#ffffff")
        edit_window.transient(self.root)
        edit_window.grab_set()

        mouse_x = self.root.winfo_pointerx()
        mouse_y = self.root.winfo_pointery()
        edit_window.geometry(f"+{mouse_x - 210}+{mouse_y - 150}")

        # æ–‡å­—è¾“å…¥
        tk.Label(edit_window, text="æ–‡å­—å†…å®¹", bg="#ffffff",
                fg="#333333", font=("å¾®è½¯é›…é»‘", 9, "bold")).pack(anchor="w", padx=15, pady=(15, 5))

        text_input = tk.Text(edit_window, height=4, bg="#f5f5f5",
                            font=("å¾®è½¯é›…é»‘", 11), relief=tk.GROOVE, bd=1, wrap=tk.WORD)
        text_input.pack(fill=tk.X, padx=15, pady=5)
        text_input.insert("1.0", box.text)
        text_input.focus_set()
        text_input.tag_add("sel", "1.0", "end")

        # å¿«æ·è®¾ç½®
        quick_frame = tk.Frame(edit_window, bg="#ffffff")
        quick_frame.pack(fill=tk.X, padx=15, pady=10)

        tk.Label(quick_frame, text="å­—å·:", bg="#ffffff", font=("å¾®è½¯é›…é»‘", 9)).pack(side=tk.LEFT)
        font_size_var = tk.StringVar(value=str(box.font_size))
        ttk.Combobox(quick_frame, textvariable=font_size_var, width=5,
                    values=["8", "10", "12", "14", "16", "18", "20", "24", "28", "32", "36", "48", "60", "72", "80", "100", "120", "150", "200"]).pack(side=tk.LEFT, padx=5)

        tk.Label(quick_frame, text="å¯¹é½:", bg="#ffffff", font=("å¾®è½¯é›…é»‘", 9)).pack(side=tk.LEFT, padx=(15, 0))
        align_var = tk.StringVar(value=box.align)
        for val, txt in [("left", "å·¦"), ("center", "ä¸­"), ("right", "å³")]:
            tk.Radiobutton(quick_frame, text=txt, variable=align_var, value=val,
                          bg="#ffffff", font=("å¾®è½¯é›…é»‘", 9)).pack(side=tk.LEFT)

        # æ ·å¼
        style_frame = tk.Frame(edit_window, bg="#ffffff")
        style_frame.pack(fill=tk.X, padx=15, pady=5)

        bold_var = tk.BooleanVar(value=box.bold)
        tk.Checkbutton(style_frame, text="åŠ ç²—", variable=bold_var,
                      bg="#ffffff", font=("å¾®è½¯é›…é»‘", 9)).pack(side=tk.LEFT)

        color_var = tk.StringVar(value=box.font_color)
        color_btn = tk.Button(style_frame, text="é¢œè‰²", bg=box.font_color, width=6,
                             command=lambda: self._pick_color_for_editor(color_btn, color_var))
        color_btn.pack(side=tk.LEFT, padx=10)

        def auto_calc():
            text = text_input.get("1.0", tk.END).strip()
            if text:
                font_size_var.set(
                    str(
                        fit_font_size_pt(
                            text,
                            box.width,
                            box.height,
                            editor=self,
                            font_name=getattr(box, "font_name", None),
                        )
                    )
                )

        tk.Button(style_frame, text="è‡ªåŠ¨å­—å·", command=auto_calc,
                 bg="#9C27B0", fg="white", font=("å¾®è½¯é›…é»‘", 9)).pack(side=tk.LEFT, padx=10)

        # æŒ‰é’®
        btn_frame = tk.Frame(edit_window, bg="#ffffff")
        btn_frame.pack(fill=tk.X, padx=15, pady=15)

        def save():
            box.text = text_input.get("1.0", tk.END).strip()
            try:
                box.font_size = int(font_size_var.get())
            except:
                pass
            box.align = align_var.get()
            box.bold = bold_var.get()
            box.font_color = color_var.get()
            edit_window.destroy()
            self.refresh_canvas()
            self.update_listbox()
            self.update_property_panel()

        tk.Button(btn_frame, text="ç¡®å®š", command=save,
                 bg="#4CAF50", fg="white", font=("å¾®è½¯é›…é»‘", 10),
                 width=10, cursor="hand2").pack(side=tk.LEFT, padx=5)
        tk.Button(btn_frame, text="å–æ¶ˆ", command=edit_window.destroy,
                 bg="#9E9E9E", fg="white", font=("å¾®è½¯é›…é»‘", 10),
                 width=10, cursor="hand2").pack(side=tk.LEFT, padx=5)

        edit_window.bind("<Control-Return>", lambda e: save())
        edit_window.bind("<Escape>", lambda e: edit_window.destroy())

    def _pick_color_for_editor(self, btn, color_var):
        color = colorchooser.askcolor(color=color_var.get(), title="é€‰æ‹©é¢œè‰²")
        if color[1]:
            color_var.set(color[1])
            btn.config(bg=color[1])

    # ==================== æ’¤é”€/é‡åš ====================

    # ==================== æ’¤é”€/é‡åšç³»ç»Ÿï¼ˆå¢å¼ºç‰ˆï¼‰====================

    def save_state(self, operation_type="textboxes", extra_data=None):
        return history_core.save_state(self, operation_type=operation_type, extra_data=extra_data)

    def undo(self):
        return history_core.undo(self)

    def redo(self):
        return history_core.redo(self)

    def _restore_state(self, state):
        return history_core.restore_state(self, state)

    # ==================== æ¡†æ“ä½œ ====================

    def delete_selected_box(self):
        indices: list[int] = []
        if self.selected_boxes:
            indices.extend(self.selected_boxes)
        if self.selected_box_index >= 0:
            indices.append(self.selected_box_index)

        indices = sorted({i for i in indices if 0 <= i < len(self.text_boxes)})
        if not indices:
            return

        self.save_state()
        for idx in sorted(indices, reverse=True):
            del self.text_boxes[idx]

        deleted_n = len(indices)
        self.selected_box_index = -1
        self.selected_boxes = []
        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()
        self.update_status(f"\u5df2\u5220\u9664 {deleted_n} \u4e2a\u6846 \u2713")

    def clear_all_boxes(self):
        if messagebox.askyesno("ç¡®è®¤", "ç¡®å®šæ¸…ç©ºæ‰€æœ‰æ–‡æœ¬æ¡†ï¼Ÿ"):
            self.save_state()
            self.text_boxes = []
            self.selected_box_index = -1
            self.selected_boxes = []
            self.refresh_canvas()
            self.update_listbox()

    def auto_font_size(self):
        if self.selected_box_index < 0:
            return
        box = self.text_boxes[self.selected_box_index]
        if not box.text:
            return
        box.font_size = fit_font_size_pt(
            box.text,
            box.width,
            box.height,
            editor=self,
            font_name=getattr(box, "font_name", None),
        )
        self.fontsize_var.set(str(box.font_size))
        self.refresh_canvas()

    def auto_font_size_all(self):
        for box in self.text_boxes:
            if not box.text:
                continue
            box.font_size = fit_font_size_pt(
                box.text,
                box.width,
                box.height,
                editor=self,
                font_name=getattr(box, "font_name", None),
            )
        self.update_property_panel()
        self.refresh_canvas()
        self.update_status("å·²ä¸ºå½“å‰é¡µæ‰€æœ‰æ¡†è®¡ç®—å­—å· âœ“")

    def align_boxes(self, align_type):
        if len(self.selected_boxes) < 2:
            self.update_status("è¯·Ctrl+ç‚¹å‡»é€‰ä¸­è‡³å°‘2ä¸ªæ¡†")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        if align_type == "left":
            min_x = min(b.x for b in boxes)
            for b in boxes: b.x = min_x
        elif align_type == "right":
            max_right = max(b.x + b.width for b in boxes)
            for b in boxes: b.x = max_right - b.width
        elif align_type == "center_h":
            avg = sum(b.x + b.width / 2 for b in boxes) / len(boxes)
            for b in boxes: b.x = int(avg - b.width / 2)
        elif align_type == "top":
            min_y = min(b.y for b in boxes)
            for b in boxes: b.y = min_y
        elif align_type == "bottom":
            max_bottom = max(b.y + b.height for b in boxes)
            for b in boxes: b.y = max_bottom - b.height
        elif align_type == "center_v":
            avg = sum(b.y + b.height / 2 for b in boxes) / len(boxes)
            for b in boxes: b.y = int(avg - b.height / 2)

        self.refresh_canvas()
        self.update_status(f"å·²å¯¹é½ {len(self.selected_boxes)} ä¸ªæ¡† âœ“")

    def batch_offset(self, dx_dir, dy_dir):
        """æ‰¹é‡ä½ç§»é€‰ä¸­çš„æ–‡æœ¬æ¡†

        Args:
            dx_dir: Xæ–¹å‘ï¼ˆ-1å·¦, 0æ— , 1å³ï¼‰
            dy_dir: Yæ–¹å‘ï¼ˆ-1ä¸Š, 0æ— , 1ä¸‹ï¼‰
        """
        # è‡³å°‘è¦æœ‰ä¸€ä¸ªé€‰ä¸­çš„æ¡†ï¼ˆåŒ…æ‹¬ä¸»é€‰ä¸­æ¡†ï¼‰
        boxes_to_move = []
        if self.selected_boxes:
            boxes_to_move = self.selected_boxes
        elif self.selected_box_index >= 0:
            boxes_to_move = [self.selected_box_index]

        if not boxes_to_move:
            self.update_status("è¯·å…ˆé€‰ä¸­è‡³å°‘ä¸€ä¸ªæ–‡æœ¬æ¡†")
            return

        # è·å–åƒç´ å€¼
        try:
            pixels = int(self.offset_px_var.get())
            if pixels <= 0:
                self.update_status("åƒç´ å€¼å¿…é¡»å¤§äº0")
                return
        except ValueError:
            self.update_status("è¯·è¾“å…¥æœ‰æ•ˆçš„åƒç´ æ•°å€¼")
            return

        # ä¿å­˜çŠ¶æ€ç”¨äºæ’¤é”€
        self.save_state()

        # è®¡ç®—å®é™…åç§»é‡
        dx = dx_dir * pixels
        dy = dy_dir * pixels

        # ç§»åŠ¨æ‰€æœ‰é€‰ä¸­çš„æ¡†
        for idx in boxes_to_move:
            if 0 <= idx < len(self.text_boxes):
                box = self.text_boxes[idx]
                box.x = max(0, box.x + dx)  # ä¸èƒ½ç§»å‡ºè¾¹ç•Œ
                box.y = max(0, box.y + dy)

        # æ›´æ–°ç•Œé¢
        self.refresh_canvas()
        self.update_property_panel()
        self.mark_unsaved()

        # æç¤ºä¿¡æ¯
        direction = ""
        if dx_dir == -1:
            direction = "å·¦"
        elif dx_dir == 1:
            direction = "å³"
        elif dy_dir == -1:
            direction = "ä¸Š"
        elif dy_dir == 1:
            direction = "ä¸‹"

        self.update_status(f"å·²å°† {len(boxes_to_move)} ä¸ªæ¡†å‘{direction}ç§»åŠ¨ {pixels} åƒç´  âœ“")

    def apply_style_to_selected(self):
        if len(self.selected_boxes) < 1:
            self.update_status("è¯·å…ˆCtrl+ç‚¹å‡»é€‰ä¸­æ¡†")
            return

        any_selected = (self.apply_fontsize_var.get() or self.apply_fontname_var.get() or
                       self.apply_color_var.get() or self.apply_bold_var.get() or
                       self.apply_italic_var.get() or self.apply_align_var.get())

        if not any_selected:
            self.update_status("è¯·å…ˆå‹¾é€‰è¦åº”ç”¨çš„å±æ€§")
            return

        self.save_state()

        try:
            font_size = int(self.fontsize_var.get())
        except:
            font_size = 16

        for idx in self.selected_boxes:
            if 0 <= idx < len(self.text_boxes):
                box = self.text_boxes[idx]
                if self.apply_fontsize_var.get(): box.font_size = font_size
                if self.apply_fontname_var.get(): box.font_name = self.fontname_var.get()
                if self.apply_bold_var.get(): box.bold = self.bold_var.get()
                if self.apply_italic_var.get(): box.italic = self.italic_var.get()
                if self.apply_align_var.get(): box.align = self.align_var.get()
                if self.apply_color_var.get(): box.font_color = self.color_btn.cget("bg")

        self.refresh_canvas()
        self.update_status(f"å·²åº”ç”¨æ ·å¼åˆ° {len(self.selected_boxes)} ä¸ªæ¡† âœ“")

    # ==================== OCR ====================

    def _prepare_image_for_ocr(self, img_path, edit_scale=1.0):
        """å‡†å¤‡OCRç”¨çš„å›¾ç‰‡ï¼Œå¦‚æœå›¾ç‰‡è¿‡å¤§åˆ™ç¼©æ”¾"""
        MAX_SIDE = 3000  # æœ€å¤§è¾¹é•¿é™åˆ¶

        img = Image.open(img_path)
        w, h = img.size

        # å…ˆåº”ç”¨ç¼–è¾‘ç¼©æ”¾
        if edit_scale < 1.0:
            w = int(w * edit_scale)
            h = int(h * edit_scale)
            img = img.resize((w, h), Image.Resampling.LANCZOS)

        # å¦‚æœè¿˜æ˜¯å¤ªå¤§ï¼Œå†ç¼©æ”¾
        if max(w, h) <= MAX_SIDE:
            # ä¿å­˜åˆ°ä¸´æ—¶æ–‡ä»¶
            temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
            temp_path = temp_file.name
            temp_file.close()
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            img.save(temp_path, quality=95)
            return temp_path, 1.0

        # è®¡ç®—é¢å¤–ç¼©æ”¾æ¯”ä¾‹
        extra_scale = MAX_SIDE / max(w, h)
        new_w = int(w * extra_scale)
        new_h = int(h * extra_scale)

        # ç¼©æ”¾å›¾ç‰‡
        resized_img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)

        # ä¿å­˜åˆ°ä¸´æ—¶æ–‡ä»¶
        temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
        temp_path = temp_file.name
        temp_file.close()

        if resized_img.mode == 'RGBA':
            resized_img = resized_img.convert('RGB')
        resized_img.save(temp_path, quality=95)

        return temp_path, extra_scale

    def auto_detect_text_regions(self):
        if not self.original_image:
            self.update_status("è¯·å…ˆåŠ è½½å›¾ç‰‡")
            return
        if not self.ocr:
            self.update_status("OCRæ¨¡å‹æœªåŠ è½½")
            return

        if self.text_boxes:
            result = messagebox.askyesnocancel("æç¤º", "æ˜¯å¦æ¸…ç©ºç°æœ‰æ¡†ï¼Ÿ\næ˜¯-æ¸…ç©º  å¦-è¿½åŠ   å–æ¶ˆ-å–æ¶ˆ")
            if result is None:
                return
            elif result:
                self.text_boxes = []

        self.update_status("æ­£åœ¨æ£€æµ‹...")

        def detect():
            try:
                # ç›´æ¥ä½¿ç”¨å½“å‰ç¼–è¾‘å›¾ç‰‡ï¼Œå®Œå…¨ä¸ç¼©æ”¾ï¼Œä¿è¯åæ ‡100%å‡†ç¡®
                # PIL Imageè½¬ä¸ºOpenCVæ ¼å¼
                img = np.array(self.original_image)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]

                # ä¿å­˜ä¸´æ—¶æ–‡ä»¶ç”¨äºOCRï¼ˆä¸ç¼©æ”¾ï¼ï¼‰
                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                cv2.imwrite(temp_path, img)

                result = self.ocr.predict(temp_path)

                # åˆ é™¤ä¸´æ—¶æ–‡ä»¶
                try:
                    os.remove(temp_path)
                except:
                    pass

                # æ–°ç‰ˆ PaddleOCR è¿”å› listï¼Œå–ç¬¬ä¸€ä¸ªç»“æœ
                if not result or len(result) == 0:
                    self.root.after(0, lambda: self.update_status("æœªæ£€æµ‹åˆ°æ–‡å­—"))
                    return

                ocr_result = result[0]
                dp = ocr_result.get("doc_preprocessor_res")
                dp_angle = None
                try:
                    dp_angle = dp.get("angle") if dp else None
                except Exception:
                    dp_angle = None
                dt_polys = ocr_result.get('dt_polys', [])
                rec_texts = ocr_result.get('rec_texts', [])

                if not dt_polys:
                    self.root.after(0, lambda: self.update_status("æœªæ£€æµ‹åˆ°æ–‡å­—"))
                    return

                new_boxes = []
                for i, poly in enumerate(dt_polys):
                    x_coords = [p[0] for p in poly]
                    y_coords = [p[1] for p in poly]

                    # å®Œå…¨ä½¿ç”¨OCRåŸå§‹åæ ‡ï¼Œä¸åšä»»ä½•è°ƒæ•´
                    x = int(min(x_coords))
                    y = int(min(y_coords))
                    w = int(max(x_coords) - min(x_coords))
                    h = int(max(y_coords) - min(y_coords))

                    if w < 10 or h < 10:
                        continue

                    box = TextBox(max(0, x), max(0, y), w, h)
                    if i < len(rec_texts):
                        box.text = rec_texts[i]
                    if box.text:
                        box.font_size = fit_font_size_pt(
                            box.text,
                            w,
                            h,
                            editor=self,
                            font_name=getattr(box, "font_name", None),
                        )
                    new_boxes.append(box)

                new_boxes.sort(key=lambda b: (b.y // 30, b.x))
                self.text_boxes.extend(new_boxes)

                self.root.after(0, self.refresh_canvas)
                self.root.after(0, self.update_listbox)
                if dp_angle not in (None, 0):
                    self.root.after(
                        0,
                        lambda a=dp_angle, n=len(new_boxes): self.update_status(
                            f"æ£€æµ‹åˆ° {n} ä¸ªæ–‡å­—åŒºåŸŸï¼ˆæç¤ºï¼šOCR æ–‡æ¡£é¢„å¤„ç†æ—‹è½¬äº†å›¾ç‰‡ {a}Â°ï¼Œå¦‚å æ¡†åç§»å¯åœ¨é…ç½®å…³é—­ç›¸å…³é¢„å¤„ç†ï¼‰"
                        ),
                    )
                else:
                    self.root.after(0, lambda n=len(new_boxes): self.update_status(f"æ£€æµ‹åˆ° {n} ä¸ªæ–‡å­—åŒºåŸŸ"))

            except Exception as e:
                err_text = str(e)
                self.root.after(0, lambda t=err_text: self.update_status(f"æ£€æµ‹å¤±è´¥: {t}"))

        threading.Thread(target=detect, daemon=True).start()

    def ocr_all_boxes(self):
        return ocr_core.ocr_all_boxes(self)

    def ocr_single_box(self):
        return ocr_core.ocr_single_box(self)

    # ==================== æ‰¹é‡æ“ä½œ ====================

    def auto_detect_all_pages(self):
        if not self.pages or not self.ocr:
            self.update_status("è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        self.save_current_page()

        def detect_all():
            total = len(self.pages)
            for i, page in enumerate(self.pages):
                self.root.after(0, lambda idx=i: self.update_status(f"æ£€æµ‹ç¬¬ {idx+1}/{total} é¡µ..."))

                try:
                    # ç›´æ¥ä½¿ç”¨è¯¥é¡µçš„ç¼–è¾‘å›¾ç‰‡ï¼Œå®Œå…¨ä¸ç¼©æ”¾
                    page_img = page["image"]
                    img = np.array(page_img)
                    img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                    # ä¿å­˜ä¸´æ—¶æ–‡ä»¶ï¼ˆä¸ç¼©æ”¾ï¼ï¼‰
                    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_path = temp_file.name
                    temp_file.close()
                    cv2.imwrite(temp_path, img)

                    result = self.ocr.predict(temp_path)

                    # åˆ é™¤ä¸´æ—¶æ–‡ä»¶
                    try:
                        os.remove(temp_path)
                    except:
                        pass

                    if not result or len(result) == 0:
                        continue

                    ocr_result = result[0]
                    dt_polys = ocr_result.get('dt_polys', [])
                    rec_texts = ocr_result.get('rec_texts', [])

                    if not dt_polys:
                        continue

                    new_boxes = []
                    for j, poly in enumerate(dt_polys):
                        x_coords = [p[0] for p in poly]
                        y_coords = [p[1] for p in poly]

                        # å®Œå…¨ä½¿ç”¨OCRåŸå§‹åæ ‡ï¼Œä¸åšä»»ä½•è°ƒæ•´
                        x = int(min(x_coords))
                        y = int(min(y_coords))
                        w = int(max(x_coords) - min(x_coords))
                        h = int(max(y_coords) - min(y_coords))

                        if w < 10 or h < 10:
                            continue

                        box_data = {
                            "x": max(0, x), "y": max(0, y), "width": w, "height": h,
                            "text": rec_texts[j] if j < len(rec_texts) else "",
                            "font_size": 16, "font_name": "å¾®è½¯é›…é»‘", "font_color": "#000000",
                            "bold": False, "italic": False, "align": "left"
                        }

                        if box_data["text"]:
                            box_data["font_size"] = fit_font_size_pt(
                                box_data["text"],
                                w,
                                h,
                                editor=self,
                                font_name=box_data.get("font_name"),
                            )

                        new_boxes.append(box_data)

                    new_boxes.sort(key=lambda b: (b["y"] // 30, b["x"]))
                    page["text_boxes"] = new_boxes

                except Exception as e:
                    print(f"ç¬¬ {i+1} é¡µæ£€æµ‹å¤±è´¥: {e}")

            self.root.after(0, self.load_current_page)
            self.root.after(0, lambda: self.update_status(f"å…¨éƒ¨æ£€æµ‹å®Œæˆï¼å…± {total} é¡µ âœ“"))

        threading.Thread(target=detect_all, daemon=True).start()

    def ocr_all_pages(self):
        if not self.pages or not self.ocr:
            return

        self.save_current_page()

        def ocr_all():
            total = len(self.pages)
            for i, page in enumerate(self.pages):
                self.root.after(0, lambda idx=i: self.update_status(f"è¯†åˆ«ç¬¬ {idx+1}/{total} é¡µ..."))

                boxes = page.get("text_boxes", [])
                if not boxes:
                    continue

                # ä½¿ç”¨è¯¥é¡µçš„ç¼–è¾‘å›¾ç‰‡
                page_img = page["image"]
                img = np.array(page_img)
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)

                img_h, img_w = img.shape[:2]

                for box_data in boxes:
                    if box_data.get("text"):
                        continue

                    x, y, w, h = box_data["x"], box_data["y"], box_data["width"], box_data["height"]
                    expand_h, expand_w = int(h * 0.3), int(w * 0.1)

                    crop_x = max(0, x - expand_w)
                    crop_y = max(0, y - expand_h)
                    crop_x2 = min(x + w + expand_w, img_w)
                    crop_y2 = min(y + h + expand_h, img_h)

                    cropped = img[crop_y:crop_y2, crop_x:crop_x2]

                    temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                    temp_path = temp_file.name
                    temp_file.close()
                    cv2.imwrite(temp_path, cropped)

                    try:
                        result = self.ocr.predict(temp_path)
                        os.remove(temp_path)

                        if result and len(result) > 0:
                            ocr_result = result[0]
                            rec_texts = ocr_result.get('rec_texts', [])
                            if rec_texts:
                                box_data["text"] = ''.join(rec_texts)
                                if box_data["text"]:
                                    box_data["font_size"] = fit_font_size_pt(
                                        box_data["text"],
                                        w,
                                        h,
                                        editor=self,
                                        font_name=box_data.get("font_name"),
                                    )
                    except:
                        try:
                            os.remove(temp_path)
                        except:
                            pass

            self.root.after(0, self.load_current_page)
            self.root.after(0, lambda: self.update_status(f"å…¨éƒ¨è¯†åˆ«å®Œæˆï¼å…± {total} é¡µ âœ“"))

        threading.Thread(target=ocr_all, daemon=True).start()

    def auto_font_size_all_pages(self):
        if not self.pages:
            return

        self.save_current_page()

        for page in self.pages:
            for box_data in page.get("text_boxes", []):
                if not box_data.get("text"):
                    continue
                h, w = box_data["height"], box_data["width"]
                box_data["font_size"] = fit_font_size_pt(
                    box_data["text"],
                    w,
                    h,
                    editor=self,
                    font_name=box_data.get("font_name"),
                )

        self.load_current_page()
        self.update_status(f"å…¨éƒ¨ {len(self.pages)} é¡µå­—å·å·²è°ƒæ•´ âœ“")

    # ==================== é¡¹ç›®ä¿å­˜/åŠ è½½ ====================

    def save_project(self):
        return project_feature.save_project(self)

    def load_project(self):
        return project_feature.load_project(self)

    # ==================== PPTç”Ÿæˆ ====================

    def generate_multi_page_ppt(self):
        return export_feature.generate_multi_page_ppt(self)

    # ==================== è®¾ç½®å¯¹è¯æ¡† ====================

    def show_settings_dialog(self):
        """æ˜¾ç¤ºè®¾ç½®å¯¹è¯æ¡†"""
        dialog = tk.Toplevel(self.root)
        dialog.title("è®¾ç½®")
        dialog.geometry("680x820")  # å¢åŠ é«˜åº¦ä»¥å®¹çº³IOPainté…ç½®
        dialog.configure(bg=COLOR_WHITE)
        dialog.transient(self.root)
        dialog.grab_set()

        # å±…ä¸­æ˜¾ç¤º
        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() - 680) // 2
        y = (dialog.winfo_screenheight() - 820) // 2
        dialog.geometry(f"+{x}+{y}")

        # æ ‡é¢˜
        title_frame = tk.Frame(dialog, bg=COLOR_THEME, height=40)
        title_frame.pack(fill=tk.X, side=tk.TOP)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text="  OCRæ¨¡å‹è®¾ç½®", bg=COLOR_THEME, fg="white",
                font=(FONT_FAMILY, 12, "bold")).pack(side=tk.LEFT, pady=8)

        # æŒ‰é’®åŒº - å›ºå®šåœ¨åº•éƒ¨
        btn_frame = tk.Frame(dialog, bg=COLOR_WHITE, pady=15)
        btn_frame.pack(fill=tk.X, side=tk.BOTTOM)

        tk.Button(btn_frame, text="ä¿å­˜å¹¶åŠ è½½OCR", command=lambda: self._save_settings(dialog),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 11, "bold"),
                 padx=30, pady=8, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=20)

        tk.Button(btn_frame, text="å–æ¶ˆ", command=dialog.destroy,
                 bg="#9E9E9E", fg="white", font=(FONT_FAMILY, 11),
                 padx=30, pady=8, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT)

        # åˆ†éš”çº¿
        tk.Frame(dialog, bg="#ddd", height=1).pack(fill=tk.X, side=tk.BOTTOM)

        # å¯æ»šåŠ¨å†…å®¹åŒº - æ”¾åœ¨ä¸­é—´
        content_container = tk.Frame(dialog, bg=COLOR_WHITE)
        content_container.pack(fill=tk.BOTH, expand=True, side=tk.TOP)

        canvas = tk.Canvas(content_container, bg=COLOR_WHITE, highlightthickness=0)
        scrollbar = tk.Scrollbar(content_container, orient=tk.VERTICAL, command=canvas.yview)

        content = tk.Frame(canvas, bg=COLOR_WHITE, padx=20, pady=15)

        canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        canvas_window = canvas.create_window((0, 0), window=content, anchor=tk.NW)

        # æ›´æ–°æ»šåŠ¨åŒºåŸŸ
        def on_frame_configure(event):
            canvas.configure(scrollregion=canvas.bbox("all"))

        content.bind("<Configure>", on_frame_configure)

        # è°ƒæ•´canvasçª—å£å®½åº¦
        def on_canvas_configure(event):
            canvas.itemconfig(canvas_window, width=event.width)

        canvas.bind("<Configure>", on_canvas_configure)

        # é¼ æ ‡æ»šè½®æ”¯æŒ
        def on_mousewheel(event):
            canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

        canvas.bind_all("<MouseWheel>", on_mousewheel)

        # === æ–¹å¼1: æŒ‡å®šå·²æœ‰æ¨¡å‹ç›®å½• ===
        tk.Label(content, text="æ–¹å¼1: æŒ‡å®šå·²æœ‰æ¨¡å‹ç›®å½•", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="å¦‚æœå·²æœ‰æ¨¡å‹æ–‡ä»¶ï¼Œç›´æ¥é€‰æ‹©æ¨¡å‹æ‰€åœ¨ç›®å½•",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        path_frame = tk.Frame(content, bg=COLOR_WHITE)
        path_frame.pack(fill=tk.X, pady=5)

        self.model_dir_var = tk.StringVar(value=self.config.get("model_dir", ""))
        path_entry = tk.Entry(path_frame, textvariable=self.model_dir_var,
                             font=(FONT_FAMILY, 10), width=45)
        path_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        browse_btn = tk.Button(path_frame, text="æµè§ˆ...", command=self._browse_model_dir,
                              bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                              padx=10, cursor="hand2", relief=tk.FLAT)
        browse_btn.pack(side=tk.LEFT, padx=(10, 0))

        # === æ–¹å¼2: ä¸‹è½½æ¨¡å‹åˆ°æŒ‡å®šç›®å½• ===
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)

        tk.Label(content, text="æ–¹å¼2: ä¸‹è½½æ¨¡å‹åˆ°æŒ‡å®šç›®å½•", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="å¦‚æœæ²¡æœ‰æ¨¡å‹ï¼Œé€‰æ‹©ä¸€ä¸ªç›®å½•åç‚¹å‡»ä¸‹è½½ï¼ˆéœ€è¦è”ç½‘ï¼Œçº¦200MBï¼‰",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        download_frame = tk.Frame(content, bg=COLOR_WHITE)
        download_frame.pack(fill=tk.X, pady=5)

        self.download_dir_var = tk.StringVar(value=os.path.join(get_base_dir(), ".paddlex", "official_models"))
        download_entry = tk.Entry(download_frame, textvariable=self.download_dir_var,
                                 font=(FONT_FAMILY, 10), width=45)
        download_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        browse_download_btn = tk.Button(download_frame, text="æµè§ˆ...",
                                       command=lambda: self._browse_download_dir(),
                                       bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9),
                                       padx=10, cursor="hand2", relief=tk.FLAT)
        browse_download_btn.pack(side=tk.LEFT, padx=(10, 0))

        # ä¸‹è½½æŒ‰é’®å’Œè¿›åº¦
        download_btn_frame = tk.Frame(content, bg=COLOR_WHITE)
        download_btn_frame.pack(fill=tk.X, pady=10)

        self.download_btn = tk.Button(download_btn_frame, text="ä¸‹è½½æ¨¡å‹",
                                     command=lambda: self._download_models(dialog),
                                     bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 10, "bold"),
                                     padx=20, pady=5, cursor="hand2", relief=tk.FLAT)
        self.download_btn.pack(side=tk.LEFT)

        self.download_status_label = tk.Label(download_btn_frame, text="", bg=COLOR_WHITE,
                                             fg="#666", font=(FONT_FAMILY, 9))
        self.download_status_label.pack(side=tk.LEFT, padx=15)

        # è¿›åº¦æ¡
        progress_frame = tk.Frame(content, bg=COLOR_WHITE)
        progress_frame.pack(fill=tk.X, pady=5)

        self.download_progress = ttk.Progressbar(progress_frame, length=400, mode='determinate')
        self.download_progress.pack(fill=tk.X)

        self.download_detail_label = tk.Label(progress_frame, text="", bg=COLOR_WHITE,
                                              fg="#999", font=(FONT_FAMILY, 8))
        self.download_detail_label.pack(anchor="w")

        # === è®¾å¤‡é€‰æ‹© ===
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)

        tk.Label(content, text="è®¾å¤‡é€‰æ‹©", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="é€‰æ‹©OCRè¿è¡Œçš„è®¾å¤‡ï¼ˆGPUéœ€è¦å®‰è£…PaddlePaddle-GPUç‰ˆæœ¬ï¼‰",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        device_frame = tk.Frame(content, bg=COLOR_WHITE)
        device_frame.pack(fill=tk.X, pady=5)

        self.device_var = tk.StringVar(value=self.config.get("ocr_device", "cpu"))

        tk.Radiobutton(device_frame, text="CPU - å…¼å®¹æ€§å¥½ï¼Œé€‚åˆæ‰€æœ‰ç”µè„‘",
                      variable=self.device_var, value="cpu",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)
        tk.Radiobutton(device_frame, text="GPU - é€Ÿåº¦å¿«ï¼Œéœ€è¦NVIDIAæ˜¾å¡",
                      variable=self.device_var, value="gpu",
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w", pady=3)

        # æç¤ºä¿¡æ¯
        tk.Label(device_frame,
                text="æç¤ºï¼šä½¿ç”¨GPUéœ€è¦å…ˆå®‰è£… paddlepaddle-gpu\nå¦‚æœªå®‰è£…ï¼Œè¯·è¿è¡Œï¼špip uninstall paddlepaddle && pip install paddlepaddle-gpu",
                bg=COLOR_WHITE, fg="#999", font=(FONT_FAMILY, 8), justify=tk.LEFT).pack(anchor="w", pady=(5, 0))

        # === IOPaint API é…ç½® ===
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)

        tk.Label(content, text="IOPaint API é…ç½®ï¼ˆèƒŒæ™¯ç”ŸæˆåŠŸèƒ½ï¼‰", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w")
        tk.Label(content, text="ç”¨äºè‡ªåŠ¨å»é™¤æ–‡å­—åŒºåŸŸï¼Œç”Ÿæˆå¹²å‡€çš„èƒŒæ™¯å›¾",
                bg=COLOR_WHITE, fg="#666", font=(FONT_FAMILY, 9)).pack(anchor="w", pady=(0, 5))

        # å¯ç”¨å¼€å…³
        inpaint_switch_frame = tk.Frame(content, bg=COLOR_WHITE)
        inpaint_switch_frame.pack(fill=tk.X, pady=5)

        self.inpaint_enabled_var = tk.BooleanVar(value=self.config.get("inpaint_enabled", True))
        tk.Checkbutton(inpaint_switch_frame, text="å¯ç”¨èƒŒæ™¯ç”ŸæˆåŠŸèƒ½",
                      variable=self.inpaint_enabled_var,
                      bg=COLOR_WHITE, font=(FONT_FAMILY, 10)).pack(anchor="w")

        # APIåœ°å€é…ç½®
        api_frame = tk.Frame(content, bg=COLOR_WHITE)
        api_frame.pack(fill=tk.X, pady=5)

        tk.Label(api_frame, text="APIåœ°å€:", bg=COLOR_WHITE, font=(FONT_FAMILY, 9)).pack(anchor="w")
        self.inpaint_api_var = tk.StringVar(value=self.config.get("inpaint_api_url", "http://127.0.0.1:8080/api/v1/inpaint"))
        api_entry = tk.Entry(api_frame, textvariable=self.inpaint_api_var,
                            font=(FONT_FAMILY, 10), width=50)
        api_entry.pack(fill=tk.X, pady=3)

        # æµ‹è¯•æŒ‰é’®
        test_btn_frame = tk.Frame(content, bg=COLOR_WHITE)
        test_btn_frame.pack(fill=tk.X, pady=5)

        tk.Button(test_btn_frame, text="æµ‹è¯•è¿æ¥", command=self._test_inpaint_api,
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9),
                 padx=15, pady=3, cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT)

        self.api_test_label = tk.Label(test_btn_frame, text="", bg=COLOR_WHITE,
                                       fg="#666", font=(FONT_FAMILY, 9))
        self.api_test_label.pack(side=tk.LEFT, padx=10)

        # è¯´æ˜ä¿¡æ¯
        info_frame = tk.Frame(content, bg="#fff3cd", padx=10, pady=8)
        info_frame.pack(fill=tk.X, pady=5)

        tk.Label(info_frame, text="ğŸ“Œ ä½¿ç”¨è¯´æ˜", bg="#fff3cd",
                font=(FONT_FAMILY, 9, "bold"), fg="#856404").pack(anchor="w")
        tk.Label(info_frame,
                text="1. å®‰è£…IOPaintï¼špip install iopaint\n"
                     "2. å¯åŠ¨æœåŠ¡ï¼šiopaint start --host 127.0.0.1 --port 8080\n"
                     "3. æˆ–ä½¿ç”¨åœ¨çº¿æœåŠ¡ï¼ˆä¿®æ”¹APIåœ°å€ï¼‰\n"
                     "4. ä½¿ç”¨å‰è¯·å…ˆæµ‹è¯•è¿æ¥",
                bg="#fff3cd", fg="#856404", font=(FONT_FAMILY, 8),
                justify=tk.LEFT).pack(anchor="w", pady=(3, 0))

        # æ¨¡å‹çŠ¶æ€æ˜¾ç¤º
        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=10)

        status_frame = tk.Frame(content, bg="#f5f5f5", padx=10, pady=10)
        status_frame.pack(fill=tk.X)

        self.model_status_label = tk.Label(status_frame, text="", bg="#f5f5f5",
                                           font=(FONT_FAMILY, 9), justify=tk.LEFT)
        self.model_status_label.pack(anchor="w")

        self._check_model_status()

        # ç»‘å®šè·¯å¾„å˜åŒ–äº‹ä»¶
        self.model_dir_var.trace_add("write", lambda *args: self._check_model_status())

    def _browse_model_dir(self):
        """æµè§ˆé€‰æ‹©æ¨¡å‹ç›®å½•"""
        current_dir = self.model_dir_var.get()
        if not os.path.exists(current_dir):
            current_dir = get_base_dir()

        dir_path = filedialog.askdirectory(
            title="é€‰æ‹©OCRæ¨¡å‹ç›®å½•ï¼ˆåŒ…å« PP-OCRv5_server_det ç­‰æ–‡ä»¶å¤¹ï¼‰",
            initialdir=current_dir
        )
        if dir_path:
            self.model_dir_var.set(dir_path)

    def _browse_download_dir(self):
        """æµè§ˆé€‰æ‹©ä¸‹è½½ç›®å½•"""
        current_dir = self.download_dir_var.get()
        if not os.path.exists(current_dir):
            current_dir = get_base_dir()

        dir_path = filedialog.askdirectory(
            title="é€‰æ‹©æ¨¡å‹ä¸‹è½½ç›®å½•",
            initialdir=current_dir
        )
        if dir_path:
            self.download_dir_var.set(dir_path)

    def _download_models(self, dialog):
        """ä¸‹è½½OCRæ¨¡å‹ - ä½¿ç”¨ç›´æ¥URLä¸‹è½½"""
        download_dir = self.download_dir_var.get()

        if not download_dir:
            messagebox.showwarning("è­¦å‘Š", "è¯·å…ˆé€‰æ‹©ä¸‹è½½ç›®å½•ï¼")
            return

        # åˆ›å»ºç›®å½•
        os.makedirs(download_dir, exist_ok=True)

        # ç¦ç”¨ä¸‹è½½æŒ‰é’®
        self.download_btn.config(state=tk.DISABLED, text="ä¸‹è½½ä¸­...")
        self.download_status_label.config(text="æ­£åœ¨å‡†å¤‡ä¸‹è½½...")
        self.download_progress['value'] = 0

        # éœ€è¦ä¸‹è½½çš„æ¨¡å‹åˆ—è¡¨
        models_to_download = [
            ("PP-OCRv5_server_det", "æ–‡å­—æ£€æµ‹æ¨¡å‹", "PP-OCRv5_server_det_infer.tar"),
            ("PP-OCRv5_server_rec", "æ–‡å­—è¯†åˆ«æ¨¡å‹", "PP-OCRv5_server_rec_infer.tar"),
            ("PP-LCNet_x1_0_doc_ori", "æ–‡æ¡£æ–¹å‘åˆ†ç±»", "PP-LCNet_x1_0_doc_ori_infer.tar"),
            ("PP-LCNet_x1_0_textline_ori", "æ–‡æœ¬è¡Œæ–¹å‘", "PP-LCNet_x1_0_textline_ori_infer.tar"),
            ("UVDoc", "æ–‡æ¡£çŸ«æ­£", "UVDoc_infer.tar"),
        ]

        base_url = "https://paddle-model-ecology.bj.bcebos.com/paddlex/official_inference_model/paddle3.0.0"

        def download_task():
            import urllib.request
            import tarfile

            total_models = len(models_to_download)
            downloaded = 0

            for model_name, desc, tar_file in models_to_download:
                model_path = os.path.join(download_dir, model_name)

                # å¦‚æœæ¨¡å‹å·²å­˜åœ¨ï¼Œè·³è¿‡
                if os.path.exists(model_path):
                    downloaded += 1
                    progress = int((downloaded / total_models) * 100)
                    dialog.after(0, lambda p=progress, d=desc: self._update_download_progress(p, f"{d} å·²å­˜åœ¨ï¼Œè·³è¿‡"))
                    continue

                url = f"{base_url}/{tar_file}"
                tar_path = os.path.join(download_dir, tar_file)

                try:
                    # æ›´æ–°çŠ¶æ€
                    dialog.after(0, lambda d=desc: self.download_status_label.config(text=f"æ­£åœ¨ä¸‹è½½: {d}"))
                    dialog.after(0, lambda d=desc: self.download_detail_label.config(text=f"ä» {url}"))

                    # ä¸‹è½½æ–‡ä»¶ï¼ˆå¸¦è¿›åº¦ï¼‰
                    def reporthook(block_num, block_size, total_size):
                        if total_size > 0:
                            downloaded_size = block_num * block_size
                            percent = min(int((downloaded_size / total_size) * 100), 100)
                            size_mb = downloaded_size / (1024 * 1024)
                            total_mb = total_size / (1024 * 1024)
                            # è®¡ç®—æ€»è¿›åº¦
                            model_progress = downloaded / total_models
                            file_progress = (downloaded_size / total_size) / total_models
                            overall = int((model_progress + file_progress) * 100)
                            dialog.after(0, lambda o=overall, s=size_mb, t=total_mb:
                                self._update_download_progress(o, f"ä¸‹è½½ä¸­: {s:.1f}MB / {t:.1f}MB"))

                    urllib.request.urlretrieve(url, tar_path, reporthook)

                    # è§£å‹
                    dialog.after(0, lambda d=desc: self.download_status_label.config(text=f"æ­£åœ¨è§£å‹: {d}"))

                    with tarfile.open(tar_path, 'r:*') as tar:
                        tar.extractall(download_dir)

                    # åˆ é™¤taræ–‡ä»¶
                    os.remove(tar_path)

                    # é‡å‘½åæ–‡ä»¶å¤¹ï¼ˆå»æ‰_inferåç¼€ï¼‰
                    infer_path = os.path.join(download_dir, f"{model_name}_infer")
                    if os.path.exists(infer_path) and not os.path.exists(model_path):
                        os.rename(infer_path, model_path)

                    downloaded += 1
                    progress = int((downloaded / total_models) * 100)
                    dialog.after(0, lambda p=progress, d=desc: self._update_download_progress(p, f"{d} ä¸‹è½½å®Œæˆ"))

                except Exception as e:
                    dialog.after(0, lambda d=desc, err=str(e):
                        self.download_status_label.config(text=f"{d} ä¸‹è½½å¤±è´¥: {err[:50]}"))
                    # æ¸…ç†å¯èƒ½çš„æ®‹ç•™æ–‡ä»¶
                    if os.path.exists(tar_path):
                        try:
                            os.remove(tar_path)
                        except:
                            pass

            # ä¸‹è½½å®Œæˆ
            dialog.after(0, lambda: self._download_complete(download_dir, dialog))

        threading.Thread(target=download_task, daemon=True).start()

    def _update_download_progress(self, progress, detail):
        """æ›´æ–°ä¸‹è½½è¿›åº¦"""
        self.download_progress['value'] = progress
        self.download_detail_label.config(text=detail)

    def _download_complete(self, download_dir, dialog):
        """ä¸‹è½½å®Œæˆå¤„ç†"""
        self.download_btn.config(state=tk.NORMAL, text="ä¸‹è½½æ¨¡å‹")
        self.download_progress['value'] = 100
        self.download_status_label.config(text="ä¸‹è½½å®Œæˆï¼")
        self.download_detail_label.config(text="")

        # è®¾ç½®æ¨¡å‹ç›®å½•
        self.model_dir_var.set(download_dir)
        self._check_model_status()

        messagebox.showinfo("æˆåŠŸ",
            f"æ¨¡å‹ä¸‹è½½å®Œæˆï¼\n\nä¸‹è½½ç›®å½•:\n{download_dir}\n\nå·²è‡ªåŠ¨è®¾ç½®ä¸ºæ¨¡å‹ç›®å½•ï¼Œç‚¹å‡»'ä¿å­˜å¹¶åŠ è½½OCR'å³å¯ä½¿ç”¨ã€‚")

    def _check_model_status(self):
        """æ£€æŸ¥æ¨¡å‹çŠ¶æ€"""
        model_dir = self.model_dir_var.get()

        required_models = [
            ("PP-OCRv5_server_det", "æ–‡å­—æ£€æµ‹æ¨¡å‹"),
            ("PP-OCRv5_server_rec", "æ–‡å­—è¯†åˆ«æ¨¡å‹"),
        ]
        optional_models = [
            ("PP-LCNet_x1_0_doc_ori", "æ–‡æ¡£æ–¹å‘åˆ†ç±»"),
            ("PP-LCNet_x1_0_textline_ori", "æ–‡æœ¬è¡Œæ–¹å‘"),
            ("UVDoc", "æ–‡æ¡£çŸ«æ­£"),
        ]

        status_lines = []

        if not model_dir:
            status_lines.append("è¯·é€‰æ‹©æˆ–ä¸‹è½½æ¨¡å‹ç›®å½•")
        elif not os.path.exists(model_dir):
            status_lines.append("ç›®å½•ä¸å­˜åœ¨ï¼Œè¯·é€‰æ‹©æœ‰æ•ˆç›®å½•æˆ–ä¸‹è½½æ¨¡å‹")
        else:
            all_required = True
            for model_name, desc in required_models:
                model_path = os.path.join(model_dir, model_name)
                if os.path.exists(model_path):
                    status_lines.append(f"[OK] {desc} ({model_name})")
                else:
                    status_lines.append(f"[X] {desc} ({model_name}) - ç¼ºå¤±!")
                    all_required = False

            for model_name, desc in optional_models:
                model_path = os.path.join(model_dir, model_name)
                if os.path.exists(model_path):
                    status_lines.append(f"[OK] {desc} ({model_name})")
                else:
                    status_lines.append(f"[  ] {desc} ({model_name}) - å¯é€‰")

            if all_required:
                status_lines.insert(0, "å½“å‰æ¨¡å‹çŠ¶æ€: å¯ç”¨\n")
            else:
                status_lines.insert(0, "å½“å‰æ¨¡å‹çŠ¶æ€: ç¼ºå°‘å¿…éœ€æ¨¡å‹!\n")

        self.model_status_label.config(text="\n".join(status_lines))

    def _test_inpaint_api(self):
        """æµ‹è¯•IOPaint APIè¿æ¥"""
        api_url = self.inpaint_api_var.get()

        if not api_url:
            self.api_test_label.config(text="âŒ è¯·è¾“å…¥APIåœ°å€", fg="red")
            return

        self.api_test_label.config(text="â³ æµ‹è¯•ä¸­...", fg="blue")

        def test():
            try:
                # åˆ›å»ºä¸€ä¸ªå°çš„æµ‹è¯•å›¾ç‰‡å’Œè’™ç‰ˆ
                test_img = Image.new("RGB", (64, 64), (255, 255, 255))
                test_mask = Image.new("L", (64, 64), 0)

                # Base64ç¼–ç 
                def to_b64(img):
                    buffer = BytesIO()
                    img.save(buffer, "PNG")
                    return base64.b64encode(buffer.getvalue()).decode()

                payload = {
                    "image": to_b64(test_img),
                    "mask": to_b64(test_mask),
                    "ldm_steps": 1,
                    "hd_strategy": "Original"
                }

                response = requests.post(api_url, json=payload, timeout=10)

                if response.status_code == 200:
                    self.root.after(0, lambda: self.api_test_label.config(
                        text="âœ“ è¿æ¥æˆåŠŸï¼", fg="green"))
                else:
                    self.root.after(0, lambda: self.api_test_label.config(
                        text=f"âŒ é”™è¯¯: {response.status_code}", fg="red"))

            except requests.exceptions.ConnectionError:
                self.root.after(0, lambda: self.api_test_label.config(
                    text="âŒ æ— æ³•è¿æ¥ï¼Œè¯·æ£€æŸ¥æœåŠ¡æ˜¯å¦å¯åŠ¨", fg="red"))
            except Exception as e:
                self.root.after(0, lambda: self.api_test_label.config(
                    text=f"âŒ æµ‹è¯•å¤±è´¥: {str(e)[:30]}", fg="red"))

        threading.Thread(target=test, daemon=True).start()

    def _save_settings(self, dialog):
        """ä¿å­˜è®¾ç½®å¹¶é‡æ–°åŠ è½½OCR"""
        new_model_dir = self.model_dir_var.get()
        new_device = self.device_var.get()  # è·å–è®¾å¤‡é€‰æ‹©
        new_inpaint_enabled = self.inpaint_enabled_var.get()
        new_inpaint_api = self.inpaint_api_var.get()

        if not new_model_dir:
            messagebox.showwarning("è­¦å‘Š", "è¯·å…ˆé€‰æ‹©æ¨¡å‹ç›®å½•ï¼")
            return

        # æ£€æŸ¥å¿…éœ€æ¨¡å‹æ˜¯å¦å­˜åœ¨
        det_model = os.path.join(new_model_dir, "PP-OCRv5_server_det")
        rec_model = os.path.join(new_model_dir, "PP-OCRv5_server_rec")

        if not os.path.exists(det_model) or not os.path.exists(rec_model):
            result = messagebox.askyesno("è­¦å‘Š",
                "æ¨¡å‹ç›®å½•ç¼ºå°‘å¿…éœ€çš„æ¨¡å‹æ–‡ä»¶ï¼\n\n"
                "éœ€è¦:\n- PP-OCRv5_server_det\n- PP-OCRv5_server_rec\n\n"
                "æ˜¯å¦ä»ç„¶ä¿å­˜ï¼Ÿï¼ˆOCRåŠŸèƒ½å°†æ— æ³•ä½¿ç”¨ï¼‰")
            if not result:
                return

        # ä¿å­˜é…ç½®
        self.config["model_dir"] = new_model_dir
        self.config["ocr_device"] = new_device  # ä¿å­˜è®¾å¤‡é€‰æ‹©
        self.config["inpaint_enabled"] = new_inpaint_enabled  # ä¿å­˜IOPaintå¼€å…³
        self.config["inpaint_api_url"] = new_inpaint_api  # ä¿å­˜APIåœ°å€
        save_config(self.config)

        # å…³é—­å¯¹è¯æ¡†
        dialog.destroy()

        # é‡æ–°åŠ è½½OCR
        self.ocr = None
        device_name = "GPU" if new_device == "gpu" else "CPU"
        self.update_status(f"æ­£åœ¨ä½¿ç”¨ {device_name} åŠ è½½OCRæ¨¡å‹...")
        threading.Thread(target=self.init_ocr, daemon=True).start()

        messagebox.showinfo("æˆåŠŸ",
            f"è®¾ç½®å·²ä¿å­˜ï¼\n\n"
            f"OCRæ¨¡å‹ç›®å½•:\n{new_model_dir}\n\n"
            f"è¿è¡Œè®¾å¤‡: {device_name}\n\n"
            f"èƒŒæ™¯ç”ŸæˆåŠŸèƒ½: {'å·²å¯ç”¨' if new_inpaint_enabled else 'å·²ç¦ç”¨'}\n"
            f"IOPaint API: {new_inpaint_api}\n\n"
            f"OCRæ¨¡å‹æ­£åœ¨åå°åŠ è½½...")



    # ==================== æ–°å¢åŠŸèƒ½ï¼šå…¨é€‰å’Œå¤åˆ¶ç²˜è´´ ====================

    def select_all_boxes(self):
        """å…¨é€‰å½“å‰é¡µæ‰€æœ‰æ–‡æœ¬æ¡†"""
        if not self.text_boxes:
            self.update_status("å½“å‰é¡µæ²¡æœ‰æ–‡æœ¬æ¡†")
            return

        # é€‰ä¸­æ‰€æœ‰æ¡†
        self.selected_boxes = list(range(len(self.text_boxes)))
        self.selected_box_index = 0 if self.text_boxes else -1

        # åˆ·æ–°ç•Œé¢
        self.refresh_canvas()
        self.update_property_panel()

        # æ›´æ–°åˆ—è¡¨æ¡†é€‰æ‹©
        self.box_listbox.selection_clear(0, tk.END)
        for idx in self.selected_boxes:
            self.box_listbox.selection_set(idx)

        self.update_status(f"å·²é€‰ä¸­å½“å‰é¡µæ‰€æœ‰ {len(self.text_boxes)} ä¸ªæ–‡æœ¬æ¡† âœ“")

    def copy_boxes(self):
        """å¤åˆ¶é€‰ä¸­çš„æ–‡æœ¬æ¡†"""
        if not self.selected_boxes:
            self.update_status("è¯·å…ˆé€‰ä¸­è¦å¤åˆ¶çš„æ–‡æœ¬æ¡†")
            return

        self.clipboard_boxes = []
        for idx in self.selected_boxes:
            if 0 <= idx < len(self.text_boxes):
                self.clipboard_boxes.append(self.text_boxes[idx].copy())

        self.update_status(f"å·²å¤åˆ¶ {len(self.clipboard_boxes)} ä¸ªæ–‡æœ¬æ¡†")

    def paste_boxes(self):
        """ç²˜è´´æ–‡æœ¬æ¡†"""
        if not self.clipboard_boxes:
            self.update_status("å‰ªè´´æ¿ä¸ºç©º")
            return

        self.save_state()

        offset = 20
        new_boxes = []
        for box in self.clipboard_boxes:
            new_box = box.copy()
            new_box.x += offset
            new_box.y += offset
            self.text_boxes.append(new_box)
            new_boxes.append(new_box)

        start_idx = len(self.text_boxes) - len(new_boxes)
        self.selected_boxes = list(range(start_idx, len(self.text_boxes)))
        self.selected_box_index = self.selected_boxes[0] if self.selected_boxes else -1

        self.refresh_canvas()
        self.update_listbox()
        self.mark_unsaved()
        self.mark_unsaved()
        self.update_status(f"å·²ç²˜è´´ {len(new_boxes)} ä¸ªæ–‡æœ¬æ¡†")

    def move_box_by_key(self, dx, dy):
        """ä½¿ç”¨æ–¹å‘é”®ç§»åŠ¨æ–‡æœ¬æ¡†"""
        if self.selected_box_index < 0:
            return

        box = self.text_boxes[self.selected_box_index]
        box.x = max(0, box.x + dx)
        box.y = max(0, box.y + dy)

        self.refresh_canvas()
        self.update_property_panel()
        self.mark_unsaved()

    # ==================== æ–°å¢åŠŸèƒ½ï¼šå®Œæ•´å¯¹é½å·¥å…· ====================

    def show_align_dialog(self):
        """æ˜¾ç¤ºå¯¹é½å·¥å…·å¯¹è¯æ¡†"""
        if len(self.selected_boxes) < 2:
            messagebox.showinfo("æç¤º", "è¯·å…ˆä½¿ç”¨Ctrl+ç‚¹å‡»é€‰ä¸­è‡³å°‘2ä¸ªæ–‡æœ¬æ¡†")
            return

        dialog = tk.Toplevel(self.root)
        dialog.title("å¯¹é½ä¸åˆ†å¸ƒå·¥å…·")
        dialog.geometry("450x550")
        dialog.configure(bg=COLOR_WHITE)
        dialog.transient(self.root)
        dialog.grab_set()

        dialog.update_idletasks()
        x = (dialog.winfo_screenwidth() - 450) // 2
        y = (dialog.winfo_screenheight() - 550) // 2
        dialog.geometry(f"+{x}+{y}")

        title_frame = tk.Frame(dialog, bg=COLOR_THEME, height=40)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)
        tk.Label(title_frame, text=f"  å¯¹é½ä¸åˆ†å¸ƒ - å·²é€‰ä¸­ {len(self.selected_boxes)} ä¸ªæ¡†",
                bg=COLOR_THEME, fg="white",
                font=(FONT_FAMILY, 11, "bold")).pack(side=tk.LEFT, pady=8)

        content = tk.Frame(dialog, bg=COLOR_WHITE, padx=20, pady=15)
        content.pack(fill=tk.BOTH, expand=True)

        # æ°´å¹³å¯¹é½
        tk.Label(content, text="æ°´å¹³å¯¹é½", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(5, 5))

        h_frame = tk.Frame(content, bg=COLOR_WHITE)
        h_frame.pack(fill=tk.X, pady=5)

        tk.Button(h_frame, text="å·¦å¯¹é½", command=lambda: self.align_boxes("left"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_frame, text="æ°´å¹³å±…ä¸­", command=lambda: self.align_boxes("center_h"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(h_frame, text="å³å¯¹é½", command=lambda: self.align_boxes("right"),
                 bg=COLOR_BLUE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # å‚ç›´å¯¹é½
        tk.Label(content, text="å‚ç›´å¯¹é½", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        v_frame = tk.Frame(content, bg=COLOR_WHITE)
        v_frame.pack(fill=tk.X, pady=5)

        tk.Button(v_frame, text="é¡¶å¯¹é½", command=lambda: self.align_boxes("top"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_frame, text="å‚ç›´å±…ä¸­", command=lambda: self.align_boxes("center_v"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(v_frame, text="åº•å¯¹é½", command=lambda: self.align_boxes("bottom"),
                 bg=COLOR_GREEN, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # åˆ†å¸ƒ
        tk.Label(content, text="å‡åŒ€åˆ†å¸ƒ (éœ€è¦3ä¸ªæˆ–ä»¥ä¸Š)", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        dist_frame = tk.Frame(content, bg=COLOR_WHITE)
        dist_frame.pack(fill=tk.X, pady=5)

        tk.Button(dist_frame, text="æ°´å¹³ç­‰é—´è·", command=lambda: self.distribute_boxes("horizontal"),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(dist_frame, text="å‚ç›´ç­‰é—´è·", command=lambda: self.distribute_boxes("vertical"),
                 bg=COLOR_ORANGE, fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # å°ºå¯¸ç»Ÿä¸€
        tk.Label(content, text="å°ºå¯¸ç»Ÿä¸€ (ä»¥ç¬¬ä¸€ä¸ªé€‰ä¸­æ¡†ä¸ºåŸºå‡†)", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        size_frame = tk.Frame(content, bg=COLOR_WHITE)
        size_frame.pack(fill=tk.X, pady=5)

        tk.Button(size_frame, text="ç»Ÿä¸€å®½åº¦", command=lambda: self.unify_size("width"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="ç»Ÿä¸€é«˜åº¦", command=lambda: self.unify_size("height"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(size_frame, text="ç»Ÿä¸€å¤§å°", command=lambda: self.unify_size("both"),
                 bg=COLOR_PURPLE, fg="white", font=(FONT_FAMILY, 9), width=10,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        # å¯¹é½åˆ°ç”»å¸ƒ
        tk.Label(content, text="å¯¹é½åˆ°ç”»å¸ƒ", bg=COLOR_WHITE,
                font=(FONT_FAMILY, 10, "bold")).pack(anchor="w", pady=(15, 5))

        canvas_frame = tk.Frame(content, bg=COLOR_WHITE)
        canvas_frame.pack(fill=tk.X, pady=5)

        tk.Button(canvas_frame, text="ç”»å¸ƒæ°´å¹³å±…ä¸­", command=lambda: self.align_to_canvas("h"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)
        tk.Button(canvas_frame, text="ç”»å¸ƒå‚ç›´å±…ä¸­", command=lambda: self.align_to_canvas("v"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=15,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        canvas_frame2 = tk.Frame(content, bg=COLOR_WHITE)
        canvas_frame2.pack(fill=tk.X, pady=5)

        tk.Button(canvas_frame2, text="ç”»å¸ƒå®Œå…¨å±…ä¸­", command=lambda: self.align_to_canvas("center"),
                 bg="#00897B", fg="white", font=(FONT_FAMILY, 9), width=32,
                 cursor="hand2", relief=tk.FLAT).pack(side=tk.LEFT, padx=2)

        tk.Frame(content, bg="#ddd", height=1).pack(fill=tk.X, pady=15)
        tk.Button(content, text="å…³é—­", command=dialog.destroy,
                 bg=COLOR_GRAY, fg="white", font=(FONT_FAMILY, 10),
                 width=15, cursor="hand2", relief=tk.FLAT).pack()

    def distribute_boxes(self, direction):
        """å‡åŒ€åˆ†å¸ƒæ–‡æœ¬æ¡†"""
        if len(self.selected_boxes) < 3:
            messagebox.showinfo("æç¤º", "å‡åŒ€åˆ†å¸ƒéœ€è¦è‡³å°‘é€‰ä¸­3ä¸ªæ–‡æœ¬æ¡†")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        if direction == "horizontal":
            boxes.sort(key=lambda b: b.x)
            first = boxes[0]
            last = boxes[-1]

            total_width = sum(b.width for b in boxes)
            total_space = (last.x + last.width) - first.x - total_width
            gap = total_space / (len(boxes) - 1) if len(boxes) > 1 else 0

            current_x = first.x + first.width
            for box in boxes[1:-1]:
                box.x = int(current_x + gap)
                current_x = box.x + box.width

        elif direction == "vertical":
            boxes.sort(key=lambda b: b.y)
            first = boxes[0]
            last = boxes[-1]

            total_height = sum(b.height for b in boxes)
            total_space = (last.y + last.height) - first.y - total_height
            gap = total_space / (len(boxes) - 1) if len(boxes) > 1 else 0

            current_y = first.y + first.height
            for box in boxes[1:-1]:
                box.y = int(current_y + gap)
                current_y = box.y + box.height

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"å·²å‡åŒ€åˆ†å¸ƒ {len(self.selected_boxes)} ä¸ªæ¡† âœ“")

    def unify_size(self, size_type):
        """ç»Ÿä¸€æ–‡æœ¬æ¡†å¤§å°"""
        if len(self.selected_boxes) < 2:
            self.update_status("è¯·Ctrl+ç‚¹å‡»é€‰ä¸­è‡³å°‘2ä¸ªæ¡†")
            return

        self.save_state()
        boxes = [self.text_boxes[i] for i in self.selected_boxes]

        base_box = boxes[0]

        for box in boxes[1:]:
            if size_type in ["width", "both"]:
                box.width = base_box.width
            if size_type in ["height", "both"]:
                box.height = base_box.height

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"å·²ç»Ÿä¸€ {len(self.selected_boxes)} ä¸ªæ¡†çš„å°ºå¯¸ âœ“")

    def align_to_canvas(self, align_type):
        """å¯¹é½åˆ°ç”»å¸ƒä¸­å¿ƒ"""
        if not self.selected_boxes or not self.original_image:
            return

        self.save_state()

        img_w, img_h = self.original_image.size
        center_x = img_w // 2
        center_y = img_h // 2

        for idx in self.selected_boxes:
            box = self.text_boxes[idx]

            if align_type == "h":
                box.x = center_x - box.width // 2
            elif align_type == "v":
                box.y = center_y - box.height // 2
            elif align_type == "center":
                box.x = center_x - box.width // 2
                box.y = center_y - box.height // 2

        self.refresh_canvas()
        self.mark_unsaved()
        self.update_status(f"å·²å¯¹é½åˆ°ç”»å¸ƒä¸­å¿ƒ âœ“")

    # ==================== æ–°å¢åŠŸèƒ½ï¼šè‡ªåŠ¨ä¿å­˜ ====================

    def start_autosave(self):
        """å¯åŠ¨è‡ªåŠ¨ä¿å­˜"""
        interval = self.config.get("autosave_interval", 300) * 1000
        self.autosave_timer = self.root.after(interval, self.auto_save)

    def stop_autosave(self):
        """åœæ­¢è‡ªåŠ¨ä¿å­˜"""
        if self.autosave_timer:
            self.root.after_cancel(self.autosave_timer)
            self.autosave_timer = None

    def auto_save(self):
        """è‡ªåŠ¨ä¿å­˜"""
        if self.has_unsaved_changes and self.pages:
            try:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                autosave_path = os.path.join(self.autosave_dir, f"autosave_{timestamp}.json")

                self.save_current_page()

                pages_data = []
                for page in self.pages:
                    pages_data.append({
                        "original_path": page["original_path"],
                        "original_size": page.get("original_size", page["image"].size),
                        "edit_scale": page.get("edit_scale", 1.0),
                        "bg_path": page.get("bg_path"),
                        "bg_original_path": page.get("bg_original_path"),
                        "text_boxes": page.get("text_boxes", []),
                        "layers": page.get("layers", []),
                    })

                with open(autosave_path, 'w', encoding='utf-8') as f:
                    json.dump({"version": 3, "pages": pages_data, "current_page": self.current_page_index},
                             f, ensure_ascii=False, indent=2)

                self.cleanup_autosave_files()
                print(f"è‡ªåŠ¨ä¿å­˜å®Œæˆ: {autosave_path}")

            except Exception as e:
                print(f"è‡ªåŠ¨ä¿å­˜å¤±è´¥: {e}")

        self.start_autosave()

    def cleanup_autosave_files(self):
        """æ¸…ç†æ—§çš„è‡ªåŠ¨ä¿å­˜æ–‡ä»¶"""
        try:
            autosave_files = [f for f in os.listdir(self.autosave_dir) if f.startswith("autosave_")]
            autosave_files.sort(reverse=True)

            for old_file in autosave_files[10:]:
                try:
                    os.remove(os.path.join(self.autosave_dir, old_file))
                except:
                    pass
        except:
            pass

    def mark_unsaved(self):
        """æ ‡è®°æœ‰æœªä¿å­˜çš„æ›´æ”¹"""
        self.has_unsaved_changes = True
        if hasattr(self, 'autosave_indicator'):
            self.autosave_indicator.config(fg="#FFC107")

    def mark_saved(self):
        """æ ‡è®°å·²ä¿å­˜"""
        self.has_unsaved_changes = False
        if hasattr(self, 'autosave_indicator'):
            self.autosave_indicator.config(fg="#4CAF50")

    def on_closing(self):
        """çª—å£å…³é—­äº‹ä»¶"""
        if self.has_unsaved_changes:
            result = messagebox.askyesnocancel(
                "æœªä¿å­˜çš„æ›´æ”¹",
                "æ˜¯å¦ä¿å­˜å½“å‰é¡¹ç›®ï¼Ÿ\n\næ˜¯ - ä¿å­˜å¹¶é€€å‡º\nå¦ - ä¸ä¿å­˜é€€å‡º\nå–æ¶ˆ - è¿”å›ç¼–è¾‘"
            )
            if result is None:
                return
            elif result:
                self.save_project()

        self.stop_autosave()
        self.root.destroy()

    # ==================== æ–°å¢åŠŸèƒ½ï¼šPDFå¯¼å…¥ ====================

    def import_pdf(self):
        """å¯¼å…¥PDFæ–‡ä»¶ - ä½¿ç”¨PyMuPDFï¼Œç®€å•å¿«é€Ÿ"""
        if not PDF_SUPPORT:
            messagebox.showerror("éœ€è¦å®‰è£…åº“",
                "PDFè½¬å›¾ç‰‡éœ€è¦å®‰è£… PyMuPDF\n\n"
                "è¯·è¿è¡Œä»¥ä¸‹å‘½ä»¤:\n"
                "pip install PyMuPDF\n\n"
                "æˆ–è€…:\n"
                "1. ä½¿ç”¨åœ¨çº¿å·¥å…·å°†PDFè½¬ä¸ºå›¾ç‰‡\n"
                "2. ç„¶åç”¨'å¯¼å…¥å›¾ç‰‡'åŠŸèƒ½å¯¼å…¥")
            return

        file_path = filedialog.askopenfilename(
            title="é€‰æ‹©PDFæ–‡ä»¶",
            filetypes=[("PDFæ–‡ä»¶", "*.pdf")]
        )
        if not file_path:
            return

        self.update_status("æ­£åœ¨è½¬æ¢PDF...")

        def convert_pdf():
            try:
                self.root.after(0, lambda: self.update_status("æ­£åœ¨è§£æPDF..."))

                # æ‰“å¼€PDF
                doc = fitz.open(file_path)
                page_count = len(doc)

                if page_count == 0:
                    self.root.after(0, lambda: messagebox.showerror("é”™è¯¯", "PDFæ–‡ä»¶ä¸ºç©º"))
                    doc.close()
                    return

                # è¯¢é—®æ˜¯å¦æ¸…ç©ºç°æœ‰é¡µé¢
                if self.pages:
                    result = messagebox.askyesnocancel(
                        "æç¤º",
                        f"PDFå…± {page_count} é¡µã€‚\n\næ˜¯å¦æ¸…ç©ºç°æœ‰é¡µé¢ï¼Ÿ\n\n"
                        "æ˜¯ - æ¸…ç©ºåå¯¼å…¥\nå¦ - è¿½åŠ åˆ°ç°æœ‰é¡µé¢\nå–æ¶ˆ - å–æ¶ˆå¯¼å…¥"
                    )
                    if result is None:
                        self.root.after(0, lambda: self.update_status("å·²å–æ¶ˆ"))
                        doc.close()
                        return
                    elif result:
                        self.root.after(0, lambda: setattr(self, 'pages', []))

                # åˆ›å»ºä¸´æ—¶ç›®å½•
                temp_dir = os.path.join(get_base_dir(), "temp_pdf_imports")
                os.makedirs(temp_dir, exist_ok=True)

                start_index = len(self.pages)

                # è½¬æ¢æ¯ä¸€é¡µ
                for page_num in range(page_count):
                    self.root.after(0, lambda idx=page_num+1, total=page_count:
                        self.update_status(f"æ­£åœ¨è½¬æ¢ç¬¬ {idx}/{total} é¡µ..."))

                    # è·å–é¡µé¢
                    page = doc[page_num]

                    # è½¬æ¢ä¸ºå›¾ç‰‡ï¼ˆ200 DPIé«˜è´¨é‡ï¼‰
                    zoom = 200 / 72  # PDFé»˜è®¤72 DPIï¼Œæå‡åˆ°200 DPI
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)

                    # ä¿å­˜ä¸ºPNG
                    pdf_basename = os.path.splitext(os.path.basename(file_path))[0]
                    temp_path = os.path.join(temp_dir, f"{pdf_basename}_page_{page_num+1:03d}.png")
                    pix.save(temp_path)

                    # è½¬æ¢ä¸ºPIL Image
                    img_data = pix.tobytes("png")
                    from io import BytesIO
                    img = Image.open(BytesIO(img_data))

                    # æ·»åŠ åˆ°é¡µé¢
                    original_size = img.size
                    edit_img, edit_scale = self._resize_image_for_edit(img)

                    page_data = {
                        "original_path": temp_path,
                        "original_size": original_size,
                        "edit_scale": edit_scale,
                        "bg_path": None,
                        "image": edit_img,
                        "text_boxes": [],
                        "layers": []
                    }
                    self.pages.append(page_data)

                # å…³é—­PDF
                doc.close()

                # æ›´æ–°ç•Œé¢
                self.root.after(0, lambda: setattr(self, 'current_page_index', start_index))
                self.root.after(0, self.load_current_page)
                self.root.after(0, self.update_page_label)
                self.root.after(0, self.update_thumbnails)
                self.root.after(0, lambda: self.placeholder_label.place_forget())
                self.root.after(0, lambda: self.update_status(f"PDFè½¬æ¢æˆåŠŸï¼å…± {page_count} é¡µ"))
                self.root.after(0, lambda: messagebox.showinfo("æˆåŠŸ",
                    f"PDFè½¬æ¢æˆåŠŸï¼\n\n"
                    f"å…±è½¬æ¢ {page_count} é¡µ\n"
                    f"å›¾ç‰‡ä¿å­˜åœ¨ï¼š{temp_dir}\n\n"
                    f"ç°åœ¨å¯ä»¥è¿›è¡ŒOCRè¯†åˆ«äº†"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"PDFè½¬æ¢å¤±è´¥:\n{error_msg}")
                self.root.after(0, lambda: messagebox.showerror("é”™è¯¯",
                    f"PDFè½¬æ¢å¤±è´¥:\n\n{str(e)}\n\n"
                    f"å»ºè®®:\n"
                    f"1. æ£€æŸ¥PDFæ–‡ä»¶æ˜¯å¦æŸå\n"
                    f"2. æˆ–ä½¿ç”¨åœ¨çº¿å·¥å…·è½¬æ¢åå¯¼å…¥å›¾ç‰‡"))
                self.root.after(0, lambda: self.update_status("PDFè½¬æ¢å¤±è´¥"))

        threading.Thread(target=convert_pdf, daemon=True).start()

    # ==================== æ–°å¢åŠŸèƒ½ï¼šPDFå¯¼å‡º ====================

    def export_as_pdf(self):
        return export_feature.export_as_pdf(self)

    # ==================== æ–°å¢åŠŸèƒ½ï¼šå›¾ç‰‡å¯¼å‡º ====================

    def export_as_images(self):
        return export_feature.export_as_images(self)

    def _show_image_format_dialog(self, folder_path):
        return export_feature._show_image_format_dialog(self, folder_path)

    def _do_export_images(self, folder_path, img_format, quality):
        return export_feature._do_export_images(self, folder_path, img_format, quality)

    # ==================== æ–°å¢åŠŸèƒ½ï¼šè‡ªå®šä¹‰æ¶‚æŠ¹æ¨¡å¼ ====================

    def toggle_inpaint_mode(self):
        """åˆ‡æ¢æ¶‚æŠ¹æ¨¡å¼"""
        return inpaint_feature.toggle_inpaint_mode(self)
        if not self.pages or not self.original_image:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        self.inpaint_mode = not self.inpaint_mode

        if self.inpaint_mode:
            # è¿›å…¥æ¶‚æŠ¹æ¨¡å¼
            self.inpaint_mode_btn.config(text="é€€å‡ºæ¶‚æŠ¹", bg="#FF5722")

            # æ˜¾ç¤ºå·¥å…·æ 
            self.inpaint_tools_frame.pack(side=tk.LEFT, after=self.inpaint_mode_btn)
            self.brush_size_frame.pack(side=tk.LEFT, after=self.inpaint_tools_frame)
            self.inpaint_actions_frame.pack(side=tk.LEFT, after=self.brush_size_frame)

            # æ£€æŸ¥æ˜¯å¦æœ‰èƒŒæ™¯å›¾
            page = self.pages[self.current_page_index]
            has_background = page.get("bg_path") and os.path.exists(page.get("bg_path", ""))

            # å†³å®šä½¿ç”¨å“ªä¸ªå›¾ä½œä¸ºåº•å›¾
            if has_background:
                # æœ‰èƒŒæ™¯å›¾ï¼ŒåŸºäºèƒŒæ™¯å›¾è¿›è¡Œè¿­ä»£ä¿®å¤
                base_image = Image.open(page["bg_path"])
                mode_desc = "èƒŒæ™¯å›¾"
            else:
                # æ²¡æœ‰èƒŒæ™¯å›¾ï¼ŒåŸºäºåŸå›¾
                base_image = self.original_image
                mode_desc = "åŸå›¾"

            # åˆå§‹åŒ–è’™ç‰ˆå±‚ï¼ˆä½¿ç”¨å½“å‰åº•å›¾çš„å°ºå¯¸ï¼‰
            if self.inpaint_mask_layer is None or \
               self.inpaint_mask_layer.size != base_image.size:
                self.inpaint_mask_layer = Image.new("L", base_image.size, 0)
                self.inpaint_draw_layer = ImageDraw.Draw(self.inpaint_mask_layer)
                self.inpaint_strokes = []

            # åˆ‡æ¢ç”»å¸ƒå…‰æ ‡
            if self.inpaint_tool == "brush":
                self.canvas.config(cursor="dot")
            else:
                self.canvas.config(cursor="tcross")

            # æ¸…ç©ºæ–‡æœ¬æ¡†é€‰ä¸­çŠ¶æ€ï¼ˆé¿å…å¹²æ‰°ï¼‰
            self.selected_box_index = -1
            self.selected_boxes = []

            self.refresh_canvas()

            # æ ¹æ®æ˜¯å¦æœ‰èƒŒæ™¯å›¾æ˜¾ç¤ºä¸åŒæç¤º
            if has_background:
                self.update_status(f"æ¶‚æŠ¹æ¨¡å¼å·²æ¿€æ´» - åŸºäºèƒŒæ™¯å›¾è¿­ä»£ä¿®å¤")
                messagebox.showinfo("æ¶‚æŠ¹æ¨¡å¼ï¼ˆè¿­ä»£ä¿®å¤ï¼‰",
                    "âœ… æ£€æµ‹åˆ°å·²æœ‰èƒŒæ™¯å›¾ï¼\n\n"
                    "å½“å‰å°†åŸºäºèƒŒæ™¯å›¾è¿›è¡Œè¿­ä»£ä¿®å¤\n\n"
                    "âœï¸ ç¬”åˆ·å·¥å…· - æ¶‚æŠ¹éœ€è¦ä¿®å¤çš„åŒºåŸŸ\n"
                    "â¬œ æ¡†é€‰å·¥å…· - æ‹‰æ¡†æ ‡è®°åŒºåŸŸ\n"
                    "ğŸ¨ ç”ŸæˆèƒŒæ™¯ - ä¿®å¤æ ‡è®°åŒºåŸŸ\n\n"
                    "ğŸ’¡ é€‚ç”¨åœºæ™¯ï¼š\n"
                    "- ä¹‹å‰ç”Ÿæˆçš„èƒŒæ™¯æœ‰é—æ¼\n"
                    "- æ•ˆæœä¸æ»¡æ„éœ€è¦è¡¥å……\n"
                    "- å¤šæ¬¡è¿­ä»£ä¼˜åŒ–èƒŒæ™¯")
            else:
                self.update_status("æ¶‚æŠ¹æ¨¡å¼å·²æ¿€æ´» - æ ‡è®°éœ€è¦å»é™¤çš„åŒºåŸŸ")
                messagebox.showinfo("æ¶‚æŠ¹æ¨¡å¼",
                    "å·²è¿›å…¥æ¶‚æŠ¹æ¨¡å¼ï¼\n\n"
                    "âœï¸ ç¬”åˆ·å·¥å…· - æ¶‚æŠ¹æ ‡è®°åŒºåŸŸ\n"
                    "â¬œ æ¡†é€‰å·¥å…· - æ‹‰æ¡†æ ‡è®°åŒºåŸŸ\n"
                    "ğŸ¨ ç‚¹å‡»ã€Œç”ŸæˆèƒŒæ™¯ã€å¤„ç†æ ‡è®°åŒºåŸŸ\n\n"
                    "æç¤ºï¼šå¯ä»¥ä¸OCRæ£€æµ‹ç»“åˆä½¿ç”¨\n"
                    "å…ˆOCRæ£€æµ‹æ–‡å­—ï¼Œå†æ‰‹åŠ¨è¡¥å……é—æ¼åŒºåŸŸ")
        else:
            # é€€å‡ºæ¶‚æŠ¹æ¨¡å¼
            self.inpaint_mode_btn.config(text="è¿›å…¥æ¶‚æŠ¹", bg="#FF6F00")

            # éšè—å·¥å…·æ 
            self.inpaint_tools_frame.pack_forget()
            self.brush_size_frame.pack_forget()
            self.inpaint_actions_frame.pack_forget()

            # æ¢å¤å…‰æ ‡
            self.canvas.config(cursor="")

            # æ¸…é™¤æ¶‚æŠ¹è§†è§‰
            self.canvas.delete("inpaint_visual")
            self.canvas.delete("inpaint_temp")

            self.update_status("å·²é€€å‡ºæ¶‚æŠ¹æ¨¡å¼")

    def switch_inpaint_tool(self, tool):
        """åˆ‡æ¢æ¶‚æŠ¹å·¥å…·"""
        return inpaint_feature.switch_inpaint_tool(self, tool)
        self.inpaint_tool = tool

        if tool == "brush":
            self.brush_btn.config(relief=tk.SUNKEN, bg="#FFE0B2")
            self.rect_btn.config(relief=tk.RAISED, bg=COLOR_RIBBON_BG)
            self.canvas.config(cursor="dot")
        else:
            self.brush_btn.config(relief=tk.RAISED, bg=COLOR_RIBBON_BG)
            self.rect_btn.config(relief=tk.SUNKEN, bg="#FFE0B2")
            self.canvas.config(cursor="tcross")

    def update_brush_size(self, val):
        """æ›´æ–°ç¬”åˆ·å¤§å°"""
        return inpaint_feature.update_brush_size(self, val)
        self.inpaint_brush_size = int(float(val))

    def handle_inpaint_press(self, x, y):
        """æ¶‚æŠ¹æ¨¡å¼ - æŒ‰ä¸‹äº‹ä»¶"""
        return inpaint_feature.handle_inpaint_press(self, x, y)
        if self.inpaint_tool == "brush":
            # ç¬”åˆ·æ¨¡å¼ - å¼€å§‹æ¶‚æŠ¹
            r = self.inpaint_brush_size // 2
            self.inpaint_draw_layer.ellipse([x-r, y-r, x+r, y+r], fill=255, outline=255)
            self.inpaint_last_pos = (x, y)

            # ç»˜åˆ¶è§†è§‰åé¦ˆ
            self.draw_inpaint_visual_brush(x, y, r)

            # è®°å½•ç¬”ç”»å¼€å§‹
            self.inpaint_strokes.append({"type": "brush", "points": [(x, y)]})

        else:
            # çŸ©å½¢æ¡†é€‰æ¨¡å¼ - è®°å½•èµ·å§‹ç‚¹
            self.inpaint_rect_start = (x, y)

    def handle_inpaint_drag(self, x, y):
        """æ¶‚æŠ¹æ¨¡å¼ - æ‹–æ‹½äº‹ä»¶"""
        return inpaint_feature.handle_inpaint_drag(self, x, y)
        if self.inpaint_tool == "brush":
            # ç¬”åˆ·æ¨¡å¼ - è¿ç»­æ¶‚æŠ¹
            r = self.inpaint_brush_size // 2
            self.inpaint_draw_layer.ellipse([x-r, y-r, x+r, y+r], fill=255, outline=255)

            # è¿çº¿ï¼ˆå¹³æ»‘ï¼‰
            if self.inpaint_last_pos:
                self.inpaint_draw_layer.line([self.inpaint_last_pos, (x, y)],
                                            fill=255, width=self.inpaint_brush_size)

            self.inpaint_last_pos = (x, y)

            # ç»˜åˆ¶è§†è§‰åé¦ˆ
            self.draw_inpaint_visual_brush(x, y, r)

            # è®°å½•ç¬”ç”»ç‚¹
            if self.inpaint_strokes and self.inpaint_strokes[-1]["type"] == "brush":
                self.inpaint_strokes[-1]["points"].append((x, y))

        else:
            # çŸ©å½¢æ¡†é€‰æ¨¡å¼ - ç»˜åˆ¶ä¸´æ—¶çŸ©å½¢
            if self.inpaint_rect_start:
                self.draw_inpaint_temp_rect(x, y)

    def handle_inpaint_release(self, x, y):
        """æ¶‚æŠ¹æ¨¡å¼ - é‡Šæ”¾äº‹ä»¶"""
        return inpaint_feature.handle_inpaint_release(self, x, y)
        if self.inpaint_tool == "brush":
            # ç¬”åˆ·æ¨¡å¼ - ç»“æŸç¬”ç”»
            self.inpaint_last_pos = None

            # ä¿å­˜å½“å‰ç¬”ç”»åˆ°å†å²ï¼ˆç¬”åˆ·å®Œæˆæ—¶ä¿å­˜ï¼‰
            if self.inpaint_strokes and self.inpaint_strokes[-1]["type"] == "brush":
                self.save_state("inpaint_stroke", {
                    "stroke": self.inpaint_strokes[-1],
                    "mask_state": self.inpaint_strokes[:-1]  # ä¹‹å‰çš„çŠ¶æ€
                })

        else:
            # çŸ©å½¢æ¡†é€‰æ¨¡å¼ - å®Œæˆæ¡†é€‰
            if self.inpaint_rect_start:
                sx, sy = self.inpaint_rect_start
                x1, y1 = min(sx, x), min(sy, y)
                x2, y2 = max(sx, x), max(sy, y)

                # å†™å…¥è’™ç‰ˆ
                self.inpaint_draw_layer.rectangle([x1, y1, x2, y2], fill=255, outline=255)

                # ç»˜åˆ¶æ°¸ä¹…è§†è§‰
                self.draw_inpaint_visual_rect(x1, y1, x2, y2)

                # æ¸…é™¤ä¸´æ—¶çŸ©å½¢
                self.canvas.delete("inpaint_temp")

                # è®°å½•çŸ©å½¢
                rect_stroke = {
                    "type": "rect",
                    "coords": (x1, y1, x2, y2)
                }
                self.inpaint_strokes.append(rect_stroke)

                # ä¿å­˜åˆ°å†å²
                self.save_state("inpaint_stroke", {
                    "stroke": rect_stroke,
                    "mask_state": self.inpaint_strokes[:-1]
                })

                self.inpaint_rect_start = None

    def draw_inpaint_visual_brush(self, x, y, radius):
        """ç»˜åˆ¶ç¬”åˆ·æ¶‚æŠ¹çš„è§†è§‰åé¦ˆ"""
        return inpaint_feature.draw_inpaint_visual_brush(self, x, y, radius)
        # è½¬æ¢ä¸ºç”»å¸ƒåæ ‡
        canvas_x = x * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_y = y * self.scale + getattr(self, 'canvas_offset_y', 0)
        canvas_r = radius * self.scale

        # åŠé€æ˜çº¢è‰²åœ†å½¢
        self.canvas.create_oval(
            canvas_x - canvas_r, canvas_y - canvas_r,
            canvas_x + canvas_r, canvas_y + canvas_r,
            fill="#ff0000", stipple="gray50", outline="",
            tags="inpaint_visual"
        )

    def draw_inpaint_temp_rect(self, x, y):
        """ç»˜åˆ¶ä¸´æ—¶çŸ©å½¢æ¡†é€‰"""
        return inpaint_feature.draw_inpaint_temp_rect(self, x, y)
        if not self.inpaint_rect_start:
            return

        sx, sy = self.inpaint_rect_start

        # è½¬æ¢ä¸ºç”»å¸ƒåæ ‡
        canvas_sx = sx * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_sy = sy * self.scale + getattr(self, 'canvas_offset_y', 0)
        canvas_x = x * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_y = y * self.scale + getattr(self, 'canvas_offset_y', 0)

        # åˆ é™¤æ—§çš„ä¸´æ—¶çŸ©å½¢
        self.canvas.delete("inpaint_temp")

        # ç»˜åˆ¶æ–°çš„ä¸´æ—¶çŸ©å½¢
        self.canvas.create_rectangle(
            canvas_sx, canvas_sy, canvas_x, canvas_y,
            outline="red", width=2, tags="inpaint_temp"
        )

    def draw_inpaint_visual_rect(self, x1, y1, x2, y2):
        """ç»˜åˆ¶çŸ©å½¢æ¡†é€‰çš„æ°¸ä¹…è§†è§‰åé¦ˆ"""
        return inpaint_feature.draw_inpaint_visual_rect(self, x1, y1, x2, y2)
        # è½¬æ¢ä¸ºç”»å¸ƒåæ ‡
        canvas_x1 = x1 * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_y1 = y1 * self.scale + getattr(self, 'canvas_offset_y', 0)
        canvas_x2 = x2 * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_y2 = y2 * self.scale + getattr(self, 'canvas_offset_y', 0)

        # åŠé€æ˜çº¢è‰²çŸ©å½¢
        self.canvas.create_rectangle(
            canvas_x1, canvas_y1, canvas_x2, canvas_y2,
            fill="#ff0000", stipple="gray25", outline="red",
            tags="inpaint_visual"
        )

    def clear_inpaint_mask(self):
        """æ¸…ç©ºæ‰€æœ‰æ¶‚æŠ¹"""
        return inpaint_feature.clear_inpaint_mask(self)
        if not self.inpaint_strokes:
            messagebox.showinfo("æç¤º", "å½“å‰æ²¡æœ‰æ¶‚æŠ¹å†…å®¹")
            return

        result = messagebox.askyesno("ç¡®è®¤", "ç¡®å®šè¦æ¸…ç©ºæ‰€æœ‰æ¶‚æŠ¹å—ï¼Ÿ")
        if not result:
            return

        # æ¸…ç©ºè’™ç‰ˆ
        self.inpaint_mask_layer = Image.new("L", self.original_image.size, 0)
        self.inpaint_draw_layer = ImageDraw.Draw(self.inpaint_mask_layer)
        self.inpaint_strokes = []

        # æ¸…é™¤è§†è§‰
        self.canvas.delete("inpaint_visual")
        self.canvas.delete("inpaint_temp")

        self.update_status("å·²æ¸…ç©ºæ‰€æœ‰æ¶‚æŠ¹")

    def rebuild_inpaint_mask(self):
        """é‡å»ºæ¶‚æŠ¹è’™ç‰ˆï¼ˆç”¨äºæ’¤é”€åï¼‰"""
        return inpaint_feature.rebuild_inpaint_mask(self)
        # é‡ç½®è’™ç‰ˆ
        self.inpaint_mask_layer = Image.new("L", self.original_image.size, 0)
        self.inpaint_draw_layer = ImageDraw.Draw(self.inpaint_mask_layer)

        # æ¸…é™¤è§†è§‰
        self.canvas.delete("inpaint_visual")

        # é‡æ–°ç»˜åˆ¶æ‰€æœ‰ç¬”ç”»
        for stroke in self.inpaint_strokes:
            if stroke["type"] == "brush":
                points = stroke["points"]
                r = self.inpaint_brush_size // 2

                for i, (x, y) in enumerate(points):
                    self.inpaint_draw_layer.ellipse([x-r, y-r, x+r, y+r], fill=255, outline=255)
                    if i > 0:
                        prev_x, prev_y = points[i-1]
                        self.inpaint_draw_layer.line([(prev_x, prev_y), (x, y)],
                                                    fill=255, width=self.inpaint_brush_size)
                    # ç»˜åˆ¶è§†è§‰
                    self.draw_inpaint_visual_brush(x, y, r)

            elif stroke["type"] == "rect":
                x1, y1, x2, y2 = stroke["coords"]
                self.inpaint_draw_layer.rectangle([x1, y1, x2, y2], fill=255, outline=255)
                self.draw_inpaint_visual_rect(x1, y1, x2, y2)

    def generate_bg_from_custom_mask(self):
        """åŸºäºè‡ªå®šä¹‰æ¶‚æŠ¹è’™ç‰ˆç”ŸæˆèƒŒæ™¯"""
        return inpaint_feature.generate_bg_from_custom_mask(self)
        if not self.pages or not self.original_image:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        if not self.inpaint_mask_layer:
            messagebox.showwarning("æç¤º", "è¯·å…ˆæ¶‚æŠ¹æ ‡è®°éœ€è¦å»é™¤çš„åŒºåŸŸ")
            return

        # æ£€æŸ¥æ˜¯å¦æœ‰æ¶‚æŠ¹å†…å®¹
        if not self.inpaint_mask_layer.getbbox():
            messagebox.showwarning("æç¤º", "å½“å‰æ²¡æœ‰æ¶‚æŠ¹å†…å®¹\n\nè¯·ä½¿ç”¨ç¬”åˆ·æˆ–æ¡†é€‰å·¥å…·æ ‡è®°éœ€è¦å»é™¤çš„åŒºåŸŸ")
            return

        if not self.config.get("inpaint_enabled", True):
            messagebox.showwarning("æç¤º", "èƒŒæ™¯ç”ŸæˆåŠŸèƒ½å·²ç¦ç”¨\n\nè¯·åœ¨è®¾ç½®ä¸­å¯ç”¨")
            return

        # æ£€æŸ¥æ˜¯å¦æœ‰èƒŒæ™¯å›¾ï¼ˆå†³å®šä½¿ç”¨å“ªä¸ªå›¾ä½œä¸ºåº•å›¾ï¼‰
        page = self.pages[self.current_page_index]
        has_background = page.get("bg_path") and os.path.exists(page.get("bg_path", ""))

        if has_background:
            base_image = Image.open(page["bg_path"])
            mode_desc = "èƒŒæ™¯å›¾ï¼ˆè¿­ä»£ä¿®å¤ï¼‰"
        else:
            base_image = page["image"]
            mode_desc = "åŸå›¾"

        # ç¡®è®¤å¯¹è¯æ¡†
        result = messagebox.askyesno("ç¡®è®¤",
            f"å³å°†åŸºäº{mode_desc}ç”Ÿæˆæ–°èƒŒæ™¯å›¾\n\n"
            f"åº•å›¾ï¼š{mode_desc}\n"
            f"æ¶‚æŠ¹åŒºåŸŸï¼šå°†è¢«AIæ™ºèƒ½å¡«å……\n\n"
            f"æ­¤æ“ä½œéœ€è¦è°ƒç”¨IOPaint APIæœåŠ¡\n"
            f"å¤„ç†æ—¶é—´çº¦ 5-30 ç§’\n\n"
            f"æ˜¯å¦ç»§ç»­ï¼Ÿ")

        if not result:
            return

        # ä¿å­˜èƒŒæ™¯ç”Ÿæˆå‰çš„çŠ¶æ€åˆ°å†å²ï¼ˆé‡è¦ï¼ï¼‰
        old_bg_path = page.get("bg_path")
        self.save_state("background", {
            "old_bg_path": old_bg_path,
            "new_bg_path": None  # å°†åœ¨ç”Ÿæˆåå¡«å……
        })

        self.update_status(f"æ­£åœ¨ç”ŸæˆèƒŒæ™¯å›¾ï¼ˆåŸºäº{mode_desc}ï¼‰...")

        def generate_bg():
            try:
                # è°ƒç”¨APIä¿®å¤ï¼ˆä½¿ç”¨åº•å›¾è€Œä¸æ˜¯åŸå›¾ï¼‰
                self.root.after(0, lambda: self.update_status(f"æ­£åœ¨è°ƒç”¨IOPaint APIä¿®å¤ï¼ˆ{mode_desc}ï¼‰..."))
                result_img = self.call_inpaint_api(base_image, self.inpaint_mask_layer)

                if result_img:
                    # ä¿å­˜åˆ°ä¸´æ—¶æ–‡ä»¶
                    temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
                    os.makedirs(temp_dir, exist_ok=True)

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    bg_path = os.path.join(temp_dir, f"bg_iter_page{self.current_page_index+1}_{timestamp}.png")
                    result_img.save(bg_path, quality=95)

                    # æ›´æ–°å†å²è®°å½•ä¸­çš„new_bg_path
                    if self.history and self.history[-1]["type"] == "background":
                        self.history[-1]["data"]["new_bg_path"] = bg_path

                    # æ›´æ–°ä¸ºæ–°çš„èƒŒæ™¯å›¾
                    page["bg_path"] = bg_path
                    self.root.after(0, lambda: setattr(self, 'clean_bg_path', bg_path))

                    # æ¸…ç©ºæ¶‚æŠ¹ï¼ˆå·²ç”ŸæˆèƒŒæ™¯ï¼‰
                    self.root.after(0, lambda: self.clear_inpaint_mask())

                    # åˆ·æ–°ç•Œé¢
                    self.root.after(0, self.update_bg_status)
                    self.root.after(0, self.update_thumbnails)
                    self.root.after(0, self.refresh_canvas)

                    self.root.after(0, lambda: self.update_status(f"ç¬¬ {self.current_page_index+1} é¡µèƒŒæ™¯ç”ŸæˆæˆåŠŸï¼"))

                    if has_background:
                        msg = (f"è¿­ä»£ä¿®å¤æˆåŠŸï¼\n\n"
                               f"âœ… å·²åŸºäºä¹‹å‰çš„èƒŒæ™¯å›¾è¿›è¡Œä¿®å¤\n"
                               f"âœ… æ¶‚æŠ¹åŒºåŸŸå·²è¢«æ™ºèƒ½å¡«å……\n"
                               f"âœ… æ–°èƒŒæ™¯å·²è‡ªåŠ¨è®¾ç½®åˆ°å½“å‰é¡µ\n\n"
                               f"ğŸ’¡ å¦‚éœ€ç»§ç»­ä¿®å¤ï¼Œå¯å†æ¬¡è¿›å…¥æ¶‚æŠ¹æ¨¡å¼\n"
                               f"ğŸ’¡ æŒ‰Ctrl+Zå¯ä»¥æ’¤é”€æ­¤æ¬¡ç”Ÿæˆ\n\n"
                               f"ä¿å­˜ä½ç½®ï¼š\n{bg_path}")
                    else:
                        msg = (f"èƒŒæ™¯å›¾ç”ŸæˆæˆåŠŸï¼\n\n"
                               f"å·²æ ¹æ®æ¶‚æŠ¹åŒºåŸŸå»é™¤å†…å®¹\n"
                               f"èƒŒæ™¯å·²è‡ªåŠ¨è®¾ç½®åˆ°å½“å‰é¡µ\n\n"
                               f"ğŸ’¡ æŒ‰Ctrl+Zå¯ä»¥æ’¤é”€æ­¤æ¬¡ç”Ÿæˆ\n\n"
                               f"ä¿å­˜ä½ç½®ï¼š\n{bg_path}")

                    self.root.after(0, lambda: messagebox.showinfo("æˆåŠŸ", msg))
                else:
                    self.root.after(0, lambda: self.update_status("èƒŒæ™¯ç”Ÿæˆå¤±è´¥"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"èƒŒæ™¯ç”Ÿæˆå¤±è´¥:\n{error_msg}")
                err_text = str(e)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"èƒŒæ™¯ç”Ÿæˆå¤±è´¥:\n{t}"))
                self.root.after(0, lambda: self.update_status("èƒŒæ™¯ç”Ÿæˆå¤±è´¥"))

        threading.Thread(target=generate_bg, daemon=True).start()

    # ==================== æ–°å¢åŠŸèƒ½ï¼šIOPaint API èƒŒæ™¯ç”Ÿæˆ ====================

    def call_inpaint_api(self, image_pil, mask_pil, crop_padding=128):
        """
        è°ƒç”¨IOPaint APIè¿›è¡Œå›¾åƒä¿®å¤

        Args:
            image_pil: PIL Imageï¼ŒåŸå›¾
            mask_pil: PIL Image (Læ¨¡å¼)ï¼Œè’™ç‰ˆï¼ˆç™½è‰²=éœ€è¦ä¿®å¤çš„åŒºåŸŸï¼‰
            crop_padding: è£åˆ‡paddingå¤§å°

        Returns:
            PIL Image æˆ– None
        """
        return inpaint_feature.call_inpaint_api(self, image_pil, mask_pil, crop_padding=crop_padding)
        try:
            api_url = self.config.get("inpaint_api_url", "http://127.0.0.1:8080/api/v1/inpaint")

            # === æ™ºèƒ½è£åˆ‡é€»è¾‘ï¼ˆåªå¤„ç†æœ‰è’™ç‰ˆçš„åŒºåŸŸï¼‰===
            mask_np = np.array(mask_pil)
            rows = np.any(mask_np, axis=1)
            cols = np.any(mask_np, axis=0)

            if not rows.any() or not cols.any():
                # æ²¡æœ‰è’™ç‰ˆåŒºåŸŸ
                return image_pil.copy()

            y_min, y_max = np.where(rows)[0][[0, -1]]
            x_min, x_max = np.where(cols)[0][[0, -1]]

            W, H = image_pil.size
            pad = crop_padding
            x1 = max(0, x_min - pad)
            y1 = max(0, y_min - pad)
            x2 = min(W, x_max + pad)
            y2 = min(H, y_max + pad)

            crop_box = (x1, y1, x2, y2)
            crop_img = image_pil.crop(crop_box)
            crop_mask = mask_pil.crop(crop_box)

            # === Base64ç¼–ç  ===
            def to_b64(img):
                buffer = BytesIO()
                img.save(buffer, "PNG")
                return base64.b64encode(buffer.getvalue()).decode()

            payload = {
                "image": to_b64(crop_img),
                "mask": to_b64(crop_mask),
                "ldm_steps": 30,
                "hd_strategy": "Original",
                "sd_sampler": "UniPC"
            }

            # === è°ƒç”¨API ===
            response = requests.post(api_url, json=payload, timeout=120)

            if response.status_code == 200:
                # ä¿®å¤æˆåŠŸï¼Œåˆæˆå›åŸå›¾
                res_crop = Image.open(BytesIO(response.content))

                # åˆ›å»ºç»“æœå›¾
                final = image_pil.copy()

                # ä½¿ç”¨é«˜æ–¯æ¨¡ç³Šå¹³æ»‘è¾¹ç¼˜
                blur_mask = crop_mask.filter(ImageFilter.GaussianBlur(3))
                orig_crop_area = final.crop(crop_box)
                blended = Image.composite(res_crop, orig_crop_area, blur_mask)
                final.paste(blended, (x1, y1))

                return final
            else:
                self.root.after(0, lambda: messagebox.showerror("APIé”™è¯¯",
                    f"IOPaint APIè¿”å›é”™è¯¯: {response.status_code}\n{response.text[:200]}"))
                return None

        except requests.exceptions.ConnectionError:
            self.root.after(0, lambda: messagebox.showerror("è¿æ¥é”™è¯¯",
                "æ— æ³•è¿æ¥åˆ°IOPaint APIæœåŠ¡ï¼\n\n"
                "è¯·ç¡®ä¿IOPaintæœåŠ¡æ­£åœ¨è¿è¡Œï¼š\n"
                f"APIåœ°å€ï¼š{api_url}\n\n"
                "å¯åŠ¨å‘½ä»¤ï¼š\n"
                "iopaint start --host 127.0.0.1 --port 8080"))
            return None
        except Exception as e:
            import traceback
            error_msg = traceback.format_exc()
            print(f"IOPaint APIè°ƒç”¨å¤±è´¥:\n{error_msg}")
            err_text = str(e)
            self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"ä¿®å¤å¤±è´¥:\n{t}"))
            return None

    def create_mask_from_boxes(self, image_size, text_boxes, padding=5):
        """
        æ ¹æ®æ–‡æœ¬æ¡†ä½ç½®åˆ›å»ºè’™ç‰ˆ

        Args:
            image_size: (width, height) å›¾ç‰‡å°ºå¯¸
            text_boxes: æ–‡æœ¬æ¡†åˆ—è¡¨
            padding: æ–‡æœ¬æ¡†æ‰©å±•è¾¹è·

        Returns:
            PIL Image (Læ¨¡å¼)ï¼Œç™½è‰²=éœ€è¦ä¿®å¤çš„åŒºåŸŸ
        """
        return inpaint_feature.create_mask_from_boxes(self, image_size, text_boxes, padding=padding)
        mask = Image.new("L", image_size, 0)  # å…¨é»‘èƒŒæ™¯
        draw = ImageDraw.Draw(mask)

        img_w, img_h = image_size

        for box in text_boxes:
            # ç¨å¾®æ‰©å¤§æ–‡æœ¬æ¡†åŒºåŸŸ
            x1 = max(0, box.x - padding)
            y1 = max(0, box.y - padding)
            x2 = min(img_w, box.x + box.width + padding)
            y2 = min(img_h, box.y + box.height + padding)

            # æ ‡è®°ä¸ºç™½è‰²ï¼ˆéœ€è¦ä¿®å¤ï¼‰
            draw.rectangle([x1, y1, x2, y2], fill=255)

        return mask

    def auto_generate_background_current(self):
        """ä¸ºå½“å‰é¡µè‡ªåŠ¨ç”Ÿæˆä¿®å¤å›¾å±‚ï¼ˆæ ¹æ®æ–‡æœ¬æ¡†ä½ç½®ï¼›ä¸æ›¿æ¢åŸå›¾/èƒŒæ™¯ï¼‰"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        if not self.text_boxes:
            messagebox.showwarning("æç¤º", "å½“å‰é¡µæ²¡æœ‰æ–‡æœ¬æ¡†\n\nè¯·å…ˆä½¿ç”¨ã€Œæ£€æµ‹ã€åŠŸèƒ½è¯†åˆ«æ–‡æœ¬åŒºåŸŸ")
            return

        if not self.config.get("inpaint_enabled", True):
            messagebox.showwarning("æç¤º", "èƒŒæ™¯ç”ŸæˆåŠŸèƒ½å·²ç¦ç”¨\n\nè¯·åœ¨è®¾ç½®ä¸­å¯ç”¨")
            return

        # ç¡®è®¤å¯¹è¯æ¡†
        result = messagebox.askyesno(
            "ç¡®è®¤",
            f"å³å°†ä¸ºç¬¬ {self.current_page_index + 1} é¡µç”Ÿæˆä¿®å¤å›¾å±‚\n\n"
            f"å½“å‰é¡µæœ‰ {len(self.text_boxes)} ä¸ªæ–‡æœ¬æ¡†\n"
            "ç³»ç»Ÿå°†è‡ªåŠ¨å¯¹è¿™äº›æ–‡å­—åŒºåŸŸè¿›è¡Œä¿®å¤\n\n"
            "æç¤ºï¼šç»“æœä¼šä½œä¸ºå›¾å±‚å åŠ ï¼Œä¸ä¼šç›´æ¥æ›¿æ¢åŸå›¾/èƒŒæ™¯\n\n"
            "æ­¤æ“ä½œéœ€è¦è°ƒç”¨ IOPaint API æœåŠ¡\n"
            "å¤„ç†æ—¶é—´çº¦ 5-30 ç§’\n\n"
            "æ˜¯å¦ç»§ç»­ï¼Ÿ",
        )

        if not result:
            return

        page = self.pages[self.current_page_index]
        # éç ´åï¼šä¿å­˜â€œå›¾å±‚å¿«ç…§â€ä¾¿äº Ctrl+Z æ’¤é”€
        self.save_state("layers")

        self.update_status("æ­£åœ¨ç”Ÿæˆä¿®å¤å›¾å±‚...")

        def generate_bg():
            try:
                # è·å–å½“å‰é¡µæ•°æ®
                img = page["image"]  # ç¼–è¾‘ç”¨çš„å›¾ç‰‡

                # åˆ›å»ºè’™ç‰ˆ
                self.root.after(0, lambda: self.update_status("æ­£åœ¨åˆ›å»ºè’™ç‰ˆ..."))
                mask = self.create_mask_from_boxes(img.size, self.text_boxes, padding=5)

                # è°ƒç”¨APIä¿®å¤
                self.root.after(0, lambda: self.update_status("æ­£åœ¨è°ƒç”¨IOPaint APIä¿®å¤..."))
                result_img = self.call_inpaint_api(img, mask)

                if result_img:
                    overlay = result_img.convert("RGBA")
                    alpha = mask.convert("L").filter(ImageFilter.GaussianBlur(3))
                    overlay.putalpha(alpha)

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    layer_name = f"IOPaintå»å­—_{timestamp}"
                    layer = self.add_image_layer(page, overlay, name=layer_name, x=0, y=0, opacity=1.0, visible=True)

                    self.root.after(0, self.update_thumbnails)
                    self.root.after(0, self.update_layer_listbox)
                    self.root.after(0, self.scroll_to_layers)
                    if layer and layer.get("id"):
                        self.root.after(0, lambda lid=layer["id"]: self.select_layer_by_id(lid))
                    self.root.after(0, self.refresh_canvas)
                    self.root.after(0, self.mark_unsaved)

                    self.root.after(0, lambda: self.update_status(f"å·²ç”Ÿæˆä¿®å¤å›¾å±‚ï¼š{layer_name}"))
                    self.root.after(
                        0,
                        lambda: messagebox.showinfo(
                            "å®Œæˆ",
                            "å¤„ç†å®Œæˆï¼\n\n"
                            f"å·²å»é™¤ {len(self.text_boxes)} ä¸ªæ–‡å­—åŒºåŸŸ\n"
                            "ç»“æœå·²ä½œä¸ºå›¾å±‚å åŠ ï¼ˆå³ä¾§å±æ€§é¢æ¿æ»šåŠ¨åˆ°åº•éƒ¨â€œå›¾å±‚â€å³å¯çœ‹åˆ°ï¼‰\n\n"
                            "æç¤ºï¼šCtrl+Z å¯ä»¥æ’¤é”€",
                        ),
                    )
                else:
                    self.root.after(0, lambda: self.update_status("ä¿®å¤å¤±è´¥"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"ä¿®å¤å¤±è´¥:\n{error_msg}")
                err_text = str(e)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"ä¿®å¤å¤±è´¥:\n{t}"))
                self.root.after(0, lambda: self.update_status("ä¿®å¤å¤±è´¥"))

        threading.Thread(target=generate_bg, daemon=True).start()

    def auto_generate_background_all(self):
        """æ‰¹é‡ä¸ºæ‰€æœ‰é¡µç”Ÿæˆä¿®å¤å›¾å±‚ï¼ˆä¸æ›¿æ¢åŸå›¾/èƒŒæ™¯ï¼‰"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        if not self.config.get("inpaint_enabled", True):
            messagebox.showwarning("æç¤º", "èƒŒæ™¯ç”ŸæˆåŠŸèƒ½å·²ç¦ç”¨\n\nè¯·åœ¨è®¾ç½®ä¸­å¯ç”¨")
            return

        # ç»Ÿè®¡æœ‰æ–‡æœ¬æ¡†çš„é¡µé¢
        pages_with_boxes = sum(1 for p in self.pages if p.get("text_boxes"))

        if pages_with_boxes == 0:
            messagebox.showwarning("æç¤º", "æ‰€æœ‰é¡µé¢éƒ½æ²¡æœ‰æ–‡æœ¬æ¡†\n\nè¯·å…ˆä½¿ç”¨ã€Œæ£€æµ‹ - å…¨éƒ¨é¡µã€åŠŸèƒ½")
            return

        # ç¡®è®¤å¯¹è¯æ¡†
        result = messagebox.askyesno(
            "æ‰¹é‡ä¿®å¤ï¼ˆIOPaintï¼‰",
            f"å³å°†ä¸º {pages_with_boxes}/{len(self.pages)} é¡µç”Ÿæˆä¿®å¤å›¾å±‚\n\n"
            "æç¤ºï¼šç»“æœä¼šä½œä¸ºå›¾å±‚å åŠ ï¼Œä¸ä¼šç›´æ¥æ›¿æ¢åŸå›¾/èƒŒæ™¯\n\n"
            "æ­¤æ“ä½œå¯èƒ½éœ€è¦è¾ƒé•¿æ—¶é—´\n"
            f"é¢„è®¡æ—¶é—´ï¼š{pages_with_boxes * 10} - {pages_with_boxes * 30} ç§’\n\n"
            "å¤„ç†æœŸé—´å¯ä»¥ç»§ç»­ç¼–è¾‘ï¼Œä½†è¯·å‹¿å…³é—­ç¨‹åº\n\n"
            "æ˜¯å¦ç»§ç»­ï¼Ÿ",
        )

        if not result:
            return

        # æ‰¹é‡æ“ä½œï¼šä¿å­˜â€œå…¨é¡µå›¾å±‚å¿«ç…§â€ä»¥ä¾¿ Ctrl+Z ä¸€æ¬¡æ’¤é”€æ•´ä¸ªæ‰¹é‡ç»“æœ
        self.save_state("pages_layers")

        self.save_current_page()
        self.update_status("å¼€å§‹æ‰¹é‡ç”Ÿæˆä¿®å¤å›¾å±‚...")

        def generate_all_bg():
            try:
                success_count = 0
                fail_count = 0

                for page_idx, page in enumerate(self.pages):
                    text_boxes = page.get("text_boxes", [])

                    if not text_boxes:
                        continue

                    self.root.after(0, lambda idx=page_idx+1, total=len(self.pages):
                        self.update_status(f"æ­£åœ¨å¤„ç†ç¬¬ {idx}/{total} é¡µ..."))

                    try:
                        # è·å–å›¾ç‰‡ï¼ˆéœ€è¦ä»dictè½¬ä¸ºTextBoxå¯¹è±¡ï¼‰
                        img = page["image"]
                        boxes = [TextBox.from_dict(b) if isinstance(b, dict) else b for b in text_boxes]

                        # åˆ›å»ºè’™ç‰ˆ
                        mask = self.create_mask_from_boxes(img.size, boxes, padding=5)

                        # è°ƒç”¨API
                        result_img = self.call_inpaint_api(img, mask)

                        if result_img:
                            overlay = result_img.convert("RGBA")
                            alpha = mask.convert("L").filter(ImageFilter.GaussianBlur(3))
                            overlay.putalpha(alpha)

                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                            layer_name = f"IOPaintæ‰¹é‡ä¿®å¤_p{page_idx+1}_{timestamp}"
                            self.add_image_layer(page, overlay, name=layer_name, x=0, y=0, opacity=1.0, visible=True)
                            success_count += 1
                        else:
                            fail_count += 1

                    except Exception as e:
                        print(f"ç¬¬ {page_idx+1} é¡µå¤„ç†å¤±è´¥: {e}")
                        fail_count += 1
                        continue

                # åˆ·æ–°ç•Œé¢
                self.root.after(0, self.load_current_page)
                self.root.after(0, self.update_thumbnails)
                self.root.after(0, self.update_layer_listbox)

                # æ˜¾ç¤ºç»“æœ
                self.root.after(0, lambda: self.update_status(
                    f"æ‰¹é‡å¤„ç†å®Œæˆï¼æˆåŠŸ {success_count} é¡µï¼Œå¤±è´¥ {fail_count} é¡µ"))

                self.root.after(0, lambda: messagebox.showinfo("å®Œæˆ",
                    f"æ‰¹é‡ä¿®å¤å®Œæˆï¼\n\n"
                    f"æˆåŠŸï¼š{success_count} é¡µ\n"
                    f"å¤±è´¥ï¼š{fail_count} é¡µ\n\n"
                    f"ç»“æœå·²ä½œä¸ºå›¾å±‚å åŠ ï¼ˆå³ä¾§â€œå›¾å±‚â€å¯éšè—/åˆ é™¤/è°ƒé¡ºåºï¼‰\n"
                    f"æç¤ºï¼šCtrl+Z å¯æ’¤é”€æ•´ä¸ªæ‰¹é‡ç»“æœ"))

            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"æ‰¹é‡ç”Ÿæˆå¤±è´¥:\n{error_msg}")
                err_text = str(e)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"æ‰¹é‡ç”Ÿæˆå¤±è´¥:\n{t}"))
                self.root.after(0, lambda: self.update_status("æ‰¹é‡ç”Ÿæˆå¤±è´¥"))

        threading.Thread(target=generate_all_bg, daemon=True).start()

    # ==================== AIå›¾ç‰‡æ›¿æ¢åŠŸèƒ½ ====================

    def toggle_ai_replace_mode(self):
        """åˆ‡æ¢AIæ›¿æ¢æ¨¡å¼"""
        return ai_replace_feature.toggle_ai_replace_mode(self)
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        self.save_current_page()
        self.ai_replace_mode = not self.ai_replace_mode

        if self.ai_replace_mode:
            self.ai_replace_mode_btn.config(text="é€€å‡ºAIæ›¿æ¢", bg="#F50057")
            if self.inpaint_mode:
                self.toggle_inpaint_mode()
            self.ai_replace_selection = None
            if self.ai_replace_rect_id:
                self.canvas.delete(self.ai_replace_rect_id)
                self.ai_replace_rect_id = None
            self.update_status("AIæ›¿æ¢æ¨¡å¼å·²æ¿€æ´» - æ¡†é€‰è¦æ›¿æ¢çš„åŒºåŸŸ")
            messagebox.showinfo("AIæ›¿æ¢æ¨¡å¼",
                "å·²è¿›å…¥AIæ›¿æ¢æ¨¡å¼ï¼\n\n"
                "ğŸ“ æ“ä½œæ­¥éª¤ï¼š\n"
                "1. ç”¨é¼ æ ‡æ¡†é€‰è¦æ›¿æ¢/ç¼–è¾‘çš„åŒºåŸŸ\n"
                "2. è¾“å…¥æç¤ºè¯æè¿°æƒ³è¦çš„æ•ˆæœ\n"
                "3. ç­‰å¾…AIç”Ÿæˆå¹¶è‡ªåŠ¨èåˆ\n\n"
                "ğŸ’¡ æç¤ºï¼š\n"
                "- å¯ä»¥åœ¨åŸå›¾æˆ–èƒŒæ™¯å›¾ä¸Šæ¡†é€‰\n"
                "- æ”¯æŒå¤šæ¬¡ç¼–è¾‘å’Œè¿­ä»£")
        else:
            self.ai_replace_mode_btn.config(text="AIæ›¿æ¢", bg="#E91E63")
            if self.ai_replace_rect_id:
                self.canvas.delete(self.ai_replace_rect_id)
                self.ai_replace_rect_id = None
            self.ai_replace_selection = None
            self.update_status("å·²é€€å‡ºAIæ›¿æ¢æ¨¡å¼")

    def handle_ai_replace_press(self, x, y):
        """AIæ›¿æ¢æ¨¡å¼ - æŒ‰ä¸‹äº‹ä»¶"""
        return ai_replace_feature.handle_ai_replace_press(self, x, y)
        self.ai_replace_rect_start = (x, y)

    def handle_ai_replace_drag(self, canvas_x, canvas_y):
        """AIæ›¿æ¢æ¨¡å¼ - æ‹–æ‹½äº‹ä»¶"""
        return ai_replace_feature.handle_ai_replace_drag(self, canvas_x, canvas_y)
        if not self.ai_replace_rect_start:
            return
        if self.ai_replace_rect_id:
            self.canvas.delete(self.ai_replace_rect_id)

        img_x, img_y = self.ai_replace_rect_start
        canvas_x1 = img_x * self.scale + getattr(self, 'canvas_offset_x', 0)
        canvas_y1 = img_y * self.scale + getattr(self, 'canvas_offset_y', 0)

        self.ai_replace_rect_id = self.canvas.create_rectangle(
            canvas_x1, canvas_y1, canvas_x, canvas_y,
            outline="#E91E63", width=3, dash=(5, 5))

    def handle_ai_replace_release(self, canvas_x, canvas_y):
        """AIæ›¿æ¢æ¨¡å¼ - é‡Šæ”¾äº‹ä»¶"""
        return ai_replace_feature.handle_ai_replace_release(self, canvas_x, canvas_y)
        if not self.ai_replace_rect_start:
            return

        img_x = (canvas_x - getattr(self, 'canvas_offset_x', 0)) / self.scale
        img_y = (canvas_y - getattr(self, 'canvas_offset_y', 0)) / self.scale

        x1, y1 = self.ai_replace_rect_start
        x1, x2 = min(x1, img_x), max(x1, img_x)
        y1, y2 = min(y1, img_y), max(y1, img_y)

        if abs(x2 - x1) < 10 or abs(y2 - y1) < 10:
            messagebox.showwarning("æç¤º", "é€‰æ¡†å¤ªå°ï¼Œè¯·é‡æ–°æ¡†é€‰")
            if self.ai_replace_rect_id:
                self.canvas.delete(self.ai_replace_rect_id)
                self.ai_replace_rect_id = None
            self.ai_replace_rect_start = None
            return

        self.ai_replace_selection = (int(x1), int(y1), int(x2), int(y2))
        self.ai_replace_rect_start = None
        self.show_ai_replace_dialog()

    def show_ai_replace_dialog(self):
        """æ˜¾ç¤ºAIæ›¿æ¢æ“ä½œå¯¹è¯æ¡†"""
        return ai_replace_feature.show_ai_replace_dialog(self)
        if not self.ai_replace_selection:
            return

        x1, y1, x2, y2 = self.ai_replace_selection

        dialog = tk.Toplevel(self.root)
        dialog.title("AIå›¾ç‰‡æ›¿æ¢/ç”Ÿæˆ")
        dialog.geometry("500x350")
        dialog.transient(self.root)
        dialog.grab_set()

        # æ ‡é¢˜
        title_frame = tk.Frame(dialog, bg="#E91E63", height=50)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)

        tk.Label(title_frame, text="AI å›¾ç‰‡æ›¿æ¢/ç”Ÿæˆ",
                bg="#E91E63", fg="white",
                font=(FONT_FAMILY, 14, "bold")).pack(pady=10)

        # å†…å®¹åŒº
        content_frame = tk.Frame(dialog, bg="white", padx=20, pady=20)
        content_frame.pack(fill=tk.BOTH, expand=True)

        # é€‰åŒºä¿¡æ¯
        info_text = f"å·²é€‰ä¸­åŒºåŸŸ: {x2-x1}Ã—{y2-y1} åƒç´ "
        tk.Label(content_frame, text=info_text,
                bg="white", fg="#666",
                font=(FONT_FAMILY, 9)).pack(anchor=tk.W, pady=(0, 10))

        # æç¤ºè¯è¾“å…¥
        tk.Label(content_frame, text="æç¤ºè¯:",
                bg="white", fg="#333",
                font=(FONT_FAMILY, 10, "bold")).pack(anchor=tk.W, pady=(10, 5))

        prompt_frame = tk.Frame(content_frame, bg="white")
        prompt_frame.pack(fill=tk.BOTH, expand=True, pady=5)

        prompt_text = tk.Text(prompt_frame, height=5, font=(FONT_FAMILY, 9),
                             relief=tk.SOLID, borderwidth=1)
        prompt_text.pack(fill=tk.BOTH, expand=True)

        api_type = self.ai_api_manager.config.get("api_type", "openai")
        use_gemini_args_var = tk.BooleanVar(value=False)
        gemini_image_size_var = tk.StringVar(value=self.ai_api_manager.config.get("gemini", {}).get("image_size", "1K"))
        # é»˜è®¤ç”¨é€‰åŒºæ¯”ä¾‹ï¼ˆæ›´å®¹æ˜“ç”ŸæˆåŒå®½é«˜æ¯”çš„ç»“æœï¼Œå‡å°‘è£åˆ‡/ç•™è¾¹ï¼‰
        gemini_aspect_ratio_var = tk.StringVar(value=self._best_ratio_label(x2 - x1, y2 - y1))

        if api_type == "gemini":
            args_frame = tk.LabelFrame(content_frame, text="Gemini å‚æ•°ï¼ˆå¯é€‰ï¼‰", bg="white", fg="#333",
                                       font=(FONT_FAMILY, 9, "bold"), padx=10, pady=6)
            args_frame.pack(fill=tk.X, pady=(10, 0))

            tk.Checkbutton(
                args_frame,
                text="å‹¾é€‰åæŒ‰æœ¬æ¬¡å‚æ•°ç”Ÿæˆ",
                variable=use_gemini_args_var,
                bg="white",
                font=(FONT_FAMILY, 9),
            ).grid(row=0, column=0, columnspan=4, sticky=tk.W, pady=(0, 6))

            tk.Label(args_frame, text="åˆ†è¾¨ç‡:", bg="white", font=(FONT_FAMILY, 9)).grid(row=1, column=0, sticky=tk.W)
            for i, val in enumerate(["1K", "2K", "4K"]):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_image_size_var,
                    bg="white",
                    font=(FONT_FAMILY, 9),
                ).grid(row=1, column=1 + i, sticky=tk.W, padx=6)

            tk.Label(args_frame, text="æ¯”ä¾‹:", bg="white", font=(FONT_FAMILY, 9)).grid(row=2, column=0, sticky=tk.W, pady=(6, 0))
            ratio_vals = ["auto", "1:1", "16:9", "9:16", "4:3", "3:4"]
            for i, val in enumerate(ratio_vals):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_aspect_ratio_var,
                    bg="white",
                    font=(FONT_FAMILY, 9),
                ).grid(row=3 + i // 4, column=i % 4, sticky=tk.W, padx=6)

        # å¿«é€Ÿæ¨¡æ¿
        tk.Label(content_frame, text="å¿«é€Ÿæ¨¡æ¿:",
                bg="white", fg="#666",
                font=(FONT_FAMILY, 9)).pack(anchor=tk.W, pady=(10, 5))

        template_frame = tk.Frame(content_frame, bg="white")
        template_frame.pack(anchor=tk.W)

        def set_prompt(template):
            prompt_text.delete("1.0", tk.END)
            prompt_text.insert("1.0", template)

        templates = [
            ("æ¢æˆè‹¹æœ", "Replace with a red apple"),
            ("å»é™¤ç‰©ä½“", "Remove this object and generate clean background"),
            ("æ²¹ç”»é£æ ¼", "Transform to oil painting style"),
            ("å¡é€šé£æ ¼", "Transform to cartoon style")
        ]

        for i, (label, template) in enumerate(templates):
            btn = tk.Button(template_frame, text=label,
                          command=lambda t=template: set_prompt(t),
                          bg="#F5F5F5", relief=tk.FLAT,
                          font=(FONT_FAMILY, 8))
            btn.grid(row=i//2, column=i%2, padx=5, pady=2, sticky=tk.W)

        # æŒ‰é’®åŒº
        button_frame = tk.Frame(dialog, bg="white", pady=15)
        button_frame.pack(fill=tk.X)

        def on_generate():
            prompt = prompt_text.get("1.0", tk.END).strip()
            if not prompt:
                messagebox.showwarning("æç¤º", "è¯·è¾“å…¥æç¤ºè¯")
                return
            dialog.destroy()
            overrides = None
            if api_type == "gemini" and use_gemini_args_var.get():
                overrides = {
                    "image_size": gemini_image_size_var.get(),
                    "aspect_ratio": gemini_aspect_ratio_var.get(),
                }
            self.execute_ai_replace(prompt, overrides=overrides)

        def on_cancel():
            if self.ai_replace_rect_id:
                self.canvas.delete(self.ai_replace_rect_id)
                self.ai_replace_rect_id = None
            self.ai_replace_selection = None
            dialog.destroy()

        tk.Button(button_frame, text="ç”Ÿæˆ/æ›¿æ¢", command=on_generate,
                 bg="#E91E63", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10, "bold"),
                 padx=30, pady=8).pack(side=tk.LEFT, padx=(20, 10))

        tk.Button(button_frame, text="å–æ¶ˆ", command=on_cancel,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10),
                 padx=30, pady=8).pack(side=tk.LEFT)

    def execute_ai_replace(self, prompt, overrides=None):
        """æ‰§è¡ŒAIæ›¿æ¢"""
        return ai_replace_feature.execute_ai_replace(self, prompt, overrides=overrides)
        if not self.ai_replace_selection:
            return

        x1, y1, x2, y2 = self.ai_replace_selection

        # è·å–å½“å‰æ˜¾ç¤ºçš„å›¾ç‰‡ï¼ˆåŸå›¾æˆ–èƒŒæ™¯å›¾ï¼‰
        current_page = self.pages[self.current_page_index]

        # ä½¿ç”¨èƒŒæ™¯å›¾ï¼ˆå¦‚æœæœ‰ï¼‰æˆ–åŸå›¾
        if current_page.get("bg_path") and os.path.exists(current_page["bg_path"]):
            base_image = Image.open(current_page["bg_path"])
        else:
            base_image = current_page["image"].copy()

        # è£å‰ªé€‰ä¸­åŒºåŸŸ
        crop_box = (x1, y1, x2, y2)
        cropped_image = base_image.crop(crop_box)

        # åˆ›å»ºè’™ç‰ˆï¼ˆé€‰ä¸­åŒºåŸŸä¸ºç™½è‰²ï¼‰
        mask = Image.new("L", base_image.size, 0)
        mask_draw = ImageDraw.Draw(mask)
        mask_draw.rectangle([x1, y1, x2, y2], fill=255)

        # è£å‰ªè’™ç‰ˆ
        cropped_mask = mask.crop(crop_box)

        # æ˜¾ç¤ºè¿›åº¦å¯¹è¯æ¡†
        progress_dialog = tk.Toplevel(self.root)
        progress_dialog.title("AIå¤„ç†ä¸­")
        progress_dialog.geometry("400x150")
        progress_dialog.transient(self.root)
        progress_dialog.grab_set()

        tk.Label(progress_dialog, text="AIæ­£åœ¨å¤„ç†å›¾ç‰‡...",
                font=(FONT_FAMILY, 11, "bold")).pack(pady=20)

        progress_label = tk.Label(progress_dialog, text="æ­£åœ¨åˆå§‹åŒ–...",
                                 font=(FONT_FAMILY, 9), fg="#666")
        progress_label.pack(pady=10)

        def update_progress(message):
            def _update():
                try:
                    if progress_label.winfo_exists():
                        progress_label.config(text=message)
                except Exception:
                    pass

            try:
                self.root.after(0, _update)
            except Exception:
                pass

        def process_in_thread():
            try:
                # è°ƒç”¨AI API
                result_image = self.ai_api_manager.image_to_image(
                    prompt,
                    cropped_image,
                    cropped_mask,
                    update_progress,
                    overrides=overrides,
                )

                if result_image:
                    # å…ˆæŠŠAIè¿”å›ç»“æœè½ç›˜ï¼ˆæ–¹ä¾¿æ’æŸ¥/å¤ç”¨ï¼‰ï¼Œå†åšæ— æ‹‰ä¼¸é€‚é…æ’å…¥
                    temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
                    os.makedirs(temp_dir, exist_ok=True)

                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    raw_path = os.path.join(temp_dir, f"ai_replace_raw_{timestamp}.png")
                    try:
                        result_image.save(raw_path)
                    except Exception:
                        try:
                            result_image.convert("RGB").save(raw_path)
                        except Exception:
                            pass

                    # æ— å˜å½¢è´´å›ï¼šç­‰æ¯”ç¼©æ”¾ + å±…ä¸­è£åˆ‡åˆ°é€‰åŒºå°ºå¯¸
                    if result_image.size != cropped_image.size:
                        result_image = self._resize_cover_no_distort(result_image, cropped_image.size)

                    # ä½œä¸ºå›¾å±‚åŠ å…¥ï¼ˆä¸èåˆåˆ°èƒŒæ™¯ï¼Œæ–¹ä¾¿æ— æŸåˆ‡æ¢/éšè—/åˆ é™¤ï¼‰
                    layer_name = f"AIæ›¿æ¢ {timestamp}"
                    layer_img = result_image.convert("RGBA") if result_image.mode != "RGBA" else result_image
                    self.add_image_layer(current_page, layer_img, name=layer_name, x=x1, y=y1, opacity=1.0, visible=True)
                    try:
                        self.layers = current_page.get("layers", [])
                    except Exception:
                        pass

                    # å…³é—­è¿›åº¦å¯¹è¯æ¡†
                    self.root.after(0, progress_dialog.destroy)

                    # åˆ·æ–°æ˜¾ç¤º
                    self.root.after(0, self.update_layer_listbox)
                    self.root.after(0, self.refresh_canvas)
                    self.root.after(0, self.mark_unsaved)

                    # æ¸…é™¤é€‰æ¡†
                    if self.ai_replace_rect_id:
                        self.root.after(0, lambda: self.canvas.delete(self.ai_replace_rect_id))
                        self.ai_replace_rect_id = None
                    self.ai_replace_selection = None

                    # æ˜¾ç¤ºæˆåŠŸæ¶ˆæ¯
                    self.root.after(0, lambda: messagebox.showinfo("æˆåŠŸ",
                        "AIæ›¿æ¢å®Œæˆï¼\n\n"
                        "âœ… å·²ä½œä¸ºå›¾å±‚å åŠ ï¼ˆå³ä¾§â€œå›¾å±‚â€å¯éšè—/åˆ é™¤/è°ƒé€æ˜åº¦ï¼‰\n"
                        f"ğŸ’¾ åŸå§‹è¿”å›å·²ä¿å­˜ï¼š{raw_path}\n\n"
                        "ğŸ’¡ å¯ç»§ç»­æ¡†é€‰å…¶ä»–åŒºåŸŸè¿›è¡Œç¼–è¾‘"))

                    self.root.after(0, lambda: self.update_status("AIæ›¿æ¢å®Œæˆ"))
                else:
                    raise Exception("AI APIæœªè¿”å›ç»“æœ")

            except Exception as e:
                print(f"AIæ›¿æ¢å¤±è´¥: {e}")
                import traceback
                traceback.print_exc()
                err_text = str(e)

                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯",
                    f"AIæ›¿æ¢å¤±è´¥:\n{t}\n\n"
                    f"è¯·æ£€æŸ¥:\n"
                    f"1. APIé…ç½®æ˜¯å¦æ­£ç¡®\n"
                    f"2. API Keyæ˜¯å¦æœ‰æ•ˆ\n"
                    f"3. ç½‘ç»œè¿æ¥æ˜¯å¦æ­£å¸¸"))
                self.root.after(0, lambda: self.update_status("AIæ›¿æ¢å¤±è´¥"))

        # åœ¨åå°çº¿ç¨‹æ‰§è¡Œ
        threading.Thread(target=process_in_thread, daemon=True).start()

    def _resize_cover_no_distort(self, img, target_size):
        """
        ç­‰æ¯”ç¼©æ”¾å¹¶å±…ä¸­è£åˆ‡ï¼Œä¿è¯å¡«æ»¡ç›®æ ‡å°ºå¯¸ä¸”ä¸æ‹‰ä¼¸å˜å½¢ã€‚
        ç”¨äºæŠŠAIè¿”å›å›¾æ— å˜å½¢è´´å›æŒ‡å®šåŒºåŸŸ/ç”»å¸ƒã€‚
        """
        target_w, target_h = target_size
        if target_w <= 0 or target_h <= 0:
            return img

        img_w, img_h = img.size
        if img_w <= 0 or img_h <= 0:
            return img

        if (img_w, img_h) == (target_w, target_h):
            return img

        scale = max(target_w / img_w, target_h / img_h)
        new_w = max(1, int(math.ceil(img_w * scale)))
        new_h = max(1, int(math.ceil(img_h * scale)))

        resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)

        left = max(0, (new_w - target_w) // 2)
        top = max(0, (new_h - target_h) // 2)
        return resized.crop((left, top, left + target_w, top + target_h))

    def _resize_contain_no_distort(self, img, target_size, fill_color=None):
        """
        ç­‰æ¯”ç¼©æ”¾å¹¶å±…ä¸­æ”¾ç½®åˆ°ç›®æ ‡ç”»å¸ƒï¼ˆä¸è£åˆ‡ã€ä¸æ‹‰ä¼¸ï¼‰ã€‚
        è‹¥æ¯”ä¾‹ä¸ä¸€è‡´ï¼Œä¼šäº§ç”Ÿç•™è¾¹ï¼›ç•™è¾¹é¢œè‰²é»˜è®¤å–å·¦ä¸Šè§’åƒç´ æˆ–é€æ˜ã€‚
        """
        target_w, target_h = target_size
        if target_w <= 0 or target_h <= 0:
            return img

        img_w, img_h = img.size
        if img_w <= 0 or img_h <= 0:
            return img

        if (img_w, img_h) == (target_w, target_h):
            return img

        scale = min(target_w / img_w, target_h / img_h)
        new_w = max(1, int(math.floor(img_w * scale)))
        new_h = max(1, int(math.floor(img_h * scale)))
        resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)

        if fill_color is None:
            if "A" in resized.mode:
                fill_color = (0, 0, 0, 0)
            else:
                try:
                    fill_color = resized.getpixel((0, 0))
                except Exception:
                    fill_color = (0, 0, 0)

        canvas = Image.new(resized.mode, (target_w, target_h), fill_color)
        left = (target_w - new_w) // 2
        top = (target_h - new_h) // 2
        if resized.mode == "RGBA":
            canvas.paste(resized, (left, top), mask=resized.split()[-1])
        else:
            canvas.paste(resized, (left, top))
        return canvas

    def _best_ratio_label(self, width, height):
        """ä»å¸¸ç”¨æ¯”ä¾‹é‡Œé€‰ä¸€ä¸ªæœ€æ¥è¿‘çš„ï¼ˆç”¨äº Gemini çš„ aspectRatioï¼‰ã€‚"""
        if width <= 0 or height <= 0:
            return "auto"
        r = width / height
        candidates = {
            "1:1": 1.0,
            "16:9": 16 / 9,
            "9:16": 9 / 16,
            "4:3": 4 / 3,
            "3:4": 3 / 4,
        }
        best = min(candidates.items(), key=lambda kv: abs(kv[1] - r))[0]
        return best

    def ai_text_to_image_layer(self):
        """æ ¹æ®æ–‡å­—æè¿°ç”Ÿæˆå›¾ç‰‡ï¼Œå¹¶ä½œä¸ºå›¾å±‚æ·»åŠ åˆ°å½“å‰é¡µ"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        current_page = self.pages[self.current_page_index]

        dialog = tk.Toplevel(self.root)
        dialog.title("AIæ–‡å­—ç”Ÿå›¾")
        dialog.geometry("520x320")
        dialog.transient(self.root)
        dialog.grab_set()

        tk.Label(dialog, text="æç¤ºè¯ï¼ˆæè¿°ä½ å¸Œæœ›ç”Ÿæˆçš„å›¾ç‰‡ï¼‰",
                font=(FONT_FAMILY, 11, "bold")).pack(anchor="w", padx=15, pady=(15, 5))

        prompt_text = tk.Text(dialog, height=8, font=(FONT_FAMILY, 10), wrap=tk.WORD)
        prompt_text.pack(fill=tk.BOTH, expand=True, padx=15, pady=5)
        prompt_text.insert("1.0", "è¯·è¾“å…¥ä½ æƒ³è¦ç”Ÿæˆçš„å›¾ç‰‡æè¿°...")

        api_type = self.ai_api_manager.config.get("api_type", "openai")
        use_gemini_args_var = tk.BooleanVar(value=False)
        gemini_image_size_var = tk.StringVar(value=self.ai_api_manager.config.get("gemini", {}).get("image_size", "1K"))
        # é»˜è®¤ 1:1 æ¯”ä¾‹
        gemini_aspect_ratio_var = tk.StringVar(value="1:1")

        if api_type == "gemini":
            args_frame = tk.LabelFrame(dialog, text="Gemini å‚æ•°ï¼ˆå¯é€‰ï¼‰", font=(FONT_FAMILY, 9, "bold"),
                                       padx=10, pady=6)
            args_frame.pack(fill=tk.X, padx=15, pady=(0, 10))

            tk.Checkbutton(
                args_frame,
                text="å‹¾é€‰åæŒ‰æœ¬æ¬¡å‚æ•°ç”Ÿæˆ",
                variable=use_gemini_args_var,
                font=(FONT_FAMILY, 9),
            ).grid(row=0, column=0, columnspan=4, sticky=tk.W, pady=(0, 6))

            tk.Label(args_frame, text="åˆ†è¾¨ç‡:", font=(FONT_FAMILY, 9)).grid(row=1, column=0, sticky=tk.W)
            for i, val in enumerate(["1K", "2K", "4K"]):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_image_size_var,
                    font=(FONT_FAMILY, 9),
                ).grid(row=1, column=1 + i, sticky=tk.W, padx=6)

            tk.Label(args_frame, text="æ¯”ä¾‹:", font=(FONT_FAMILY, 9)).grid(row=2, column=0, sticky=tk.W, pady=(6, 0))
            ratio_vals = ["1:1", "16:9", "9:16", "4:3", "3:4"]
            for i, val in enumerate(ratio_vals):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_aspect_ratio_var,
                    font=(FONT_FAMILY, 9),
                ).grid(row=3 + i // 4, column=i % 4, sticky=tk.W, padx=6)

        btn_frame = tk.Frame(dialog)
        btn_frame.pack(fill=tk.X, padx=15, pady=12)

        def on_cancel():
            dialog.destroy()

        def on_generate():
            prompt = prompt_text.get("1.0", tk.END).strip()
            if not prompt or prompt == "è¯·è¾“å…¥ä½ æƒ³è¦ç”Ÿæˆçš„å›¾ç‰‡æè¿°...":
                messagebox.showwarning("æç¤º", "è¯·è¾“å…¥æœ‰æ•ˆçš„æç¤ºè¯")
                return
            dialog.destroy()
            overrides = None
            if api_type == "gemini" and use_gemini_args_var.get():
                overrides = {
                    "image_size": gemini_image_size_var.get(),
                    "aspect_ratio": gemini_aspect_ratio_var.get(),
                }
            self._execute_ai_text_to_image(prompt, overrides=overrides)

        tk.Button(btn_frame, text="ç”Ÿæˆå¹¶ä½œä¸ºå›¾å±‚", command=on_generate,
                 bg="#7B1FA2", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10, "bold"),
                 padx=20, pady=8).pack(side=tk.LEFT)

        tk.Button(btn_frame, text="å–æ¶ˆ", command=on_cancel,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10),
                 padx=20, pady=8).pack(side=tk.LEFT, padx=10)

    def _execute_ai_text_to_image(self, prompt, overrides=None):
        """åå°æ‰§è¡Œçº¯æ–‡å­—ç”Ÿæˆå›¾ç‰‡"""
        if not self.pages:
            return

        current_page = self.pages[self.current_page_index]

        progress_dialog = tk.Toplevel(self.root)
        progress_dialog.title("AIå¤„ç†ä¸­")
        progress_dialog.geometry("420x160")
        progress_dialog.transient(self.root)
        progress_dialog.grab_set()

        tk.Label(progress_dialog, text="AIæ­£åœ¨æ ¹æ®æ–‡å­—ç”Ÿæˆå›¾ç‰‡...",
                font=(FONT_FAMILY, 11, "bold")).pack(pady=20)

        progress_label = tk.Label(progress_dialog, text="æ­£åœ¨åˆå§‹åŒ–...",
                                 font=(FONT_FAMILY, 9), fg="#666")
        progress_label.pack(pady=10)

        def update_progress(message):
            def _update():
                try:
                    if progress_label.winfo_exists():
                        progress_label.config(text=message)
                except Exception:
                    pass

            try:
                self.root.after(0, _update)
            except Exception:
                pass

        def worker():
            try:
                # çº¯æ–‡å­—ç”Ÿæˆï¼Œä¸ä¼ å…¥æºå›¾ç‰‡
                result_image = self.ai_api_manager.generate_image(
                    prompt,
                    source_image=None,
                    mask_image=None,
                    progress_callback=update_progress,
                    overrides=overrides,
                )
                if not result_image:
                    raise Exception("AI APIæœªè¿”å›ç»“æœ")

                temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
                os.makedirs(temp_dir, exist_ok=True)

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                raw_path = os.path.join(temp_dir, f"ai_text_to_image_raw_{timestamp}.png")
                try:
                    result_image.save(raw_path)
                except Exception:
                    try:
                        result_image.convert("RGB").save(raw_path)
                    except Exception:
                        pass

                # å°†ç”Ÿæˆçš„å›¾ç‰‡æ·»åŠ ä¸ºå›¾å±‚ï¼Œä¿æŒåŸå§‹å°ºå¯¸
                layer_name = f"AIæ–‡å­—ç”Ÿå›¾ {timestamp}"
                layer_img = result_image.convert("RGBA") if result_image.mode != "RGBA" else result_image
                # å°†å›¾å±‚å±…ä¸­æ”¾ç½®
                page_w, page_h = current_page.get("image").size if current_page.get("image") else (0, 0)
                img_w, img_h = layer_img.size
                x = max(0, (page_w - img_w) // 2)
                y = max(0, (page_h - img_h) // 2)
                self.add_image_layer(current_page, layer_img, name=layer_name, x=x, y=y, opacity=1.0, visible=True)
                try:
                    self.layers = current_page.get("layers", [])
                except Exception:
                    pass

                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, self.update_layer_listbox)
                self.root.after(0, self.refresh_canvas)
                self.root.after(0, lambda: self.update_status("æ–‡å­—ç”Ÿå›¾å®Œæˆï¼ˆå›¾å±‚å·²æ·»åŠ ï¼‰"))
                self.root.after(0, lambda: self.mark_unsaved())
                self.root.after(0, lambda: messagebox.showinfo("æˆåŠŸ", f"æ–‡å­—ç”Ÿå›¾å®Œæˆï¼\n\nå·²ä½œä¸ºå›¾å±‚æ·»åŠ ã€‚\nåŸå§‹è¿”å›å·²ä¿å­˜ï¼š\n{raw_path}"))

            except Exception as e:
                import traceback
                traceback.print_exc()
                err_text = str(e)
                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"æ–‡å­—ç”Ÿå›¾å¤±è´¥:\n{t}"))
                self.root.after(0, lambda: self.update_status("æ–‡å­—ç”Ÿå›¾å¤±è´¥"))

        threading.Thread(target=worker, daemon=True).start()

    def ai_generate_fullpage_background(self):
        """æŠŠå½“å‰é¡µæ•´å›¾å‘é€ç»™AIç”Ÿæˆï¼Œè¿”å›ç»“æœä½œä¸ºå›¾å±‚å åŠ åˆ°å½“å‰é¡µ"""
        if not self.pages:
            messagebox.showwarning("æç¤º", "è¯·å…ˆå¯¼å…¥å›¾ç‰‡")
            return

        current_page = self.pages[self.current_page_index]

        dialog = tk.Toplevel(self.root)
        dialog.title("AIæ•´é¡µç”ŸæˆèƒŒæ™¯")
        dialog.geometry("520x320")
        dialog.transient(self.root)
        dialog.grab_set()

        tk.Label(dialog, text="æç¤ºè¯ï¼ˆæè¿°ä½ å¸Œæœ›æ•´é¡µç”Ÿæˆçš„æ•ˆæœï¼‰",
                font=(FONT_FAMILY, 11, "bold")).pack(anchor="w", padx=15, pady=(15, 5))

        prompt_text = tk.Text(dialog, height=8, font=(FONT_FAMILY, 10), wrap=tk.WORD)
        prompt_text.pack(fill=tk.BOTH, expand=True, padx=15, pady=5)
        prompt_text.insert("1.0", "åœ¨ä¿æŒæ•´ä½“é£æ ¼ä¸€è‡´çš„å‰æä¸‹ï¼Œç”Ÿæˆä¸€å¼ å¯ç”¨ä½œèƒŒæ™¯çš„å›¾ç‰‡ã€‚")

        api_type = self.ai_api_manager.config.get("api_type", "openai")
        use_gemini_args_var = tk.BooleanVar(value=False)
        gemini_image_size_var = tk.StringVar(value=self.ai_api_manager.config.get("gemini", {}).get("image_size", "1K"))
        # é»˜è®¤ç”¨å½“å‰é¡µæ¯”ä¾‹ï¼ˆæ›´å®¹æ˜“ä¸åŸå›¾å¯¹é½ï¼‰
        page_w, page_h = current_page.get("image").size if current_page.get("image") else (0, 0)
        gemini_aspect_ratio_var = tk.StringVar(value=self._best_ratio_label(page_w, page_h))

        if api_type == "gemini":
            args_frame = tk.LabelFrame(dialog, text="Gemini å‚æ•°ï¼ˆå¯é€‰ï¼‰", font=(FONT_FAMILY, 9, "bold"),
                                       padx=10, pady=6)
            args_frame.pack(fill=tk.X, padx=15, pady=(0, 10))

            tk.Checkbutton(
                args_frame,
                text="å‹¾é€‰åæŒ‰æœ¬æ¬¡å‚æ•°ç”Ÿæˆ",
                variable=use_gemini_args_var,
                font=(FONT_FAMILY, 9),
            ).grid(row=0, column=0, columnspan=4, sticky=tk.W, pady=(0, 6))

            tk.Label(args_frame, text="åˆ†è¾¨ç‡:", font=(FONT_FAMILY, 9)).grid(row=1, column=0, sticky=tk.W)
            for i, val in enumerate(["1K", "2K", "4K"]):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_image_size_var,
                    font=(FONT_FAMILY, 9),
                ).grid(row=1, column=1 + i, sticky=tk.W, padx=6)

            tk.Label(args_frame, text="æ¯”ä¾‹:", font=(FONT_FAMILY, 9)).grid(row=2, column=0, sticky=tk.W, pady=(6, 0))
            ratio_vals = ["auto", "1:1", "16:9", "9:16", "4:3", "3:4"]
            for i, val in enumerate(ratio_vals):
                tk.Radiobutton(
                    args_frame,
                    text=val,
                    value=val,
                    variable=gemini_aspect_ratio_var,
                    font=(FONT_FAMILY, 9),
                ).grid(row=3 + i // 4, column=i % 4, sticky=tk.W, padx=6)

        btn_frame = tk.Frame(dialog)
        btn_frame.pack(fill=tk.X, padx=15, pady=12)

        def on_cancel():
            dialog.destroy()

        def on_generate():
            prompt = prompt_text.get("1.0", tk.END).strip()
            if not prompt:
                messagebox.showwarning("æç¤º", "è¯·è¾“å…¥æç¤ºè¯")
                return
            dialog.destroy()
            overrides = None
            if api_type == "gemini" and use_gemini_args_var.get():
                overrides = {
                    "image_size": gemini_image_size_var.get(),
                    "aspect_ratio": gemini_aspect_ratio_var.get(),
                }
            self._execute_ai_fullpage(prompt, overrides=overrides)

        tk.Button(btn_frame, text="ç”Ÿæˆå¹¶ä½œä¸ºå›¾å±‚", command=on_generate,
                 bg="#6A1B9A", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10, "bold"),
                 padx=20, pady=8).pack(side=tk.LEFT)

        tk.Button(btn_frame, text="å–æ¶ˆ", command=on_cancel,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 10),
                 padx=20, pady=8).pack(side=tk.LEFT, padx=10)

    def _execute_ai_fullpage(self, prompt, overrides=None):
        """åå°æ‰§è¡Œæ•´é¡µAIç”Ÿæˆ"""
        if not self.pages:
            return

        current_page = self.pages[self.current_page_index]

        if current_page.get("bg_path") and os.path.exists(current_page["bg_path"]):
            base_image = Image.open(current_page["bg_path"])
        else:
            base_image = current_page["image"].copy()

        # è‡ªåŠ¨é€‰æ‹©æ›´é«˜åˆ†è¾¨ç‡ï¼Œå°½é‡é¿å…â€œç”Ÿæˆå›¾è¢«æ”¾å¤§åå˜ç³Šâ€ï¼ˆç”¨æˆ· overrides ä¼˜å…ˆï¼‰
        try:
            auto_overrides = self.ai_api_manager.suggest_overrides(*base_image.size)
            overrides = {**auto_overrides, **(overrides or {})}
        except Exception:
            overrides = overrides

        progress_dialog = tk.Toplevel(self.root)
        progress_dialog.title("AIå¤„ç†ä¸­")
        progress_dialog.geometry("420x160")
        progress_dialog.transient(self.root)
        progress_dialog.grab_set()

        tk.Label(progress_dialog, text="AIæ­£åœ¨ç”Ÿæˆæ•´é¡µèƒŒæ™¯...",
                font=(FONT_FAMILY, 11, "bold")).pack(pady=20)

        progress_label = tk.Label(progress_dialog, text="æ­£åœ¨åˆå§‹åŒ–...",
                                 font=(FONT_FAMILY, 9), fg="#666")
        progress_label.pack(pady=10)

        def update_progress(message):
            def _update():
                try:
                    if progress_label.winfo_exists():
                        progress_label.config(text=message)
                except Exception:
                    pass

            try:
                self.root.after(0, _update)
            except Exception:
                pass

        def worker():
            try:
                result_image = self.ai_api_manager.generate_image(
                    prompt,
                    source_image=base_image,
                    mask_image=None,
                    progress_callback=update_progress,
                    overrides=overrides,
                )
                if not result_image:
                    raise Exception("AI APIæœªè¿”å›ç»“æœ")

                temp_dir = os.path.join(get_base_dir(), "temp_backgrounds")
                os.makedirs(temp_dir, exist_ok=True)

                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                raw_path = os.path.join(temp_dir, f"ai_fullpage_raw_{timestamp}.png")
                try:
                    result_image.save(raw_path)
                except Exception:
                    try:
                        result_image.convert("RGB").save(raw_path)
                    except Exception:
                        pass

                # ç”Ÿæˆç»“æœéœ€è¦ä¸ç¼–è¾‘å›¾åŒå°ºå¯¸ï¼Œä¸”ä¸å…è®¸æ‹‰ä¼¸å˜å½¢
                if result_image.size != base_image.size:
                    # æ•´é¡µèƒŒæ™¯ä¼˜å…ˆä¸è£åˆ‡ï¼ˆé¿å…â€œå¯¹ç…§ä¸ä¸Šâ€ï¼‰ï¼šç­‰æ¯”ç¼©æ”¾+ç•™è¾¹
                    result_image = self._resize_contain_no_distort(result_image, base_image.size)

                layer_name = f"AIæ•´é¡µ {timestamp}"
                layer_img = result_image.convert("RGBA") if result_image.mode != "RGBA" else result_image
                self.add_image_layer(current_page, layer_img, name=layer_name, x=0, y=0, opacity=1.0, visible=True)
                try:
                    self.layers = current_page.get("layers", [])
                except Exception:
                    pass

                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, self.update_layer_listbox)
                self.root.after(0, self.refresh_canvas)
                self.root.after(0, lambda: self.update_status("æ•´é¡µç”Ÿæˆå®Œæˆï¼ˆå›¾å±‚å·²æ·»åŠ ï¼‰"))
                self.root.after(0, lambda: self.mark_unsaved())
                self.root.after(0, lambda: messagebox.showinfo("æˆåŠŸ", f"æ•´é¡µç”Ÿæˆå®Œæˆï¼\n\nå·²ä½œä¸ºå›¾å±‚å åŠ ã€‚\nåŸå§‹è¿”å›å·²ä¿å­˜ï¼š\n{raw_path}"))

            except Exception as e:
                import traceback
                traceback.print_exc()
                err_text = str(e)
                self.root.after(0, progress_dialog.destroy)
                self.root.after(0, lambda t=err_text: messagebox.showerror("é”™è¯¯", f"æ•´é¡µç”Ÿæˆå¤±è´¥:\n{t}"))
                self.root.after(0, lambda: self.update_status("æ•´é¡µç”Ÿæˆå¤±è´¥"))

        threading.Thread(target=worker, daemon=True).start()

    def open_ai_api_settings(self):
        """æ‰“å¼€AI APIé…ç½®å¯¹è¯æ¡†"""
        dialog = tk.Toplevel(self.root)
        dialog.title("AIå›¾ç‰‡ç”ŸæˆAPIé…ç½®")
        dialog.geometry("650x650")
        dialog.transient(self.root)
        dialog.grab_set()

        # æ ‡é¢˜
        title_frame = tk.Frame(dialog, bg="#9C27B0", height=60)
        title_frame.pack(fill=tk.X)
        title_frame.pack_propagate(False)

        tk.Label(title_frame, text="AI å›¾ç‰‡ç”Ÿæˆ API é…ç½®",
                bg="#9C27B0", fg="white",
                font=(FONT_FAMILY, 16, "bold")).pack(pady=15)

        # ä¸»å†…å®¹
        main_frame = tk.Frame(dialog, bg="white", padx=30, pady=20)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # APIç±»å‹é€‰æ‹©
        tk.Label(main_frame, text="APIç±»å‹:",
                bg="white", fg="#333",
                font=(FONT_FAMILY, 11, "bold")).grid(row=0, column=0, sticky=tk.W, pady=(0, 10))

        api_type_var = tk.StringVar(value=self.ai_api_manager.config.get("api_type", "openai"))

        api_frame = tk.Frame(main_frame, bg="white")
        api_frame.grid(row=0, column=1, sticky=tk.W, pady=(0, 10))

        # æµå¼ä¼ è¾“ï¼ˆä»…OpenAIï¼‰
        stream_var = tk.BooleanVar(value=self.ai_api_manager.config.get("openai", {}).get("stream", True))

        def on_api_type_change():
            api_type = api_type_var.get()
            provider_cfg = self.ai_api_manager.config.get(api_type, {})

            key_var.set(provider_cfg.get("api_key", ""))
            host_var.set(provider_cfg.get("api_host", ""))
            if api_type == "openai":
                model_var.set(provider_cfg.get("model", "gpt-4o"))
                stream_var.set(self.ai_api_manager.config.get("openai", {}).get("stream", True))
            else:
                model_var.set(provider_cfg.get("model", "gemini-2.0-flash-exp-image-generation"))

        tk.Radiobutton(
            api_frame,
            text="OpenAIæ ¼å¼",
            variable=api_type_var,
            value="openai",
            bg="white",
            font=(FONT_FAMILY, 10),
            command=on_api_type_change,
        ).pack(side=tk.LEFT, padx=10)
        tk.Radiobutton(
            api_frame,
            text="Geminiæ ¼å¼",
            variable=api_type_var,
            value="gemini",
            bg="white",
            font=(FONT_FAMILY, 10),
            command=on_api_type_change,
        ).pack(side=tk.LEFT, padx=10)

        # API Key
        tk.Label(main_frame, text="API Key:",
                bg="white", fg="#333",
                font=(FONT_FAMILY, 11, "bold")).grid(row=1, column=0, sticky=tk.W, pady=(15, 5))

        key_var = tk.StringVar(value=self.ai_api_manager.config.get(
            api_type_var.get(), {}).get("api_key", ""))

        key_entry = tk.Entry(main_frame, textvariable=key_var, width=45,
                            font=(FONT_FAMILY, 10), show="*")
        key_entry.grid(row=1, column=1, sticky=tk.W, pady=(15, 5))

        # API Host
        tk.Label(main_frame, text="API Host:",
                bg="white", fg="#333",
                font=(FONT_FAMILY, 11, "bold")).grid(row=2, column=0, sticky=tk.W, pady=(10, 5))

        host_var = tk.StringVar(value=self.ai_api_manager.config.get(
            api_type_var.get(), {}).get("api_host", ""))

        tk.Entry(main_frame, textvariable=host_var, width=45,
                font=(FONT_FAMILY, 10)).grid(row=2, column=1, sticky=tk.W, pady=(10, 5))

        # æ¨¡å‹åç§°
        tk.Label(main_frame, text="æ¨¡å‹:",
                bg="white", fg="#333",
                font=(FONT_FAMILY, 11, "bold")).grid(row=3, column=0, sticky=tk.W, pady=(10, 5))

        model_var = tk.StringVar(value=self.ai_api_manager.config.get(
            api_type_var.get(), {}).get("model", "gemini-3-pro-image-preview"))

        tk.Entry(main_frame, textvariable=model_var, width=45,
                font=(FONT_FAMILY, 10)).grid(row=3, column=1, sticky=tk.W, pady=(10, 5))

        tk.Checkbutton(
            main_frame,
            text="å¯ç”¨æµå¼ä¼ è¾“ï¼ˆä»…OpenAIæ ¼å¼æœ‰æ•ˆï¼‰",
            variable=stream_var,
            bg="white",
            font=(FONT_FAMILY, 10),
        ).grid(row=4, column=1, sticky=tk.W, pady=(15, 5))

        # è¯´æ˜æ–‡å­—
        info_text = (
            "è·å–API Key:\n"
            "â€¢ OpenAI: https://platform.openai.com/api-keys\n"
            "â€¢ Gemini: https://makersuite.google.com/app/apikey\n\n"
            "ä½¿ç”¨ä»£ç†:\n"
            "å¦‚ä½¿ç”¨APIä»£ç†ï¼Œè¯·ä¿®æ”¹API Hoståœ°å€"
        )
        tk.Label(main_frame, text=info_text, bg="#F5F5F5",
                fg="#666", font=(FONT_FAMILY, 8),
                justify=tk.LEFT, anchor=tk.W,
                padx=10, pady=10).grid(row=5, column=0, columnspan=2,
                                       sticky=tk.W+tk.E, pady=(20, 0))

        # æŒ‰é’®åŒº
        button_frame = tk.Frame(dialog, bg="white", pady=15)
        button_frame.pack(fill=tk.X)

        def save_and_close():
            # ä¿å­˜é…ç½®
            api_type = api_type_var.get()
            self.ai_api_manager.config["api_type"] = api_type
            self.ai_api_manager.config[api_type]["api_key"] = key_var.get()
            self.ai_api_manager.config[api_type]["api_host"] = host_var.get()
            self.ai_api_manager.config[api_type]["model"] = model_var.get()
            if api_type == "openai":
                self.ai_api_manager.config["openai"]["stream"] = stream_var.get()

            self.config = self.ai_api_manager.save_config(self.config)
            save_config(self.config)
            messagebox.showinfo("æˆåŠŸ", "APIé…ç½®å·²ä¿å­˜")
            dialog.destroy()

        def test_connection():
            api_type = api_type_var.get()
            messagebox.showinfo("æç¤º", f"{api_type.upper()} API é…ç½®å·²è®¾ç½®\n\nè¯·åœ¨å®é™…ä½¿ç”¨ä¸­æµ‹è¯•")

        tk.Button(button_frame, text="ä¿å­˜é…ç½®", command=save_and_close,
                 bg="#9C27B0", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 11, "bold"),
                 padx=30, pady=10).pack(side=tk.LEFT, padx=(30, 10))

        tk.Button(button_frame, text="å–æ¶ˆ", command=dialog.destroy,
                 bg="#999", fg="white", relief=tk.FLAT,
                 font=(FONT_FAMILY, 11),
                 padx=30, pady=10).pack(side=tk.LEFT, padx=10)

        # åˆå§‹åŒ–ä¸€æ¬¡å­—æ®µï¼Œä¿è¯é»˜è®¤æ¨¡å‹ç­‰å­—æ®µæ­£ç¡®
        on_api_type_change()


if __name__ == "__main__":
    root = tk.Tk()
    app = ModernPPTEditor(root)
    root.mainloop()
